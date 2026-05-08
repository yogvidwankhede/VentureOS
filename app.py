from agents.pivot_agent import run_pivot_suggester
from agents.prototype_agent import run_prototype_generator
from agents.orchestrator import run_ventureOS_stream, get_llm
from flask import Flask, request, jsonify, render_template, Response, stream_with_context, send_file
from dotenv import load_dotenv
import os
import json
import re
import uuid
import io
import random
import math
import tempfile
import base64
import html
import threading
from io import BytesIO
from urllib.parse import quote
from urllib.request import Request, urlopen
from PIL import Image, ImageDraw, ImageEnhance, ImageFilter, ImageOps

try:
    from google import genai as google_genai
    from google.genai import types as google_genai_types
except ImportError:  # pragma: no cover - optional runtime dependency
    google_genai = None
    google_genai_types = None

load_dotenv()


IMAGE_GENERATION_DISABLED = os.getenv(
    'VENTUREOS_DISABLE_IMAGEGEN', ''
).strip().lower() in {'1', 'true', 'yes', 'on'}
IMAGE_GENERATION_AVAILABLE = True
IMAGE_GENERATION_IMPORT_ERROR = 'Using built-in illustrated slide visuals for this deployment.'
IMAGE_GENERATION_IS_FALLBACK = True
FALLBACK_IMAGE_MODEL_KEY = 'flux-studio'
FALLBACK_IMAGE_MODEL_LABEL = 'VentureOS Neon Editorial'
FALLBACK_VECTOR_MODEL_LABEL = 'VentureOS Neon Editorial'
HOSTED_FALLBACK_TIMEOUT_SECONDS = float(os.getenv('VENTUREOS_HOSTED_IMAGE_TIMEOUT_SECONDS', '45'))
HOSTED_FALLBACK_ENABLED = os.getenv(
    'VENTUREOS_HOSTED_IMAGE_ENABLED', 'true'
).strip().lower() not in {'0', 'false', 'no', 'off'}
GOOGLE_IMAGE_ENABLED = (
    bool(os.getenv('GOOGLE_API_KEY', '').strip())
    and google_genai is not None
    and google_genai_types is not None
    and os.getenv('VENTUREOS_GOOGLE_IMAGE_ENABLED', 'true').strip().lower() not in {'0', 'false', 'no', 'off'}
)
_HOSTED_IMAGE_CACHE = {}
_HOSTED_IMAGE_CACHE_LOCK = threading.Lock()
_GOOGLE_IMAGE_CLIENT = None
_GOOGLE_IMAGE_CLIENT_LOCK = threading.Lock()

HOSTED_FALLBACK_MODEL_VARIANTS = {
    'ventureos-editorial': {
        'label': 'VentureOS Neon Editorial',
        'repo_id': 'builtin/ventureos-neon-editorial',
    },
    'google-imagen': {
        'label': 'Google Imagen',
        'repo_id': 'google/imagen-4.0-generate-001',
    },
    'google-imagen-fast': {
        'label': 'Google Imagen Fast',
        'repo_id': 'google/imagen-4.0-fast-generate-001',
    },
    'editorial-mix': {
        'label': 'Editorial Mix',
        'repo_id': 'pollinations/editorial-mix',
    },
    'flux-studio': {
        'label': 'Flux Studio',
        'repo_id': 'pollinations/flux',
    },
    'qwen-editorial': {
        'label': 'Qwen Editorial',
        'repo_id': 'pollinations/qwen-image',
    },
    'gptimage-editorial': {
        'label': 'GPTImage Editorial',
        'repo_id': 'pollinations/gptimage',
    },
}

HOSTED_PROVIDER_LABELS = {
    'google-imagen': 'Google Imagen',
    'google-imagen-fast': 'Google Imagen Fast',
    'flux': 'Flux Studio',
    'qwen-image': 'Qwen Editorial',
    'gptimage': 'GPTImage Editorial',
}

HOSTED_PROVIDER_REPOS = {
    'google-imagen': 'google/imagen-4.0-generate-001',
    'google-imagen-fast': 'google/imagen-4.0-fast-generate-001',
    'flux': 'pollinations/flux',
    'qwen-image': 'pollinations/qwen-image',
    'gptimage': 'pollinations/gptimage',
}

PREMIUM_FALLBACK_MODEL_LABEL = 'VentureOS Premium Scene Library'
PREMIUM_FALLBACK_MODEL_REPO = 'local/premium-scene-library'
PREMIUM_FALLBACK_DIRECT_TOPICS = {'gaming', 'housing'}
PREMIUM_FALLBACK_LIBRARY = {
    'gaming': [
        'gaming-hero-a.jpg',
        'gaming-hero-b.jpg',
        'gaming-hero-c.jpg',
        'gaming-hero-d.jpg',
    ],
    'software': [
        'gaming-hero-b.jpg',
        'gaming-hero-c.jpg',
        'gaming-hero-a.jpg',
        'gaming-hero-d.jpg',
    ],
    'sports': [
        'gaming-hero-d.jpg',
        'gaming-hero-a.jpg',
        'gaming-hero-c.jpg',
    ],
    'housing': [
        'housing-atrium-a.jpg',
        'housing-atrium-b.jpg',
        'housing-atrium-c.jpg',
        'interior-sunlit-a.jpg',
        'interior-sunlit-b.jpg',
    ],
    'education': [
        'interior-sunlit-b.jpg',
        'housing-atrium-a.jpg',
        'interior-sunlit-c.jpg',
        'housing-atrium-c.jpg',
        'housing-atrium-b.jpg',
    ],
    'restaurant': [
        'interior-sunlit-c.jpg',
        'housing-atrium-a.jpg',
        'interior-sunlit-d.jpg',
        'interior-sunlit-a.jpg',
    ],
    'health': [
        'interior-sunlit-a.jpg',
        'interior-sunlit-b.jpg',
        'interior-sunlit-d.jpg',
        'housing-atrium-b.jpg',
    ],
    'finance': [
        'interior-sunlit-b.jpg',
        'gaming-hero-b.jpg',
        'interior-sunlit-d.jpg',
        'interior-sunlit-a.jpg',
    ],
    'logistics': [
        'interior-sunlit-c.jpg',
        'interior-sunlit-b.jpg',
        'interior-sunlit-d.jpg',
        'gaming-hero-d.jpg',
    ],
    'generic': [
        'gaming-hero-a.jpg',
        'housing-atrium-a.jpg',
        'interior-sunlit-a.jpg',
        'gaming-hero-b.jpg',
        'interior-sunlit-b.jpg',
        'housing-atrium-c.jpg',
        'housing-atrium-b.jpg',
        'interior-sunlit-d.jpg',
    ],
}

PREMIUM_FALLBACK_ASSET_TAGS = {
    'gaming-hero-a.jpg': ['gaming', 'creator', 'developer', 'studio', 'collectibles', 'shelves', 'premium'],
    'gaming-hero-b.jpg': ['gaming', 'software', 'dashboard', 'workflow', 'interface', 'analytics', 'creator'],
    'gaming-hero-c.jpg': ['gaming', 'creator', 'studio', 'craft', 'tools', 'workspace', 'premium'],
    'gaming-hero-d.jpg': ['gaming', 'strategy', 'operations', 'display', 'workspace', 'premium'],
    'housing-atrium-a.jpg': ['housing', 'apartment', 'interior', 'student', 'atrium', 'urban', 'lounge'],
    'housing-atrium-b.jpg': ['housing', 'apartment', 'interior', 'residential', 'sunlight', 'modern', 'lounge'],
    'housing-atrium-c.jpg': ['housing', 'apartment', 'interior', 'student', 'residential', 'study', 'modern'],
    'interior-sunlit-a.jpg': ['interior', 'housing', 'education', 'health', 'workspace', 'clean', 'modern'],
    'interior-sunlit-b.jpg': ['interior', 'software', 'finance', 'workspace', 'clean', 'modern', 'premium'],
    'interior-sunlit-c.jpg': ['interior', 'restaurant', 'operations', 'inventory', 'warm', 'supply', 'premium'],
    'interior-sunlit-d.jpg': ['interior', 'logistics', 'operations', 'structured', 'workspace', 'modern', 'premium'],
}


def _fallback_supported_models():
    ordered_keys = ['flux-studio', 'ventureos-editorial']
    return [
        {
            'key': key,
            'repo_id': HOSTED_FALLBACK_MODEL_VARIANTS[key]['repo_id'],
            'label': HOSTED_FALLBACK_MODEL_VARIANTS[key]['label'],
        }
        for key in ordered_keys
        if key in HOSTED_FALLBACK_MODEL_VARIANTS
    ]


def _fallback_style_options():
    return [
        {'key': 'deck-illustration', 'label': 'Presentation Illustration'},
        {'key': 'animated-scene', 'label': 'Animated Scene'},
        {'key': 'cartoon', 'label': 'Cartoon'},
        {'key': 'abstract', 'label': 'Neon Editorial'},
    ]


def _fallback_coverage_options():
    return [
        {'key': 'hero-only', 'label': 'Hero Only (Fast)'},
        {'key': 'key-slides', 'label': 'Key Slides'},
        {'key': 'image-heavy', 'label': 'Image Heavy'},
        {'key': 'all', 'label': 'Every Slide'},
    ]


def _fallback_clean_text(value, fallback=''):
    if value is None:
        return fallback
    text = str(value).strip()
    return text or fallback


def _fallback_hash_seed(*parts):
    joined = '|'.join(_fallback_clean_text(part) for part in parts)
    total = 0
    for index, char in enumerate(joined):
        total += (index + 1) * ord(char)
    return total or 1


def _fallback_model_variant(model_key):
    return HOSTED_FALLBACK_MODEL_VARIANTS.get(
        _fallback_clean_text(model_key, FALLBACK_IMAGE_MODEL_KEY).lower(),
        HOSTED_FALLBACK_MODEL_VARIANTS[FALLBACK_IMAGE_MODEL_KEY],
    )


def _fallback_provider_candidates(model_key, slide_type, index):
    model_key = _fallback_clean_text(model_key, FALLBACK_IMAGE_MODEL_KEY).lower()
    slide_type = _fallback_clean_text(slide_type, 'story').lower()

    if model_key == 'ventureos-editorial':
        return []
    if model_key in {'google-imagen', 'google-imagen-fast'}:
        return [model_key]
    if model_key == 'flux-studio':
        return ['flux']
    if model_key == 'qwen-editorial':
        return ['qwen-image']
    if model_key == 'gptimage-editorial':
        return ['gptimage']

    thematic_order = {
        'hook': ['flux', 'gptimage', 'qwen-image'],
        'problem': ['flux', 'qwen-image', 'gptimage'],
        'stakes': ['flux', 'qwen-image', 'gptimage'],
        'solution': ['flux', 'gptimage', 'qwen-image'],
        'how_it_works': ['qwen-image', 'flux', 'gptimage'],
        'impact': ['flux', 'gptimage', 'qwen-image'],
        'proof': ['flux', 'gptimage', 'qwen-image'],
        'business_model': ['qwen-image', 'flux', 'gptimage'],
        'vision': ['flux', 'gptimage', 'qwen-image'],
        'call_to_action': ['flux', 'gptimage', 'qwen-image'],
    }
    ordered = thematic_order.get(slide_type, ['flux', 'gptimage', 'qwen-image'])
    rotation = index % len(ordered)
    return ordered[rotation:] + ordered[:rotation]


def _fallback_provider_label(provider_model):
    return HOSTED_PROVIDER_LABELS.get(provider_model, FALLBACK_IMAGE_MODEL_LABEL)


def _fallback_provider_repo(provider_model):
    return HOSTED_PROVIDER_REPOS.get(provider_model, 'builtin/ventureos-neon-editorial')


def _google_image_model_name(provider_model):
    model_names = {
        'google-imagen': 'imagen-4.0-generate-001',
        'google-imagen-fast': 'imagen-4.0-fast-generate-001',
    }
    return model_names.get(provider_model, 'imagen-4.0-generate-001')


def _get_google_image_client():
    if not GOOGLE_IMAGE_ENABLED:
        raise RuntimeError('Google image generation is unavailable')

    global _GOOGLE_IMAGE_CLIENT
    with _GOOGLE_IMAGE_CLIENT_LOCK:
        if _GOOGLE_IMAGE_CLIENT is None:
            _GOOGLE_IMAGE_CLIENT = google_genai.Client(
                api_key=os.getenv('GOOGLE_API_KEY', '').strip()
            )
    return _GOOGLE_IMAGE_CLIENT


def _fallback_image_palette(style_key):
    palettes = {
        'deck-illustration': {
            'bg0': '#0B1020', 'bg1': '#1A2540', 'bg2': '#25355A',
            'accent': '#7DD3FC', 'accent2': '#F472B6', 'line': '#D8E6FF',
            'panel': '#14203A', 'panel2': '#1C2B4A'
        },
        'animated-scene': {
            'bg0': '#0D1326', 'bg1': '#1E2A52', 'bg2': '#2F3F77',
            'accent': '#60A5FA', 'accent2': '#A78BFA', 'line': '#E6F0FF',
            'panel': '#18254A', 'panel2': '#233463'
        },
        'cartoon': {
            'bg0': '#101827', 'bg1': '#22415F', 'bg2': '#395B7A',
            'accent': '#F59E0B', 'accent2': '#FB7185', 'line': '#FFF7ED',
            'panel': '#17314A', 'panel2': '#21415F'
        },
        'abstract': {
            'bg0': '#070915', 'bg1': '#0B1024', 'bg2': '#1A1033',
            'accent': '#FF8A3D', 'accent2': '#57C7FF', 'line': '#FFD9A1',
            'panel': '#12152F', 'panel2': '#1B2147'
        },
    }
    return palettes.get(style_key) or palettes['abstract']


def _fallback_style_family(style_key):
    style_key = _fallback_clean_text(style_key, 'deck-illustration').lower()
    if style_key in {'abstract', 'auto'}:
        return 'abstract'
    if style_key in {'deck-illustration', 'illustration', 'product-mockup', '3d-render'}:
        return 'presentation'
    if style_key == 'animated-scene':
        return 'animated'
    if style_key == 'cartoon':
        return 'cartoon'
    return 'presentation'


def _fallback_style_model_meta(style_key):
    family = _fallback_style_family(style_key)
    mapping = {
        'abstract': ('VentureOS Neon Editorial', 'builtin/ventureos-neon-editorial'),
        'presentation': ('VentureOS Presentation Illustration', 'builtin/ventureos-presentation-illustration'),
        'animated': ('VentureOS Animated Scene', 'builtin/ventureos-animated-scene'),
        'cartoon': ('VentureOS Cartoon Illustration', 'builtin/ventureos-cartoon-scene'),
    }
    return mapping.get(family, mapping['abstract'])


def _fallback_select_image_slide_indices(slides, image_options):
    slides = slides or []
    if not slides:
        return set()

    coverage = _fallback_clean_text((image_options or {}).get('coverage'), 'key-slides').lower()
    type_lookup = {}
    for index, slide in enumerate(slides):
        slide_type = _fallback_clean_text((slide or {}).get('type')).lower()
        if slide_type and slide_type not in type_lookup:
            type_lookup[slide_type] = index

    key_order = [
        type_lookup.get('hook', 0),
        type_lookup.get('problem', 1 if len(slides) > 1 else 0),
        type_lookup.get('solution', min(3, len(slides) - 1)),
        type_lookup.get('how_it_works', min(4, len(slides) - 1)),
        type_lookup.get('impact', min(5, len(slides) - 1)),
        type_lookup.get('proof', min(6, len(slides) - 1)),
        type_lookup.get('vision', min(len(slides) - 2, len(slides) - 1)),
        type_lookup.get('call_to_action', len(slides) - 1),
    ]

    ordered = []
    seen = set()
    for index in key_order:
        if index is None:
            continue
        safe_index = max(0, min(len(slides) - 1, index))
        if safe_index not in seen:
            seen.add(safe_index)
            ordered.append(safe_index)

    if coverage == 'all':
        return set(range(len(slides)))
    if coverage == 'image-heavy':
        return set(range(min(len(slides), max(7, len(ordered)))))
    if coverage == 'key-slides':
        return set(ordered[:5])
    return set(ordered[:2] or [0])


def _fallback_svg_data_url(svg_markup):
    encoded = base64.b64encode(svg_markup.encode('utf-8')).decode('ascii')
    return f'data:image/svg+xml;base64,{encoded}'


def _fallback_binary_data_url(binary_data, mime_type):
    encoded = base64.b64encode(binary_data).decode('ascii')
    return f'data:{mime_type};base64,{encoded}'


def _fallback_topic_bucket(idea, slide):
    slide = slide or {}
    corpus = ' '.join([
        _fallback_clean_text(idea),
        _fallback_clean_text(slide.get('title')),
        _fallback_clean_text(slide.get('subtitle')),
        _fallback_clean_text(slide.get('visual_suggestion')),
        _fallback_clean_text(slide.get('objective')),
    ]).lower()

    weighted_buckets = {
        'housing': {
            'housing': 6, 'sublease': 8, 'sublessee': 8, 'subleaser': 8, 'rent': 4, 'rental': 4, 'rentals': 4,
            'apartment': 6, 'apartments': 6, 'student housing': 9, 'lease': 5, 'leasing': 5, 'tenant': 5,
            'tenants': 5, 'roommate': 5, 'roommates': 5, 'landlord': 5, 'property': 4, 'real estate': 5, 'move-in': 5,
        },
        'health': {
            'health': 4, 'medical': 8, 'medicine': 8, 'med school': 8, 'medical school': 10, 'patient': 6, 'patients': 6,
            'clinic': 6, 'biotech': 6, 'dermatitis': 7, 'care': 4, 'diagnostic': 6, 'doctor': 6, 'nurse': 6,
            'hospital': 7, 'anatomy': 7, 'physiology': 7, 'clinical': 7,
        },
        'education': {
            'education': 5, 'learning': 5, 'student': 4, 'students': 4, 'school': 4, 'college': 4, 'course': 5,
            'training': 5, 'research': 4, 'tutor': 7, 'tutoring': 7, 'teacher': 5, 'classroom': 4, 'study': 5,
        },
        'gaming': {
            'game': 6, 'gaming': 6, 'indie': 5, 'developer': 5, 'dev': 4, 'studio': 4, 'creator': 4, 'player': 3, 'stream': 2, 'steam': 4,
        },
        'restaurant': {
            'restaurant': 7, 'kitchen': 5, 'food': 4, 'menu': 4, 'inventory': 6, 'dining': 5, 'hospitality': 5, 'grocery': 4, 'ingredients': 5, 'chef': 5,
        },
        'sports': {
            'sports': 6, 'football': 6, 'nfl': 7, 'team': 4, 'season': 4, 'stadium': 5, 'league': 4,
        },
        'finance': {
            'finance': 6, 'fintech': 7, 'payment': 6, 'payments': 6, 'billing': 5, 'revenue': 4, 'bank': 5, 'credit': 4, 'insurance': 5, 'subscription': 4,
        },
        'logistics': {
            'logistics': 7, 'warehouse': 6, 'supply': 5, 'shipping': 6, 'delivery': 6, 'fleet': 5, 'procurement': 5, 'routing': 5,
        },
        'software': {
            'ai': 3, 'agent': 4, 'software': 6, 'saas': 6, 'platform': 5, 'automation': 6, 'analytics': 5, 'dashboard': 5, 'workflow': 5, 'app': 3, 'tool': 3,
        },
    }
    medical_terms = ('medical', 'medicine', 'med school', 'medical school', 'clinical', 'anatomy', 'physiology', 'doctor', 'nurse', 'hospital')
    learning_terms = ('student', 'students', 'school', 'college', 'learning', 'tutor', 'tutoring', 'study', 'training')

    best_bucket = 'generic'
    best_score = 0
    for bucket, keywords in weighted_buckets.items():
        score = sum(weight for keyword, weight in keywords.items() if keyword in corpus)
        if bucket == 'health' and any(term in corpus for term in medical_terms) and any(term in corpus for term in learning_terms):
            score += 10
        if bucket == 'education' and 'student housing' in corpus:
            score -= 6
        if score > best_score:
            best_bucket = bucket
            best_score = score
    return best_bucket


def _premium_fallback_asset_candidates(topic, idea='', slide=None):
    topic = _fallback_clean_text(topic, 'generic').lower()
    candidates = list(PREMIUM_FALLBACK_LIBRARY.get(topic) or [])
    if topic not in {'generic'}:
        candidates.extend(PREMIUM_FALLBACK_LIBRARY.get('generic', []))

    focus_terms = _fallback_extract_focus_terms(idea, slide or {}, limit=8)
    focus_tokens = set()
    for term in focus_terms:
        focus_tokens.update(re.findall(r'[a-z][a-z0-9-]{1,}', _fallback_clean_text(term).lower()))
    focus_tokens.add(topic)
    focus_tokens.update({
        'housing': {'apartment', 'interior', 'student', 'lease', 'sublease', 'renter', 'residential'},
        'gaming': {'creator', 'studio', 'developer', 'dev', 'collectible', 'tool'},
        'software': {'workflow', 'dashboard', 'automation', 'platform', 'product'},
        'finance': {'payment', 'billing', 'ledger', 'market', 'exchange'},
        'restaurant': {'inventory', 'kitchen', 'dining', 'supply', 'operations'},
        'health': {'care', 'clinical', 'diagnostic', 'wellness', 'medical'},
        'education': {'campus', 'learning', 'research', 'student', 'academic'},
        'logistics': {'routing', 'warehouse', 'fleet', 'distribution', 'supply'},
        'sports': {'stadium', 'team', 'strategy', 'performance', 'athletic'},
    }.get(topic, set()))

    seen = set()
    ordered = []
    for item in candidates:
        if item and item not in seen:
            seen.add(item)
            ordered.append(item)

    def asset_score(asset_name):
        tags = set(PREMIUM_FALLBACK_ASSET_TAGS.get(asset_name, []))
        score = 0
        if topic in tags:
            score += 12
        score += sum(3 for token in focus_tokens if token in tags)
        if topic in {'housing', 'education', 'restaurant', 'health'} and 'interior' in tags:
            score += 2
        if topic in {'software', 'finance', 'gaming'} and {'dashboard', 'workflow', 'studio', 'creator'} & tags:
            score += 2
        return score

    return [
        asset_name
        for _, _, asset_name in sorted(
            [(asset_score(asset_name), idx, asset_name) for idx, asset_name in enumerate(ordered)],
            key=lambda item: (-item[0], item[1])
        )
    ]


def _premium_fallback_asset_path(asset_name):
    if not asset_name:
        return ''
    return os.path.join(app.root_path, 'static', 'premium_fallback', asset_name)


def _premium_fallback_variant_from_style(style_key):
    family = _fallback_style_family(style_key)
    if family == 'animated':
        return 'animated'
    if family == 'cartoon':
        return 'cartoon'
    return 'presentation'


def _premium_fallback_image_meta(idea, slide, index, image_options=None):
    slide = slide or {}
    style_key = _fallback_clean_text((image_options or {}).get('style'), 'deck-illustration').lower()
    topic = _fallback_topic_bucket(idea, slide)
    if topic not in PREMIUM_FALLBACK_DIRECT_TOPICS:
        return None
    slide_type = _fallback_clean_text(slide.get('type'), 'story').lower()
    asset_candidates = _premium_fallback_asset_candidates(topic, idea=idea, slide=slide)
    if not asset_candidates:
        return None

    asset_name = asset_candidates[index % len(asset_candidates)]
    asset_path = _premium_fallback_asset_path(asset_name)
    if not asset_path or not os.path.exists(asset_path):
        return None

    palette = _fallback_image_palette(style_key)
    variant = _premium_fallback_variant_from_style(style_key)
    focus_x, focus_y = _fallback_scene_focus_for_slide_type(slide_type)
    with Image.open(asset_path) as source_image:
        scene = source_image.convert('RGB')
        scene = _fallback_apply_scene_treatment(scene, palette, variant)
        scene = _fallback_resize_cover(scene, (1024, 1024), focus_x=focus_x, focus_y=focus_y)

    subject_brief = _fallback_slide_subject_brief(topic, slide_type, index=index)
    return {
        'image_url': _fallback_image_to_data_url(scene, fmt='JPEG', quality=92),
        'image_prompt': subject_brief,
        'image_model': PREMIUM_FALLBACK_MODEL_LABEL,
        'image_repo_id': PREMIUM_FALLBACK_MODEL_REPO,
        'image_status': 'generated',
        'image_error': '',
    }


def _best_available_fallback_image_meta(idea, slide, index, image_options=None, hosted_error=''):
    premium_meta = _premium_fallback_image_meta(idea, slide, index, image_options=image_options)
    if premium_meta:
        if hosted_error:
            premium_meta['image_error'] = hosted_error
        return premium_meta
    return _fallback_generate_vector_image_meta(
        idea=idea,
        slide=slide,
        index=index,
        image_options=image_options,
        hosted_error=hosted_error,
    )


def _fallback_topic_subject(topic):
    subjects = {
        'gaming': 'indie game creator studio, premium collectibles, glowing dev setup, atmospheric shelves, cinematic environment art',
        'restaurant': 'refined restaurant operations scene, chef pass, inventory shelves, premium hospitality workspace, cinematic still life',
        'housing': 'student housing search environment, elevated apartment interiors, urban architecture, leasing journey concept art',
        'health': 'premium medical learning and healthcare concept scene, anatomy study cues, clinical tutoring atmosphere, thoughtful editorial composition',
        'sports': 'dramatic sports operations environment, strategy room, stadium energy, premium athletic editorial scene',
        'finance': 'premium fintech operations scene, sophisticated financial workspace, screens and objects, editorial mood',
        'education': 'modern learning studio, thoughtful academic environment, premium educational editorial scene',
        'logistics': 'warehouse command center, supply chain operations environment, premium industrial editorial scene',
        'software': 'product strategy workspace, premium startup studio, modern software operations environment',
        'generic': 'premium editorial concept art scene for a startup pitch deck, atmospheric and polished',
    }
    return subjects.get(topic, subjects['generic'])


def _fallback_topic_abstract_subject(topic):
    subjects = {
        'gaming': 'stylized indie game studio editorial scene with creator energy, premium neon rim light, and polished synthwave mood',
        'restaurant': 'stylized restaurant operations editorial scene with supply cues, premium neon accents, and cinematic rhythm',
        'housing': 'stylized student housing editorial scene with modern apartment architecture, search cues, and premium neon wayfinding light',
        'health': 'stylized medical learning editorial scene with anatomy cues, clinical confidence, and premium neon highlights',
        'sports': 'stylized sports strategy editorial scene with route energy, performance atmosphere, and premium neon accents',
        'finance': 'stylized fintech editorial scene with exchange energy, elegant market cues, and premium neon highlights',
        'education': 'stylized learning editorial scene with research cues, academic aspiration, and premium neon accents',
        'logistics': 'stylized logistics editorial scene with routing energy, distribution cues, and premium neon highlights',
        'software': 'stylized product-tech editorial scene with workflow cards, intelligent automation cues, and premium neon accents',
        'generic': 'premium neon editorial scene with one clear subject, atmospheric depth, and polished presentation energy',
    }
    return subjects.get(topic, subjects['generic'])


_FALLBACK_VISUAL_STOPWORDS = {
    'about', 'across', 'after', 'agent', 'agents', 'analysis', 'because', 'brand', 'business', 'call',
    'chart', 'company', 'customer', 'customers', 'deck', 'for', 'future', 'growth', 'idea', 'impact', 'include',
    'investor', 'key', 'market', 'model', 'more', 'need', 'next', 'offer', 'pitch', 'platform', 'problem',
    'product', 'proof', 'reliable', 'revenue', 'risk', 'roadmap', 'show', 'slide', 'slides', 'solution',
    'startup', 'story', 'system', 'team', 'their', 'these', 'this', 'title', 'traction', 'value', 'vision',
    'what', 'with', 'your', 'users'
}

_FALLBACK_VISUAL_MOTIFS = {
    'gaming': [
        'controller silhouettes', 'concept art panels', 'creator workbench', 'asset shelves', 'code monitors', 'dev tool glow'
    ],
    'restaurant': [
        'ingredient crates', 'scanner pulses', 'prep surfaces', 'supply shelves', 'menu-card silhouettes', 'service counters'
    ],
    'housing': [
        'apartment facades', 'map pins', 'keys', 'room cards', 'leasing routes', 'urban residential interiors'
    ],
        'health': [
            'care glyphs', 'pulse rings', 'diagnostic panels', 'clinic interiors', 'anatomy study cues', 'medical device silhouettes'
        ],
    'sports': [
        'stadium arcs', 'play routes', 'training dashboards', 'team strategy boards', 'score pulses', 'athletic path lines'
    ],
    'finance': [
        'payment cards', 'ledger routes', 'market graphs', 'wallet tokens', 'trading panels', 'exchange signals'
    ],
    'education': [
        'campus silhouettes', 'learning cards', 'book forms', 'idea bulbs', 'student pathways', 'research markers'
    ],
    'logistics': [
        'route maps', 'package nodes', 'warehouse racks', 'tracking lines', 'fleet silhouettes', 'distribution hubs'
    ],
    'software': [
        'workflow cards', 'automation routes', 'product panels', 'node graphs', 'data tiles', 'clean UI blocks'
    ],
    'generic': [
        'premium interface cards', 'signal paths', 'editorial depth', 'clean product silhouettes', 'system nodes', 'ambient glow'
    ],
}

_FALLBACK_SLIDE_SUBJECTS = {
    'hook': {
        'housing': 'student housing discovery scene with modern apartment architecture and trusted search cues',
        'gaming': 'indie game creator studio scene with polished dev energy and collectible craft',
        'restaurant': 'restaurant operations command scene with premium hospitality rhythm and inventory control',
        'health': 'medical learning mission scene with clinical confidence and guided study atmosphere',
        'sports': 'sports intelligence scene with strategy energy and performance momentum',
        'finance': 'fintech mission scene with premium transaction energy and market confidence',
        'education': 'learning platform mission scene with aspirational academic energy',
        'logistics': 'logistics mission scene with global routing energy and operational precision',
        'software': 'software platform mission scene with product confidence and intelligent automation',
        'generic': 'premium startup mission scene with a clear modern focal subject',
    },
    'problem': {
        'housing': 'student renters facing fragmented apartment search and unreliable sublease discovery',
        'gaming': 'indie creators struggling with scattered tools and limited distribution support',
        'restaurant': 'restaurant teams navigating stock gaps and operational friction',
        'health': 'medical learners facing fragmented study support and low-confidence mastery',
        'sports': 'teams dealing with fragmented decision data and unclear performance signals',
        'finance': 'finance operators facing fragmented ledgers and trust friction',
        'education': 'learners facing fragmented guidance and low clarity',
        'logistics': 'operators navigating broken routing visibility and supply friction',
        'software': 'teams facing disconnected systems and manual workflow friction',
        'generic': 'teams facing fragmented workflows and avoidable operational friction',
    },
    'stakes': {
        'generic': 'the scale and urgency of the market opportunity with strong outcome pressure'
    },
    'solution': {
        'housing': 'verified housing-matching product scene with trusted listings and guided move-in flow',
        'gaming': 'creator platform scene with curated resources, launch tools, and premium support',
        'restaurant': 'inventory intelligence scene with cleaner kitchen operations and decision support',
        'health': 'medical tutoring scene with guided mastery, anatomy cues, and confident study support',
        'sports': 'sports intelligence scene with clearer coaching signals and confident decisions',
        'finance': 'fintech solution scene with smoother flows and confident transaction design',
        'education': 'learning solution scene with guided discovery and progress clarity',
        'logistics': 'supply-chain control scene with clean routing and tracking precision',
        'software': 'automation platform scene with elegant workflow orchestration and product clarity',
        'generic': 'premium solution scene with a confident product-led outcome',
    },
    'how_it_works': {
        'generic': 'connected workflow scene showing how the system moves from input to outcome'
    },
    'impact': {
        'generic': 'outcome scene showing measurable progress, adoption, and value creation'
    },
    'proof': {
        'generic': 'credible traction scene with demand signals, usage proof, and trustworthy momentum'
    },
    'business_model': {
        'generic': 'commercial system scene with pricing layers, exchange flows, and monetization clarity'
    },
    'vision': {
        'generic': 'future-state scene with aspirational scale, elegance, and category leadership'
    },
    'call_to_action': {
        'generic': 'closing brand scene with a focused emblem, decisive energy, and premium finality'
    },
}


def _fallback_extract_focus_terms(idea, slide, limit=5):
    slide = slide or {}
    corpus = ' '.join([
        _fallback_clean_text(idea),
        _fallback_clean_text(slide.get('title')),
        _fallback_clean_text(slide.get('subtitle')),
        _fallback_clean_text(slide.get('objective')),
        _fallback_clean_text(slide.get('visual_suggestion')),
        ' '.join(_fallback_clean_text(item) for item in (slide.get('content') or [])),
    ]).lower()
    tokens = re.findall(r'[a-z][a-z0-9-]{2,}', corpus)
    terms = []
    seen = set()
    for token in tokens:
        normalized = token.replace('-', ' ')
        if normalized in seen or normalized in _FALLBACK_VISUAL_STOPWORDS:
            continue
        if normalized.endswith('s') and len(normalized) > 4:
            singular = normalized[:-1]
            if singular not in seen and singular not in _FALLBACK_VISUAL_STOPWORDS:
                normalized = singular
        if normalized in seen or normalized in _FALLBACK_VISUAL_STOPWORDS:
            continue
        seen.add(normalized)
        terms.append(normalized)
        if len(terms) >= limit:
            break
    return terms


def _fallback_topic_motifs(topic, index=0, limit=3):
    motifs = _FALLBACK_VISUAL_MOTIFS.get(topic) or _FALLBACK_VISUAL_MOTIFS['generic']
    if not motifs:
        return []
    rotation = index % len(motifs)
    ordered = motifs[rotation:] + motifs[:rotation]
    return ordered[:limit]


def _fallback_slide_subject_brief(topic, slide_type, index=0):
    slide_type = _fallback_clean_text(slide_type, 'story').lower()
    topic_subjects = _FALLBACK_SLIDE_SUBJECTS.get(slide_type) or {}
    subject = topic_subjects.get(topic) or topic_subjects.get('generic')
    if subject:
        return subject

    generic_subjects = {
        'hook': 'premium mission-driven startup scene with one dominant focal subject',
        'problem': 'friction-filled scene that makes the core problem instantly legible',
        'stakes': 'market-scale scene that communicates urgency and consequence',
        'solution': 'clean product-world scene that feels organized and trustworthy',
        'how_it_works': 'connected system scene with modular steps and clear flow',
        'impact': 'success-state scene with outcome-driven momentum and value',
        'proof': 'credible traction scene with anchored evidence and trust',
        'business_model': 'commercial system scene with monetization layers and routes',
        'vision': 'future-state category scene with scale and aspiration',
        'call_to_action': 'closing brand scene with decisive energy and confidence',
    }
    return generic_subjects.get(slide_type, 'premium topic-faithful scene with a clear central subject')


def _presentation_visual_summary(idea, slide, index=0):
    slide = slide or {}
    topic = _fallback_topic_bucket(idea, slide)
    slide_type = _fallback_clean_text(slide.get('type'), 'story').lower()
    summaries = {
        'gaming': {
            'hook': 'Creator studio hero scene',
            'problem': 'Scattered tooling friction scene',
            'solution': 'Curated creator platform scene',
            'how_it_works': 'Creator workflow scene',
            'impact': 'Momentum and adoption scene',
            'proof': 'Traction proof scene',
            'generic': 'Topic-faithful gaming scene',
        },
        'housing': {
            'hook': 'Student housing discovery scene',
            'problem': 'Apartment search friction scene',
            'solution': 'Trusted listings product scene',
            'how_it_works': 'Rental workflow scene',
            'impact': 'Move-in outcome scene',
            'proof': 'Marketplace traction scene',
            'generic': 'Topic-faithful housing scene',
        },
        'software': {
            'hook': 'Product mission scene',
            'problem': 'Workflow bottleneck scene',
            'solution': 'Automation platform scene',
            'how_it_works': 'Connected workflow scene',
            'impact': 'Product value scene',
            'proof': 'Adoption proof scene',
            'generic': 'Topic-faithful software scene',
        },
        'restaurant': {
            'hook': 'Operations control scene',
            'problem': 'Inventory friction scene',
            'solution': 'Kitchen intelligence scene',
            'how_it_works': 'Service workflow scene',
            'generic': 'Topic-faithful restaurant scene',
        },
        'health': {
            'hook': 'Care mission scene',
            'problem': 'Signal fragmentation scene',
            'solution': 'Clinical decision scene',
            'generic': 'Topic-faithful health scene',
        },
        'finance': {
            'hook': 'Fintech mission scene',
            'problem': 'Payment friction scene',
            'solution': 'Transaction flow scene',
            'generic': 'Topic-faithful fintech scene',
        },
        'logistics': {
            'hook': 'Routing mission scene',
            'problem': 'Supply friction scene',
            'solution': 'Control tower scene',
            'generic': 'Topic-faithful logistics scene',
        },
        'education': {
            'hook': 'Learning mission scene',
            'problem': 'Guidance gap scene',
            'solution': 'Learning platform scene',
            'generic': 'Topic-faithful education scene',
        },
        'sports': {
            'hook': 'Performance mission scene',
            'problem': 'Decision friction scene',
            'solution': 'Strategy intelligence scene',
            'generic': 'Topic-faithful sports scene',
        },
        'generic': {
            'hook': 'Mission-led hero scene',
            'problem': 'Core friction scene',
            'stakes': 'Market opportunity scene',
            'solution': 'Product answer scene',
            'how_it_works': 'Connected workflow scene',
            'impact': 'Outcome-led scene',
            'proof': 'Proof and traction scene',
            'business_model': 'Commercial system scene',
            'vision': 'Future-state scene',
            'call_to_action': 'Closing brand scene',
            'generic': 'Topic-faithful premium scene',
        },
    }
    topic_summaries = summaries.get(topic) or summaries['generic']
    return topic_summaries.get(slide_type) or topic_summaries.get('generic') or summaries['generic']['generic']


def _fallback_style_prompt(style_key):
    prompts = {
        'deck-illustration': 'premium scene illustration, cinematic environment art, polished keynote realism, topic-faithful composition, clean atmospheric storytelling',
        'animated-scene': 'premium stylized scene illustration, animated feature-film lighting, cinematic environment art, scene-first composition, richly lit and detailed',
        'cartoon': 'premium stylized cartoon scene, cel-shaded environment art, clean shapes, rich lighting, polished and detailed',
        'abstract': 'high-end neon editorial illustration, dark premium background, electric orange and cyan rim light, topic-faithful subject, poster-grade clarity, premium presentation artwork',
    }
    return prompts.get(style_key, prompts['deck-illustration'])


def _fallback_slide_direction(slide_type, style_key=None):
    style_key = _fallback_clean_text(style_key).lower()
    if style_key == 'abstract':
        directions = {
            'hook': 'hero poster with one unmistakable subject, strong negative space, and premium neon drama',
            'problem': 'express friction through contrast, obstruction, broken flow, or mismatched elements in a cohesive scene',
            'stakes': 'show urgency with layered depth, signal intensity, and visible consequence',
            'solution': 'resolve into a cleaner product-world scene with calmer geometry and confident luminous structure',
            'how_it_works': 'show a connected step-by-step system using cards, nodes, routes, or staged objects',
            'impact': 'express progress through upward energy, success-state symbols, and brighter outcome cues',
            'proof': 'credible proof scene with refined data cues, anchored signals, and grounded editorial trust',
            'business_model': 'commercial energy through pricing layers, exchange cues, and premium modular structure',
            'vision': 'aspirational future-state poster with scale, atmosphere, and elegant premium ambition',
            'call_to_action': 'decisive closing poster with a focused emblem, radiant center, and memorable finish',
        }
        return directions.get(slide_type, 'premium abstract composition with strong hierarchy, dramatic depth, and elegant restraint')
    directions = {
        'hook': 'one dominant hero subject, strong negative space, memorable cover composition',
        'problem': 'show friction and tension through environment and objects, not infographic shapes',
        'stakes': 'high-importance atmosphere, urgency and scale, premium visual storytelling',
        'solution': 'organized and optimistic product-world scene, premium clarity and confidence',
        'how_it_works': 'structured operational scene with layers and flow, premium system storytelling',
        'impact': 'outcome-led editorial scene with proof-driven atmosphere, not charts or labels',
        'proof': 'credible traction mood with refined objects, screens, and environment, no visible text',
        'business_model': 'commercial scene with polished objects and transactional energy, premium still life',
        'vision': 'aspirational future-state environment, elegant and premium',
        'call_to_action': 'confident closing visual, polished and cinematic with decisive mood',
    }
    return directions.get(slide_type, 'premium editorial composition, strong focal hierarchy, visually rich but clean')


def _fallback_prompt_composition(index, slide_type, style_key=None):
    slide_type = _fallback_clean_text(slide_type, 'story').lower()
    style_key = _fallback_clean_text(style_key).lower()
    if style_key == 'abstract':
        type_cues = {
            'hook': 'dominant hero subject, subtle interface cues, luxury poster balance, and clear storytelling focus',
            'problem': 'tension-filled composition, asymmetrical balance, broken flow, and premium moody atmosphere',
            'stakes': 'layered depth, expanding scale cues, dramatic signal intensity, and poster-grade focus',
            'solution': 'ordered product-world scene, clean central structure, optimistic orange and cyan glow',
            'how_it_works': 'connected cards, nodes, route lines, or staged objects with clear flow',
            'impact': 'ascending energy, pulse rings, brighter outcome cues, and a strong focal symbol',
            'proof': 'evidence-driven composition, trust cues, refined panels, and anchored momentum',
            'business_model': 'exchange cards, pricing cues, routing connectors, and premium modular order',
            'vision': 'wide future-state poster with layered horizon glow and elegant ambition',
            'call_to_action': 'focused closing emblem, radiant center, and premium final cadence',
        }
        rotation_cues = [
            'editorial poster balance with one dominant luminous anchor',
            'slightly elevated composition with layered circuit planes',
            'cinematic side glow with floating interface elements',
            'gallery-like composition with strong negative space and spark details',
            'quiet premium composition with restrained signal trails and atmospheric depth',
        ]
        base = type_cues.get(slide_type, 'premium abstract composition with clean hierarchy and atmospheric depth')
        return f'{base}, {rotation_cues[index % len(rotation_cues)]}'
    type_cues = {
        'hook': 'wide establishing shot, strong silhouette subject, architectural depth, luxury editorial mood',
        'problem': 'off-axis composition, tactile objects, narrative friction, cinematic tension',
        'stakes': 'dramatic depth, layered environment storytelling, elevated scale cues',
        'solution': 'organized hero composition, cleaner geometry, optimistic lighting, product-world clarity',
        'how_it_works': 'multi-zone operational composition, process depth, clean pathways, premium systems feel',
        'impact': 'success-state composition, premium environmental storytelling, confident lighting',
        'proof': 'credible editorial still life, refined screens and objects, grounded realism without readable text',
        'business_model': 'commercial still life composition, refined money-flow cues, premium materials',
        'vision': 'future-facing composition, expansive horizon, aspirational architecture, elegant atmosphere',
        'call_to_action': 'closing campaign image, decisive focal point, polished premium finish',
    }
    rotation_cues = [
        'camera at eye level with soft depth of field',
        'slightly elevated perspective with generous foreground shape language',
        'cinematic side light and asymmetrical balance',
        'gallery-like composition with premium restraint and strong negative space',
        'editorial magazine composition with layered depth and subtle atmosphere',
    ]
    base = type_cues.get(slide_type, 'premium editorial composition with clean hierarchy and atmospheric depth')
    return f'{base}, {rotation_cues[index % len(rotation_cues)]}'


def _fallback_hosted_image_size(slide, index, image_options=None):
    slide_type = _fallback_clean_text((slide or {}).get('type'), 'story').lower()
    coverage = _fallback_clean_text((image_options or {}).get('coverage'), 'key-slides').lower()
    if slide_type in {'hook', 'vision', 'call_to_action'}:
        return 768
    if coverage == 'all':
        return 512
    if slide_type in {'solution', 'how_it_works', 'impact', 'proof'}:
        return 640
    return 512


def _fallback_hosted_image_prompt(idea, slide, style_key, index=0):
    slide = slide or {}
    topic = _fallback_topic_bucket(idea, slide)
    title = _fallback_clean_text(slide.get('title'))
    visual = _fallback_clean_text(slide.get('visual_suggestion'))
    subtitle = _fallback_clean_text(slide.get('subtitle'))
    slide_type = _fallback_clean_text(slide.get('type'), 'story').lower()
    abstract_mode = style_key == 'abstract'
    focus_terms = _fallback_extract_focus_terms(idea, slide, limit=5)
    motifs = _fallback_topic_motifs(topic, index=index, limit=3)
    subject_brief = _fallback_slide_subject_brief(topic, slide_type, index=index)

    prompt_parts = [
        f'topic-faithful premium editorial illustration of {subject_brief}',
        _fallback_topic_abstract_subject(topic) if abstract_mode else _fallback_topic_subject(topic),
        _fallback_style_prompt(style_key),
        _fallback_slide_direction(slide_type, style_key),
        _fallback_prompt_composition(index, slide_type, style_key),
    ]
    if motifs:
        prompt_parts.append(f'include motifs such as {", ".join(motifs)}')
    if focus_terms:
        prompt_parts.append(f'topic cues: {", ".join(focus_terms)}')
    if title:
        prompt_parts.append(f'theme inspired by "{title}"')
    if visual:
        prompt_parts.append(visual)
    elif subtitle:
        prompt_parts.append(subtitle)
    if idea:
        prompt_parts.append(f'for startup idea "{_fallback_clean_text(idea)[:120]}"')

    prompt_parts.extend([
        'square composition for a premium presentation image panel',
        'no text, no letters, no words, no watermark, no logo, no captions, no readable UI text',
        'no signature, no footer text, no bottom strip, no edge labels, no stray glyphs, no tiny text artifacts',
        'no signs, no posters, no billboards, no labels, no typographic panels',
        'if interface elements appear, keep them icon-only and textless with clean graphic blocks',
        'not a generic infographic, not a plain chart, not a literal dashboard screenshot',
        'prefer one cohesive scene instead of a collage of unrelated panels',
        'make the startup topic immediately obvious and avoid random unrelated symbols',
    ])
    if abstract_mode:
        prompt_parts.extend([
            'allow symbolic topic objects, dashboards, gears, icons, city forms, and industry motifs when relevant',
            'keep it stylistic and poster-like, but still grounded in the actual topic and startup domain',
        ])
    return ', '.join(part for part in prompt_parts if part)


def _fallback_fetch_hosted_image_meta(idea, slide, index, image_options=None):
    if not HOSTED_FALLBACK_ENABLED:
        raise RuntimeError('Hosted fallback image generation is disabled')

    slide = slide or {}
    style_key = _fallback_clean_text((image_options or {}).get('style'), 'deck-illustration').lower()
    model_key = _fallback_clean_text((image_options or {}).get('model_key') or (image_options or {}).get('model'), FALLBACK_IMAGE_MODEL_KEY)
    provider_model = _fallback_provider_candidates(model_key, slide.get('type'), index)[0]
    image_size = _fallback_hosted_image_size(slide, index, image_options=image_options)
    seed = _fallback_hash_seed(idea, slide.get('title'), slide.get('subtitle'), slide.get('type'), style_key, index)
    prompt = _fallback_hosted_image_prompt(idea, slide, style_key, index=index)
    cache_key = f"{provider_model}|{image_size}|{seed}|{prompt}"

    with _HOSTED_IMAGE_CACHE_LOCK:
        cached = _HOSTED_IMAGE_CACHE.get(cache_key)
    if cached:
        return {
            'image_url': cached,
            'image_prompt': prompt,
            'image_model': _fallback_provider_label(provider_model),
            'image_repo_id': _fallback_provider_repo(provider_model),
            'image_status': 'cached',
        }

    data_url = _fallback_fetch_provider_image_data_url(
        prompt=prompt,
        provider_model=provider_model,
        seed=seed,
        image_size=image_size,
        cache_key=cache_key,
    )

    return {
        'image_url': data_url,
        'image_prompt': prompt,
        'image_model': _fallback_provider_label(provider_model),
        'image_repo_id': _fallback_provider_repo(provider_model),
        'image_status': 'generated',
    }


def _fallback_fetch_google_image_data_url(prompt, provider_model, seed, image_size, cache_key=None):
    if not GOOGLE_IMAGE_ENABLED:
        raise RuntimeError('Google image generation is unavailable')
    prompt = _fallback_clean_text(prompt)
    provider_model = _fallback_clean_text(provider_model, 'google-imagen')
    cache_key = cache_key or f"{provider_model}|{image_size}|{seed}|{prompt}"

    with _HOSTED_IMAGE_CACHE_LOCK:
        cached = _HOSTED_IMAGE_CACHE.get(cache_key)
    if cached:
        return cached

    client = _get_google_image_client()
    model_name = _google_image_model_name(provider_model)
    config = google_genai_types.GenerateImagesConfig(
        number_of_images=1,
        aspect_ratio='1:1',
        person_generation='allow_adult',
    )
    response = client.models.generate_images(
        model=model_name,
        prompt=prompt,
        config=config,
    )
    generated_images = getattr(response, 'generated_images', None) or []
    if not generated_images:
        raise RuntimeError('Google image response returned no images')

    image_obj = getattr(generated_images[0], 'image', None)
    if image_obj is None:
        raise RuntimeError('Google image response was missing image payload')

    image_bytes = getattr(image_obj, 'image_bytes', None)
    if image_bytes is None:
        raise RuntimeError('Google image payload was empty')
    if isinstance(image_bytes, str):
        binary_data = base64.b64decode(image_bytes)
    else:
        binary_data = bytes(image_bytes)
    mime_type = getattr(image_obj, 'mime_type', None) or 'image/png'
    data_url = _fallback_binary_data_url(binary_data, mime_type)

    with _HOSTED_IMAGE_CACHE_LOCK:
        _HOSTED_IMAGE_CACHE[cache_key] = data_url
    return data_url


def _fallback_fetch_hosted_image_data_url(prompt, provider_model, seed, image_size, cache_key=None):
    if not HOSTED_FALLBACK_ENABLED:
        raise RuntimeError('Hosted fallback image generation is disabled')
    prompt = _fallback_clean_text(prompt)
    provider_model = _fallback_clean_text(provider_model, 'flux')
    seed = int(seed or 1)
    image_size = max(384, int(image_size or 512))
    cache_key = cache_key or f"{provider_model}|{image_size}|{seed}|{prompt}"

    with _HOSTED_IMAGE_CACHE_LOCK:
        cached = _HOSTED_IMAGE_CACHE.get(cache_key)
    if cached:
        return cached

    request_url = (
        'https://image.pollinations.ai/prompt/'
        f'{quote(prompt)}?width={image_size}&height={image_size}&seed={seed}&model={provider_model}&nologo=true&safe=true'
    )
    request = Request(
        request_url,
        headers={
            'User-Agent': 'VentureOS/1.0',
            'Accept': 'image/*',
        }
    )
    with urlopen(request, timeout=HOSTED_FALLBACK_TIMEOUT_SECONDS) as response:
        mime_type = response.headers.get_content_type() or 'image/jpeg'
        if not mime_type.startswith('image/'):
            raise RuntimeError(f'Unexpected hosted image response type: {mime_type}')
        binary_data = response.read()
        if len(binary_data) < 4096:
            raise RuntimeError('Hosted image response was unexpectedly small')

    data_url = _fallback_binary_data_url(binary_data, mime_type)
    with _HOSTED_IMAGE_CACHE_LOCK:
        _HOSTED_IMAGE_CACHE[cache_key] = data_url

    return data_url


def _fallback_fetch_provider_image_data_url(prompt, provider_model, seed, image_size, cache_key=None):
    provider_model = _fallback_clean_text(provider_model, 'flux')
    if provider_model.startswith('google-imagen') and GOOGLE_IMAGE_ENABLED:
        return _fallback_fetch_google_image_data_url(
            prompt=prompt,
            provider_model=provider_model,
            seed=seed,
            image_size=image_size,
            cache_key=cache_key,
        )
    if provider_model.startswith('google-imagen'):
        provider_model = 'flux'
    return _fallback_fetch_hosted_image_data_url(
        prompt=prompt,
        provider_model=provider_model,
        seed=seed,
        image_size=image_size,
        cache_key=cache_key,
    )


def _fallback_image_from_data_url(data_url):
    if not data_url or ',' not in data_url:
        raise ValueError('Invalid image data URL')
    header, raw = data_url.split(',', 1)
    if ';base64' not in header:
        raise ValueError('Unsupported data URL encoding')
    return Image.open(BytesIO(base64.b64decode(raw))).convert('RGB')


def _fallback_image_to_data_url(image, fmt='JPEG', quality=90):
    buffer = BytesIO()
    save_format = 'JPEG' if fmt.upper() == 'JPEG' else 'PNG'
    save_kwargs = {'format': save_format}
    mime = 'image/jpeg' if save_format == 'JPEG' else 'image/png'
    if save_format == 'JPEG':
        if image.mode != 'RGB':
            image = image.convert('RGB')
        save_kwargs.update({'quality': quality, 'optimize': True, 'progressive': True})
    image.save(buffer, **save_kwargs)
    return _fallback_binary_data_url(buffer.getvalue(), mime)


def _fallback_hex_to_rgba(hex_color, alpha=255):
    hex_color = _fallback_clean_text(hex_color, '#FFFFFF')
    if not hex_color.startswith('#') or len(hex_color) != 7:
        hex_color = '#FFFFFF'
    return tuple(int(hex_color[i:i+2], 16) for i in (1, 3, 5)) + (alpha,)


def _fallback_resize_cover(image, size, focus_x=0.5, focus_y=0.5):
    target_w, target_h = size
    if target_w <= 0 or target_h <= 0:
        raise ValueError('Target size must be positive')

    src_w, src_h = image.size
    scale = max(target_w / src_w, target_h / src_h)
    resized = image.resize(
        (max(int(src_w * scale), 1), max(int(src_h * scale), 1)),
        Image.LANCZOS,
    )
    res_w, res_h = resized.size
    left = int(max(0, min(res_w - target_w, (res_w - target_w) * focus_x)))
    top = int(max(0, min(res_h - target_h, (res_h - target_h) * focus_y)))
    return resized.crop((left, top, left + target_w, top + target_h))


def _fallback_make_rounded_card(image, size, radius, focus_x=0.5, focus_y=0.5, border=None):
    card = _fallback_resize_cover(image, size, focus_x=focus_x, focus_y=focus_y).convert('RGBA')
    mask = Image.new('L', size, 0)
    mask_draw = ImageDraw.Draw(mask)
    mask_draw.rounded_rectangle((0, 0, size[0] - 1, size[1] - 1), radius=radius, fill=255)
    card.putalpha(mask)

    if border:
        border_overlay = Image.new('RGBA', size, (0, 0, 0, 0))
        border_draw = ImageDraw.Draw(border_overlay)
        border_draw.rounded_rectangle(
            (1, 1, size[0] - 2, size[1] - 2),
            radius=max(radius - 1, 0),
            outline=border,
            width=2,
        )
        card = Image.alpha_composite(card, border_overlay)

    return card


def _fallback_paste_card(canvas, card, position, shadow_alpha=52, shadow_offset=(0, 18)):
    px, py = position
    shadow = Image.new('RGBA', canvas.size, (0, 0, 0, 0))
    shadow_mask = Image.new('L', card.size, 0)
    shadow_draw = ImageDraw.Draw(shadow_mask)
    shadow_draw.rounded_rectangle(
        (0, 0, card.size[0] - 1, card.size[1] - 1),
        radius=36,
        fill=shadow_alpha,
    )
    shadow_mask = shadow_mask.filter(ImageFilter.GaussianBlur(radius=22))
    shadow.paste(
        Image.new('RGBA', card.size, (0, 0, 0, 255)),
        (px + shadow_offset[0], py + shadow_offset[1]),
        shadow_mask,
    )
    canvas.alpha_composite(shadow)
    canvas.alpha_composite(card, dest=(px, py))


def _fallback_add_glow(canvas, center, radius, color, alpha=90, blur=56):
    glow = Image.new('RGBA', canvas.size, (0, 0, 0, 0))
    glow_draw = ImageDraw.Draw(glow)
    cx, cy = center
    glow_draw.ellipse(
        (cx - radius, cy - radius, cx + radius, cy + radius),
        fill=_fallback_hex_to_rgba(color, alpha),
    )
    glow = glow.filter(ImageFilter.GaussianBlur(radius=blur))
    canvas.alpha_composite(glow)


def _fallback_anchor_group_for_slide(slide, index):
    slide_type = _fallback_clean_text((slide or {}).get('type')).lower()
    groups = {
        'hook': 'hero',
        'problem': 'tension',
        'stakes': 'tension',
        'solution': 'product',
        'how_it_works': 'process',
        'impact': 'traction',
        'proof': 'traction',
        'business_model': 'traction',
        'vision': 'future',
        'call_to_action': 'future',
    }
    if slide_type in groups:
        return groups[slide_type]
    return ['hero', 'tension', 'product', 'traction', 'future'][index % 5]


def _fallback_compose_remixed_image(source_image, slide, index, image_options=None):
    slide = slide or {}
    slide_type = _fallback_clean_text(slide.get('type'), 'story').lower()
    group_key = _fallback_anchor_group_for_slide(slide, index)
    style_key = _fallback_clean_text((image_options or {}).get('style'), 'deck-illustration').lower()
    palette = _fallback_image_palette(style_key)
    line_color = _fallback_hex_to_rgba(palette['line'], 104)
    accent_color = _fallback_hex_to_rgba(palette['accent'], 180)
    accent_two_color = _fallback_hex_to_rgba(palette['accent2'], 148)
    panel_fill = _fallback_hex_to_rgba(palette['panel'], 204)
    panel_soft = _fallback_hex_to_rgba(palette['panel2'], 170)
    border_color = _fallback_hex_to_rgba(palette['line'], 54)
    canvas_size = (768, 768)

    focus_map = {
        'hero': (0.56, 0.38),
        'tension': (0.30, 0.42),
        'product': (0.64, 0.44),
        'process': (0.48, 0.52),
        'traction': (0.66, 0.32),
        'future': (0.52, 0.24),
    }
    focus_x, focus_y = focus_map.get(group_key, (0.5, 0.5))
    background = _fallback_resize_cover(source_image, canvas_size, focus_x=focus_x, focus_y=focus_y)
    background = background.filter(ImageFilter.GaussianBlur(radius=20)).convert('RGBA')
    background = ImageEnhance.Brightness(background).enhance(0.48 if group_key != 'future' else 0.62)
    background = ImageEnhance.Contrast(background).enhance(1.10)

    tint = Image.new('RGBA', canvas_size, _fallback_hex_to_rgba(palette['bg0'], 132))
    canvas = Image.alpha_composite(background, tint)
    canvas_draw = ImageDraw.Draw(canvas)
    canvas_draw.rounded_rectangle((20, 20, 748, 748), radius=44, outline=border_color, width=2)
    canvas_draw.line((126, 110, 642, 110), fill=line_color, width=2)

    if group_key == 'hero':
        _fallback_add_glow(canvas, (612, 186), 158, palette['accent'], alpha=96, blur=72)
        hero_card = _fallback_make_rounded_card(
            source_image,
            (474, 604),
            radius=40,
            focus_x=0.58,
            focus_y=0.38,
            border=border_color,
        )
        _fallback_paste_card(canvas, hero_card, (246, 104))
    elif group_key == 'tension':
        _fallback_add_glow(canvas, (188, 164), 132, palette['accent2'], alpha=72, blur=68)
        main_card = _fallback_make_rounded_card(
            source_image,
            (352, 504),
            radius=36,
            focus_x=0.24,
            focus_y=0.44,
            border=border_color,
        )
        side_card = _fallback_make_rounded_card(
            source_image,
            (236, 276),
            radius=32,
            focus_x=0.78,
            focus_y=0.32,
            border=border_color,
        )
        _fallback_paste_card(canvas, main_card, (108, 152), shadow_alpha=58)
        _fallback_paste_card(canvas, side_card, (430, 218), shadow_alpha=46, shadow_offset=(0, 12))
        canvas_draw.rounded_rectangle((94, 538, 674, 626), radius=28, fill=panel_fill)
        canvas_draw.line((116, 582, 646, 582), fill=accent_color, width=4)
        canvas_draw.ellipse((620, 554, 648, 582), fill=accent_two_color)
    elif group_key == 'product':
        _fallback_add_glow(canvas, (600, 198), 138, palette['accent'], alpha=82, blur=62)
        main_card = _fallback_make_rounded_card(
            source_image,
            (430, 360),
            radius=34,
            focus_x=0.64,
            focus_y=0.38,
            border=border_color,
        )
        support_left = _fallback_make_rounded_card(
            source_image,
            (182, 224),
            radius=28,
            focus_x=0.22,
            focus_y=0.40,
            border=border_color,
        )
        support_right = _fallback_make_rounded_card(
            source_image,
            (182, 224),
            radius=28,
            focus_x=0.82,
            focus_y=0.62,
            border=border_color,
        )
        _fallback_paste_card(canvas, main_card, (168, 138))
        _fallback_paste_card(canvas, support_left, (74, 442), shadow_alpha=40, shadow_offset=(0, 10))
        _fallback_paste_card(canvas, support_right, (512, 442), shadow_alpha=40, shadow_offset=(0, 10))
        canvas_draw.line((256, 512, 384, 350), fill=accent_color, width=5)
        canvas_draw.line((508, 512, 384, 350), fill=accent_two_color, width=5)
        for cx, cy, fill in ((256, 512, accent_color), (384, 350, accent_two_color), (508, 512, accent_color)):
            canvas_draw.ellipse((cx - 10, cy - 10, cx + 10, cy + 10), fill=fill)
    elif group_key == 'process':
        _fallback_add_glow(canvas, (388, 366), 112, palette['accent2'], alpha=88, blur=54)
        cards = [
            ((64, 234), (180, 226), (0.18, 0.34)),
            ((294, 122), (180, 226), (0.52, 0.46)),
            ((524, 234), (180, 226), (0.80, 0.60)),
        ]
        centers = []
        for position, size, focus in cards:
            step_card = _fallback_make_rounded_card(
                source_image,
                size,
                radius=28,
                focus_x=focus[0],
                focus_y=focus[1],
                border=border_color,
            )
            _fallback_paste_card(canvas, step_card, position, shadow_alpha=34, shadow_offset=(0, 8))
            centers.append((position[0] + size[0] // 2, position[1] + size[1] // 2))
        canvas_draw.line((centers[0][0] + 68, centers[0][1], centers[1][0] - 68, centers[1][1]), fill=accent_color, width=6)
        canvas_draw.line((centers[1][0] + 68, centers[1][1], centers[2][0] - 68, centers[2][1]), fill=accent_two_color, width=6)
        for cx, cy in centers:
            canvas_draw.ellipse((cx - 12, cy - 12, cx + 12, cy + 12), fill=panel_soft, outline=accent_color, width=3)
    elif group_key == 'traction':
        _fallback_add_glow(canvas, (642, 158), 128, palette['accent'], alpha=84, blur=66)
        stripe_specs = [
            (108, 128, 156, 500, 0.18, 0.24),
            (306, 98, 156, 560, 0.50, 0.42),
            (504, 158, 156, 470, 0.82, 0.30),
        ]
        for x, y, width, height, fx, fy in stripe_specs:
            stripe_card = _fallback_make_rounded_card(
                source_image,
                (width, height),
                radius=28,
                focus_x=fx,
                focus_y=fy,
                border=border_color,
            )
            _fallback_paste_card(canvas, stripe_card, (x, y), shadow_alpha=30, shadow_offset=(0, 8))
        chart_points = [(112, 584), (244, 530), (376, 552), (508, 458), (640, 420)]
        canvas_draw.line(chart_points, fill=accent_color, width=6, joint='curve')
        for cx, cy in chart_points:
            canvas_draw.ellipse((cx - 10, cy - 10, cx + 10, cy + 10), fill=accent_two_color)
    elif group_key == 'future':
        horizon = _fallback_resize_cover(source_image, canvas_size, focus_x=0.52, focus_y=0.18)
        horizon = horizon.filter(ImageFilter.GaussianBlur(radius=8)).convert('RGBA')
        horizon = ImageEnhance.Brightness(horizon).enhance(0.74)
        canvas = Image.alpha_composite(horizon, Image.new('RGBA', canvas_size, _fallback_hex_to_rgba(palette['bg1'], 116)))
        canvas_draw = ImageDraw.Draw(canvas)
        canvas_draw.rounded_rectangle((20, 20, 748, 748), radius=44, outline=border_color, width=2)
        canvas_draw.line((126, 110, 642, 110), fill=line_color, width=2)
        _fallback_add_glow(canvas, (628, 170), 184, palette['accent2'], alpha=92, blur=80)
        spotlight = _fallback_make_rounded_card(
            source_image,
            (310, 404),
            radius=38,
            focus_x=0.46,
            focus_y=0.28,
            border=border_color,
        )
        _fallback_paste_card(canvas, spotlight, (396, 208), shadow_alpha=44, shadow_offset=(0, 12))

    if slide_type in {'problem', 'stakes'}:
        canvas_draw.rounded_rectangle((86, 84, 208, 130), radius=22, fill=panel_soft)
        canvas_draw.ellipse((108, 98, 122, 112), fill=accent_color)
    elif slide_type in {'solution', 'how_it_works', 'business_model'}:
        canvas_draw.rounded_rectangle((540, 82, 682, 128), radius=22, fill=panel_soft)
        canvas_draw.ellipse((562, 96, 576, 110), fill=accent_two_color)
    elif slide_type in {'impact', 'proof'}:
        canvas_draw.rounded_rectangle((562, 602, 698, 646), radius=20, fill=panel_fill)
        canvas_draw.line((584, 624, 672, 624), fill=accent_color, width=4)

    return canvas.convert('RGB')


def _fallback_variant_recipe(slide_type, index):
    slide_type = _fallback_clean_text(slide_type, 'story').lower()
    recipes = {
        'hook': {'crop': (0.0, 0.0, 1.0, 1.0), 'brightness': 1.0, 'contrast': 1.05, 'color': 1.0, 'blur': 0.0, 'sharpness': 1.05, 'overlay': None, 'mirror': False},
        'problem': {'crop': (0.0, 0.12, 0.70, 0.94), 'brightness': 0.72, 'contrast': 1.18, 'color': 0.70, 'blur': 0.5, 'sharpness': 0.95, 'overlay': ('#291626', 86), 'mirror': False},
        'stakes': {'crop': (0.18, 0.0, 0.96, 0.74), 'brightness': 0.86, 'contrast': 1.22, 'color': 0.88, 'blur': 0.0, 'sharpness': 1.08, 'overlay': ('#1E223C', 62), 'mirror': True},
        'solution': {'crop': (0.30, 0.02, 0.98, 0.82), 'brightness': 1.08, 'contrast': 1.12, 'color': 1.10, 'blur': 0.0, 'sharpness': 1.10, 'overlay': ('#10253C', 24), 'mirror': False},
        'how_it_works': {'crop': (0.18, 0.18, 0.86, 0.98), 'brightness': 0.96, 'contrast': 1.22, 'color': 0.94, 'blur': 0.0, 'sharpness': 1.14, 'overlay': ('#12233B', 44), 'mirror': True},
        'impact': {'crop': (0.10, 0.22, 0.94, 1.0), 'brightness': 1.00, 'contrast': 1.20, 'color': 1.00, 'blur': 0.0, 'sharpness': 1.12, 'overlay': ('#182745', 28), 'mirror': False},
        'proof': {'crop': (0.24, 0.10, 1.0, 0.86), 'brightness': 0.90, 'contrast': 1.28, 'color': 0.92, 'blur': 0.0, 'sharpness': 1.14, 'overlay': ('#1F2239', 44), 'mirror': True},
        'business_model': {'crop': (0.08, 0.24, 0.88, 1.0), 'brightness': 0.98, 'contrast': 1.18, 'color': 0.92, 'blur': 0.0, 'sharpness': 1.08, 'overlay': ('#1A223F', 38), 'mirror': False},
        'vision': {'crop': (0.04, 0.0, 1.0, 0.72), 'brightness': 1.10, 'contrast': 1.10, 'color': 1.12, 'blur': 0.0, 'sharpness': 1.10, 'overlay': ('#173055', 18), 'mirror': False},
        'call_to_action': {'crop': (0.0, 0.06, 1.0, 0.80), 'brightness': 1.04, 'contrast': 1.14, 'color': 1.04, 'blur': 0.0, 'sharpness': 1.08, 'overlay': ('#2A1830', 20), 'mirror': True},
    }
    recipe = dict(recipes.get(slide_type, {'crop': (0.05, 0.05, 0.95, 0.95), 'brightness': 0.98, 'contrast': 1.12, 'color': 1.0, 'blur': 0.0, 'sharpness': 1.06, 'overlay': ('#18233F', 30), 'mirror': False}))
    # Small deterministic crop drift so variants don't feel identical.
    drift = ((index % 3) - 1) * 0.03
    left, top, right, bottom = recipe['crop']
    recipe['crop'] = (
        max(0.0, min(0.18, left + drift)),
        max(0.0, min(0.2, top + (drift * 0.5))),
        min(1.0, max(0.82, right + drift)),
        min(1.0, max(0.78, bottom + (drift * 0.4))),
    )
    return recipe


def _fallback_apply_variant(image, slide, index):
    slide_type = _fallback_clean_text((slide or {}).get('type'), 'story').lower()
    recipe = _fallback_variant_recipe(slide_type, index)
    width, height = image.size
    left, top, right, bottom = recipe['crop']
    crop_box = (
        int(width * left),
        int(height * top),
        max(int(width * right), int(width * left) + 32),
        max(int(height * bottom), int(height * top) + 32),
    )
    variant = image.crop(crop_box).resize((768, 768), Image.LANCZOS)

    if recipe['blur'] > 0:
        variant = variant.filter(ImageFilter.GaussianBlur(radius=recipe['blur']))
    variant = ImageEnhance.Brightness(variant).enhance(recipe['brightness'])
    variant = ImageEnhance.Contrast(variant).enhance(recipe['contrast'])
    variant = ImageEnhance.Color(variant).enhance(recipe['color'])
    variant = ImageEnhance.Sharpness(variant).enhance(recipe.get('sharpness', 1.0))
    if recipe.get('mirror'):
        variant = variant.transpose(Image.FLIP_LEFT_RIGHT)

    overlay = recipe.get('overlay')
    if overlay:
        overlay_hex, alpha = overlay
        overlay_rgb = tuple(int(overlay_hex[i:i+2], 16) for i in (1, 3, 5))
        tint = Image.new('RGBA', variant.size, overlay_rgb + (alpha,))
        variant = Image.alpha_composite(variant.convert('RGBA'), tint).convert('RGB')

    return variant


def _fallback_derive_image_meta(source_meta, slide, index, image_options=None):
    source_image = _fallback_image_from_data_url(source_meta.get('image_url'))
    variant = _fallback_compose_remixed_image(
        source_image,
        slide,
        index,
        image_options=image_options,
    )
    return {
        'image_url': _fallback_image_to_data_url(variant, fmt='JPEG', quality=90),
        'image_prompt': source_meta.get('image_prompt') or _fallback_clean_text((slide or {}).get('visual_suggestion')),
        'image_model': FALLBACK_IMAGE_MODEL_LABEL,
        'image_repo_id': 'pollinations/flux-remixed',
        'image_status': 'remixed',
    }


def _fallback_subject_layers(topic, palette, rnd):
    accent = palette['accent']
    accent2 = palette['accent2']
    panel = palette['panel']
    panel2 = palette['panel2']
    line = palette['line']

    if topic == 'gaming':
        return f"""
        <rect x="176" y="184" width="672" height="520" rx="44" fill="{panel}" opacity="0.84"/>
        <rect x="214" y="222" width="596" height="356" rx="30" fill="{panel2}" opacity="0.96"/>
        <rect x="252" y="262" width="212" height="14" rx="7" fill="{line}" opacity="0.16"/>
        <rect x="252" y="298" width="168" height="12" rx="6" fill="{line}" opacity="0.10"/>
        <path d="M296 470 C 360 402, 430 408, 500 336 S 642 270, 720 330" fill="none" stroke="{accent}" stroke-width="12" stroke-linecap="round" opacity="0.76"/>
        <path d="M298 520 C 378 462, 456 474, 530 432 S 654 388, 734 436" fill="none" stroke="{accent2}" stroke-width="10" stroke-linecap="round" opacity="0.54"/>
        <circle cx="538" cy="372" r="92" fill="{accent}" opacity="0.12"/>
        <circle cx="632" cy="338" r="58" fill="{accent2}" opacity="0.16"/>
        <rect x="438" y="602" width="148" height="20" rx="10" fill="{line}" opacity="0.14"/>
        <rect x="476" y="622" width="72" height="90" rx="20" fill="{panel2}" opacity="0.94"/>
        <path d="M624 654 C 604 622, 566 624, 548 648 L 534 678 C 524 700, 536 726, 560 734 L 596 748 C 620 754, 644 740, 652 716 L 666 688 C 676 664, 664 640, 642 632 Z" fill="{accent2}" opacity="0.34"/>
        <circle cx="578" cy="688" r="18" fill="{panel}" opacity="0.78"/>
        <circle cx="630" cy="680" r="14" fill="{panel}" opacity="0.78"/>
        """

    if topic == 'restaurant':
        return f"""
        <rect x="168" y="180" width="688" height="560" rx="46" fill="{panel}" opacity="0.84"/>
        <rect x="218" y="228" width="588" height="464" rx="34" fill="{panel2}" opacity="0.92"/>
        <rect x="258" y="300" width="508" height="16" rx="8" fill="{line}" opacity="0.18"/>
        <rect x="258" y="436" width="508" height="16" rx="8" fill="{line}" opacity="0.18"/>
        <rect x="258" y="572" width="508" height="16" rx="8" fill="{line}" opacity="0.18"/>
        <rect x="292" y="248" width="92" height="108" rx="20" fill="{accent}" opacity="0.22"/>
        <rect x="406" y="248" width="118" height="108" rx="20" fill="{panel}" opacity="0.76"/>
        <rect x="548" y="248" width="72" height="108" rx="20" fill="{accent2}" opacity="0.22"/>
        <rect x="644" y="248" width="74" height="108" rx="20" fill="{panel}" opacity="0.68"/>
        <rect x="288" y="486" width="146" height="70" rx="18" fill="{panel}" opacity="0.82"/>
        <rect x="454" y="486" width="104" height="70" rx="18" fill="{accent2}" opacity="0.24"/>
        <rect x="582" y="484" width="150" height="150" rx="26" fill="{panel}" opacity="0.94"/>
        <rect x="612" y="520" width="90" height="12" rx="6" fill="{line}" opacity="0.16"/>
        <path d="M612 580 C 638 556, 660 556, 682 528" fill="none" stroke="{accent}" stroke-width="12" stroke-linecap="round" opacity="0.78"/>
        <circle cx="682" cy="528" r="12" fill="{accent}" opacity="0.86"/>
        """

    if topic == 'housing':
        return f"""
        <rect x="200" y="278" width="206" height="420" rx="28" fill="{panel}" opacity="0.92"/>
        <rect x="438" y="218" width="286" height="480" rx="34" fill="{panel2}" opacity="0.96"/>
        <rect x="750" y="320" width="70" height="378" rx="24" fill="{panel}" opacity="0.64"/>
        <path d="M548 150 C 578 108, 640 108, 670 150 C 670 202, 608 224, 608 272 C 608 224, 548 202, 548 150 Z" fill="{accent}" opacity="0.74"/>
        <circle cx="608" cy="154" r="24" fill="{panel}" opacity="0.62"/>
        <path d="M120 786 C 290 676, 414 706, 566 632 S 842 556, 1024 640" fill="none" stroke="{line}" stroke-width="10" opacity="0.18"/>
        <rect x="238" y="326" width="48" height="58" rx="10" fill="{line}" opacity="0.12"/>
        <rect x="320" y="326" width="48" height="58" rx="10" fill="{line}" opacity="0.12"/>
        <rect x="238" y="420" width="48" height="58" rx="10" fill="{line}" opacity="0.12"/>
        <rect x="320" y="420" width="48" height="58" rx="10" fill="{line}" opacity="0.12"/>
        <rect x="238" y="514" width="48" height="58" rx="10" fill="{line}" opacity="0.12"/>
        <rect x="320" y="514" width="48" height="58" rx="10" fill="{line}" opacity="0.12"/>
        <rect x="484" y="280" width="58" height="70" rx="12" fill="{line}" opacity="0.12"/>
        <rect x="580" y="280" width="58" height="70" rx="12" fill="{line}" opacity="0.12"/>
        <rect x="484" y="392" width="58" height="70" rx="12" fill="{line}" opacity="0.12"/>
        <rect x="580" y="392" width="58" height="70" rx="12" fill="{line}" opacity="0.12"/>
        <rect x="484" y="504" width="58" height="70" rx="12" fill="{line}" opacity="0.12"/>
        <rect x="580" y="504" width="58" height="70" rx="12" fill="{line}" opacity="0.12"/>
        """

    if topic == 'health':
        return f"""
        <circle cx="340" cy="404" r="148" fill="{accent}" opacity="0.16"/>
        <circle cx="340" cy="404" r="102" fill="{panel}" opacity="0.88"/>
        <circle cx="340" cy="404" r="36" fill="{accent2}" opacity="0.42"/>
        <rect x="470" y="238" width="324" height="430" rx="38" fill="{panel}" opacity="0.94"/>
        <rect x="512" y="290" width="112" height="18" rx="9" fill="{line}" opacity="0.14"/>
        <rect x="512" y="336" width="226" height="14" rx="7" fill="{line}" opacity="0.12"/>
        <rect x="512" y="370" width="196" height="14" rx="7" fill="{line}" opacity="0.12"/>
        <path d="M560 560 C 560 496, 612 452, 612 392 C 612 344, 648 308, 696 308 C 744 308, 778 344, 778 392 C 778 452, 726 496, 726 560 Z" fill="{accent}" opacity="0.20"/>
        <path d="M584 560 C 584 510, 622 470, 622 420 C 622 386, 650 360, 686 360 C 722 360, 750 386, 750 420 C 750 470, 788 510, 788 560 Z" fill="{accent2}" opacity="0.22"/>
        <rect x="636" y="594" width="96" height="18" rx="9" fill="{line}" opacity="0.12"/>
        """

    if topic == 'finance':
        return f"""
        <rect x="188" y="208" width="648" height="468" rx="42" fill="{panel}" opacity="0.92"/>
        <rect x="230" y="250" width="212" height="160" rx="28" fill="{panel2}" opacity="0.96"/>
        <rect x="470" y="250" width="326" height="252" rx="30" fill="{panel2}" opacity="0.92"/>
        <rect x="230" y="438" width="212" height="196" rx="28" fill="{accent2}" opacity="0.18"/>
        <path d="M518 454 C 560 434, 610 428, 650 388 S 724 310, 768 290" fill="none" stroke="{accent}" stroke-width="12" stroke-linecap="round" opacity="0.84"/>
        <circle cx="518" cy="454" r="12" fill="{accent2}" opacity="0.78"/>
        <circle cx="650" cy="388" r="12" fill="{accent}" opacity="0.84"/>
        <circle cx="768" cy="290" r="12" fill="{accent2}" opacity="0.78"/>
        <rect x="272" y="294" width="128" height="20" rx="10" fill="{line}" opacity="0.18"/>
        <rect x="272" y="336" width="84" height="14" rx="7" fill="{line}" opacity="0.12"/>
        <circle cx="326" cy="542" r="56" fill="{accent}" opacity="0.24"/>
        <path d="M326 506 L326 578" stroke="{line}" stroke-width="12" stroke-linecap="round" opacity="0.78"/>
        <path d="M294 524 C 294 502, 358 502, 358 524 C 358 548, 294 548, 294 572 C 294 594, 358 594, 358 572" fill="none" stroke="{line}" stroke-width="10" stroke-linecap="round" opacity="0.72"/>
        """

    if topic == 'sports':
        return f"""
        <rect x="184" y="214" width="656" height="448" rx="42" fill="{panel}" opacity="0.90"/>
        <rect x="216" y="246" width="592" height="384" rx="30" fill="{panel2}" opacity="0.94"/>
        <rect x="250" y="278" width="524" height="320" rx="24" fill="none" stroke="{line}" stroke-width="4" opacity="0.16"/>
        <line x1="512" y1="278" x2="512" y2="598" stroke="{line}" stroke-width="4" opacity="0.16"/>
        <circle cx="512" cy="438" r="64" fill="none" stroke="{line}" stroke-width="4" opacity="0.16"/>
        <path d="M672 314 C 702 300, 738 304, 760 332 C 782 360, 778 398, 750 420 C 722 442, 686 438, 664 410 C 642 382, 644 342, 672 314 Z" fill="{accent}" opacity="0.34"/>
        <path d="M690 334 C 718 360, 734 390, 736 420" fill="none" stroke="{panel}" stroke-width="8" stroke-linecap="round" opacity="0.55"/>
        <path d="M300 560 C 376 520, 448 512, 534 462 S 688 374, 746 328" fill="none" stroke="{accent2}" stroke-width="12" stroke-linecap="round" opacity="0.74"/>
        """

    if topic == 'education':
        return f"""
        <rect x="182" y="200" width="660" height="472" rx="44" fill="{panel}" opacity="0.90"/>
        <rect x="226" y="244" width="614" height="264" rx="30" fill="{panel2}" opacity="0.94"/>
        <rect x="270" y="292" width="280" height="18" rx="9" fill="{line}" opacity="0.16"/>
        <rect x="270" y="334" width="190" height="14" rx="7" fill="{line}" opacity="0.10"/>
        <rect x="270" y="530" width="188" height="104" rx="24" fill="{accent}" opacity="0.18"/>
        <rect x="486" y="530" width="308" height="104" rx="24" fill="{panel2}" opacity="0.76"/>
        <path d="M622 320 L734 374 L622 428 L510 374 Z" fill="{accent}" opacity="0.30"/>
        <path d="M622 374 L622 470" stroke="{line}" stroke-width="8" stroke-linecap="round" opacity="0.54"/>
        """

    if topic == 'logistics':
        return f"""
        <rect x="180" y="196" width="664" height="488" rx="46" fill="{panel}" opacity="0.92"/>
        <rect x="228" y="242" width="614" height="234" rx="32" fill="{panel2}" opacity="0.96"/>
        <rect x="246" y="534" width="248" height="120" rx="26" fill="{panel2}" opacity="0.78"/>
        <rect x="526" y="534" width="270" height="120" rx="26" fill="{accent}" opacity="0.16"/>
        <rect x="286" y="292" width="114" height="72" rx="18" fill="{accent}" opacity="0.20"/>
        <rect x="424" y="292" width="114" height="72" rx="18" fill="{accent2}" opacity="0.18"/>
        <rect x="562" y="292" width="114" height="72" rx="18" fill="{panel}" opacity="0.78"/>
        <path d="M290 594 L396 594" stroke="{line}" stroke-width="14" stroke-linecap="round" opacity="0.18"/>
        <circle cx="318" cy="634" r="18" fill="{line}" opacity="0.24"/>
        <circle cx="372" cy="634" r="18" fill="{line}" opacity="0.24"/>
        <path d="M566 592 C 612 550, 670 560, 726 516" fill="none" stroke="{accent2}" stroke-width="12" stroke-linecap="round" opacity="0.74"/>
        <circle cx="726" cy="516" r="12" fill="{accent}" opacity="0.82"/>
        """

    if topic == 'software':
        return f"""
        <rect x="182" y="186" width="660" height="520" rx="46" fill="{panel}" opacity="0.90"/>
        <rect x="226" y="230" width="616" height="198" rx="30" fill="{panel2}" opacity="0.94"/>
        <rect x="226" y="458" width="282" height="212" rx="30" fill="{panel2}" opacity="0.78"/>
        <rect x="534" y="458" width="276" height="212" rx="30" fill="{accent}" opacity="0.14"/>
        <rect x="266" y="280" width="240" height="16" rx="8" fill="{line}" opacity="0.16"/>
        <rect x="266" y="318" width="182" height="12" rx="6" fill="{line}" opacity="0.10"/>
        <circle cx="666" cy="322" r="92" fill="{accent}" opacity="0.12"/>
        <circle cx="612" cy="552" r="16" fill="{accent2}" opacity="0.74"/>
        <circle cx="676" cy="606" r="14" fill="{accent}" opacity="0.80"/>
        <circle cx="752" cy="548" r="16" fill="{accent2}" opacity="0.74"/>
        <path d="M612 552 C 642 526, 652 526, 676 606 S 718 612, 752 548" fill="none" stroke="{line}" stroke-width="8" stroke-linecap="round" opacity="0.38"/>
        """

    return f"""
    <rect x="190" y="190" width="644" height="644" rx="48" fill="{panel}" opacity="0.88"/>
    <rect x="236" y="236" width="552" height="552" rx="34" fill="{panel2}" opacity="0.84"/>
    <path d="M190 656 C 316 566, 442 560, 548 470 S 742 330, 834 360" fill="none" stroke="{accent}" stroke-width="18" stroke-linecap="round" opacity="0.30"/>
    <path d="M236 734 C 360 648, 474 652, 586 574 S 734 480, 788 504" fill="none" stroke="{accent2}" stroke-width="14" stroke-linecap="round" opacity="0.34"/>
    <circle cx="{334 + rnd.randint(-36, 42)}" cy="{356 + rnd.randint(-28, 34)}" r="72" fill="{accent}" opacity="0.18"/>
    <circle cx="{674 + rnd.randint(-44, 36)}" cy="{574 + rnd.randint(-34, 30)}" r="118" fill="{accent2}" opacity="0.16"/>
    """


def _fallback_scene_layers(slide_type, palette, rnd):
    accent = palette['accent']
    accent2 = palette['accent2']
    panel = palette['panel']
    panel2 = palette['panel2']
    line = palette['line']

    if slide_type in {'hook', 'vision', 'call_to_action'}:
        return f"""
        <circle cx="{770 + rnd.randint(-40, 70)}" cy="{235 + rnd.randint(-35, 55)}" r="{170 + rnd.randint(-20, 40)}" fill="{accent}" opacity="0.22"/>
        <path d="M0 760 C 140 660, 260 695, 420 615 S 770 530, 1024 690 L1024 1024 L0 1024 Z" fill="{panel2}" opacity="0.92"/>
        <path d="M0 830 C 180 730, 320 785, 500 700 S 820 640, 1024 760 L1024 1024 L0 1024 Z" fill="{panel}" opacity="0.98"/>
        <path d="M160 690 L320 490 L470 690 Z" fill="{accent2}" opacity="0.25"/>
        <path d="M470 710 L690 420 L890 710 Z" fill="{accent}" opacity="0.20"/>
        <circle cx="{250 + rnd.randint(-80, 80)}" cy="{220 + rnd.randint(-50, 40)}" r="4" fill="{line}" opacity="0.75"/>
        <circle cx="{340 + rnd.randint(-60, 60)}" cy="{170 + rnd.randint(-60, 30)}" r="3" fill="{line}" opacity="0.55"/>
        <circle cx="{420 + rnd.randint(-40, 90)}" cy="{260 + rnd.randint(-50, 40)}" r="2.5" fill="{line}" opacity="0.48"/>
        """

    if slide_type in {'problem', 'stakes'}:
        return f"""
        <rect x="180" y="188" width="286" height="198" rx="30" fill="{panel2}" opacity="0.92" transform="rotate(-8 323 287)"/>
        <rect x="548" y="224" width="290" height="210" rx="34" fill="{panel}" opacity="0.96" transform="rotate(7 693 329)"/>
        <path d="M180 620 C 330 540, 450 700, 610 590 S 845 470, 930 565" fill="none" stroke="{accent2}" stroke-width="20" stroke-linecap="round" opacity="0.42"/>
        <path d="M180 620 C 330 540, 450 700, 610 590 S 845 470, 930 565" fill="none" stroke="{line}" stroke-width="4" stroke-dasharray="12 14" stroke-linecap="round" opacity="0.56"/>
        <circle cx="310" cy="314" r="42" fill="{accent}" opacity="0.24"/>
        <circle cx="702" cy="323" r="34" fill="{accent2}" opacity="0.28"/>
        """

    if slide_type in {'solution', 'business_model'}:
        return f"""
        <rect x="190" y="180" width="646" height="430" rx="42" fill="{panel}" opacity="0.94"/>
        <rect x="250" y="250" width="220" height="126" rx="28" fill="{panel2}" opacity="0.96"/>
        <rect x="510" y="250" width="266" height="180" rx="30" fill="{panel2}" opacity="0.96"/>
        <rect x="250" y="406" width="526" height="124" rx="26" fill="{panel2}" opacity="0.76"/>
        <circle cx="316" cy="314" r="24" fill="{accent}" opacity="0.82"/>
        <rect x="356" y="292" width="78" height="14" rx="7" fill="{line}" opacity="0.68"/>
        <rect x="356" y="324" width="96" height="12" rx="6" fill="{line}" opacity="0.38"/>
        <path d="M485 340 C 535 318, 565 318, 610 340" fill="none" stroke="{accent2}" stroke-width="10" stroke-linecap="round" opacity="0.52"/>
        <circle cx="606" cy="338" r="18" fill="{accent2}" opacity="0.75"/>
        <circle cx="686" cy="338" r="18" fill="{accent}" opacity="0.62"/>
        <circle cx="646" cy="392" r="18" fill="{line}" opacity="0.44"/>
        """

    if slide_type in {'how_it_works'}:
        return f"""
        <circle cx="220" cy="515" r="96" fill="{accent}" opacity="0.24"/>
        <circle cx="512" cy="362" r="112" fill="{accent2}" opacity="0.20"/>
        <circle cx="806" cy="515" r="96" fill="{accent}" opacity="0.18"/>
        <path d="M288 492 C 372 432, 430 432, 512 362 S 644 432, 736 492" fill="none" stroke="{line}" stroke-width="8" stroke-dasharray="14 14" opacity="0.62"/>
        <rect x="156" y="446" width="128" height="138" rx="34" fill="{panel}" opacity="0.98"/>
        <rect x="448" y="293" width="128" height="138" rx="34" fill="{panel2}" opacity="0.98"/>
        <rect x="742" y="446" width="128" height="138" rx="34" fill="{panel}" opacity="0.98"/>
        """

    if slide_type in {'impact', 'proof'}:
        return f"""
        <rect x="158" y="238" width="708" height="430" rx="38" fill="{panel}" opacity="0.96"/>
        <rect x="236" y="514" width="84" height="96" rx="20" fill="{accent}" opacity="0.74"/>
        <rect x="370" y="430" width="84" height="180" rx="20" fill="{accent2}" opacity="0.68"/>
        <rect x="504" y="346" width="84" height="264" rx="20" fill="{accent}" opacity="0.84"/>
        <rect x="638" y="276" width="84" height="334" rx="20" fill="{line}" opacity="0.44"/>
        <path d="M220 424 C 330 396, 372 438, 460 382 S 620 242, 722 264" fill="none" stroke="{line}" stroke-width="6" opacity="0.72"/>
        <circle cx="220" cy="424" r="10" fill="{accent2}"/>
        <circle cx="460" cy="382" r="10" fill="{accent}"/>
        <circle cx="722" cy="264" r="10" fill="{accent2}"/>
        """

    return f"""
    <circle cx="{278 + rnd.randint(-60, 60)}" cy="{284 + rnd.randint(-40, 40)}" r="124" fill="{accent}" opacity="0.20"/>
    <circle cx="{760 + rnd.randint(-70, 60)}" cy="{608 + rnd.randint(-70, 50)}" r="164" fill="{accent2}" opacity="0.16"/>
    <rect x="200" y="224" width="612" height="420" rx="42" fill="{panel}" opacity="0.95"/>
    <rect x="270" y="292" width="474" height="56" rx="20" fill="{line}" opacity="0.14"/>
    <rect x="270" y="386" width="390" height="56" rx="20" fill="{accent}" opacity="0.16"/>
    <rect x="270" y="480" width="290" height="56" rx="20" fill="{accent2}" opacity="0.20"/>
    """


def _fallback_mix_rgb(color_a, color_b, ratio):
    ratio = max(0.0, min(1.0, float(ratio)))
    return tuple(
        int(color_a[idx] + (color_b[idx] - color_a[idx]) * ratio)
        for idx in range(3)
    )


def _fallback_gradient_canvas(size, top_hex, bottom_hex, mid_hex=None):
    width, height = size
    top = _fallback_hex_to_rgba(top_hex)[:3]
    bottom = _fallback_hex_to_rgba(bottom_hex)[:3]
    mid = _fallback_hex_to_rgba(mid_hex, 255)[:3] if mid_hex else None
    image = Image.new('RGBA', size, (0, 0, 0, 255))
    draw = ImageDraw.Draw(image)

    for y in range(height):
        position = y / max(height - 1, 1)
        if mid:
            if position <= 0.58:
                color = _fallback_mix_rgb(top, mid, position / 0.58)
            else:
                color = _fallback_mix_rgb(mid, bottom, (position - 0.58) / 0.42)
        else:
            color = _fallback_mix_rgb(top, bottom, position)
        draw.line((0, y, width, y), fill=color + (255,))
    return image


def _fallback_round_image(image, radius, border=None):
    rounded = image.convert('RGBA')
    mask = Image.new('L', rounded.size, 0)
    mask_draw = ImageDraw.Draw(mask)
    mask_draw.rounded_rectangle(
        (0, 0, rounded.size[0] - 1, rounded.size[1] - 1),
        radius=radius,
        fill=255,
    )
    rounded.putalpha(mask)

    if border:
        border_layer = Image.new('RGBA', rounded.size, (0, 0, 0, 0))
        border_draw = ImageDraw.Draw(border_layer)
        border_draw.rounded_rectangle(
            (1, 1, rounded.size[0] - 2, rounded.size[1] - 2),
            radius=max(radius - 1, 0),
            outline=border,
            width=2,
        )
        rounded = Image.alpha_composite(rounded, border_layer)
    return rounded


def _fallback_scene_family(topic):
    if topic in {'gaming', 'software', 'finance'}:
        return 'workspace'
    if topic in {'housing', 'education'}:
        return 'interior'
    if topic in {'restaurant', 'logistics'}:
        return 'operations'
    if topic == 'health':
        return 'lab'
    if topic == 'sports':
        return 'arena'
    return 'studio'


def _fallback_add_scene_finish(scene, palette):
    vignette = Image.new('RGBA', scene.size, (0, 0, 0, 0))
    vignette_draw = ImageDraw.Draw(vignette)
    vignette_draw.ellipse(
        (-scene.size[0] * 0.05, -scene.size[1] * 0.1, scene.size[0] * 1.05, scene.size[1] * 1.1),
        fill=(0, 0, 0, 160),
    )
    vignette = vignette.filter(ImageFilter.GaussianBlur(radius=64))
    scene = Image.alpha_composite(scene, vignette)

    highlight = Image.new('RGBA', scene.size, (0, 0, 0, 0))
    highlight_draw = ImageDraw.Draw(highlight)
    highlight_draw.ellipse(
        (scene.size[0] * 0.42, -scene.size[1] * 0.05, scene.size[0] * 0.95, scene.size[1] * 0.55),
        fill=_fallback_hex_to_rgba(palette['accent'], 62),
    )
    highlight = highlight.filter(ImageFilter.GaussianBlur(radius=70))
    scene = Image.alpha_composite(scene, highlight)
    return scene


def _fallback_draw_workspace_scene(scene, draw, palette, rnd, topic):
    width, height = scene.size
    line = _fallback_hex_to_rgba(palette['line'], 38)
    panel = _fallback_hex_to_rgba(palette['panel'], 210)
    panel2 = _fallback_hex_to_rgba(palette['panel2'], 236)
    accent = _fallback_hex_to_rgba(palette['accent'], 212)
    accent2 = _fallback_hex_to_rgba(palette['accent2'], 186)

    for idx in range(3):
        x0 = 34 + idx * 72
        draw.rounded_rectangle((x0, 64, x0 + 54, height - 142), radius=18, fill=line)

    sunlight = Image.new('RGBA', scene.size, (0, 0, 0, 0))
    sunlight_draw = ImageDraw.Draw(sunlight)
    sunlight_draw.polygon(
        [(86, 82), (238, 82), (480, height - 120), (304, height - 120)],
        fill=_fallback_hex_to_rgba('#F7C48B', 62),
    )
    sunlight = sunlight.filter(ImageFilter.GaussianBlur(radius=26))
    scene.alpha_composite(sunlight)

    shelf_y = 56
    for idx in range(6):
        x0 = 246 + idx * 92
        draw.rounded_rectangle((x0, shelf_y, x0 + 70, shelf_y + 84), radius=16, fill=panel)
        draw.rounded_rectangle((x0 + 10, shelf_y + 18, x0 + 58, shelf_y + 32), radius=7, fill=line)
        draw.rounded_rectangle((x0 + 10, shelf_y + 42, x0 + 46, shelf_y + 54), radius=6, fill=_fallback_hex_to_rgba(palette['accent'], 78))

    desk_top = height - 142
    draw.polygon(
        [(64, desk_top), (width - 86, desk_top - 18), (width - 40, height - 72), (96, height - 52)],
        fill=_fallback_hex_to_rgba('#6D4938', 224),
    )
    draw.rounded_rectangle((54, height - 118, width - 76, height - 70), radius=22, fill=_fallback_hex_to_rgba('#2A1817', 186))

    draw.rounded_rectangle((220, 174, 514, 434), radius=34, fill=panel2)
    for idx in range(4):
        draw.rounded_rectangle((248 + idx * 62, 202, 292 + idx * 62, 214), radius=6, fill=line)
    draw.arc((286, 246, 590, 484), start=196, end=318, fill=accent, width=12)
    draw.arc((282, 312, 642, 506), start=202, end=330, fill=accent2, width=10)
    draw.ellipse((484, 258, 604, 378), fill=_fallback_hex_to_rgba(palette['accent2'], 52))

    monitor_specs = [
        (170, 472, 128, 88, panel),
        (334, 422, 170, 104, panel2),
        (548, 454, 132, 92, panel),
    ]
    for x, y, w, h, fill in monitor_specs:
        draw.rounded_rectangle((x, y, x + w, y + h), radius=18, fill=fill)
        draw.rounded_rectangle((x + 12, y + 12, x + w - 12, y + h - 12), radius=12, fill=_fallback_hex_to_rgba(palette['bg1'], 238))
        draw.rounded_rectangle((x + 24, y + 26, x + w - 24, y + 38), radius=5, fill=line)
        draw.rounded_rectangle((x + 24, y + 50, x + w - 34, y + 60), radius=5, fill=_fallback_hex_to_rgba(palette['accent'], 92))

    draw.rounded_rectangle((width - 116, 120, width - 74, height - 188), radius=20, fill=_fallback_hex_to_rgba('#22382B', 182))
    draw.ellipse((width - 112, height - 252, width - 76, height - 214), fill=_fallback_hex_to_rgba('#4DB875', 190))
    draw.ellipse((width - 136, height - 222, width - 92, height - 176), fill=_fallback_hex_to_rgba('#3D9963', 180))

    if topic == 'finance':
        draw.ellipse((width - 220, 180, width - 120, 280), fill=_fallback_hex_to_rgba(palette['accent2'], 64))
        draw.line((width - 234, height - 180, width - 92, height - 244), fill=accent, width=8)
        for idx in range(4):
            cx = width - 234 + idx * 46
            cy = height - 180 - (idx % 2) * 18 - idx * 12
            draw.ellipse((cx - 11, cy - 11, cx + 11, cy + 11), fill=accent2)


def _fallback_draw_interior_scene(scene, draw, palette, rnd, topic):
    width, height = scene.size
    line = _fallback_hex_to_rgba(palette['line'], 48)
    panel = _fallback_hex_to_rgba(palette['panel'], 216)
    panel2 = _fallback_hex_to_rgba(palette['panel2'], 238)

    for idx in range(4):
        x0 = 46 + idx * 80
        draw.rounded_rectangle((x0, 54, x0 + 64, height - 150), radius=18, fill=line)

    sunlight = Image.new('RGBA', scene.size, (0, 0, 0, 0))
    sunlight_draw = ImageDraw.Draw(sunlight)
    sunlight_draw.polygon(
        [(80, 64), (238, 64), (458, height - 122), (260, height - 96)],
        fill=_fallback_hex_to_rgba('#EEC88D', 72),
    )
    sunlight = sunlight.filter(ImageFilter.GaussianBlur(radius=32))
    scene.alpha_composite(sunlight)

    draw.rounded_rectangle((width * 0.52, 58, width - 64, height * 0.42), radius=28, fill=panel2)
    draw.rounded_rectangle((width * 0.56, height * 0.18, width - 94, height * 0.22), radius=14, fill=line)
    for idx in range(4):
        x0 = int(width * 0.58) + idx * 62
        draw.rounded_rectangle((x0, int(height * 0.26), x0 + 42, int(height * 0.58)), radius=12, fill=panel)
        for shelf in range(4):
            y = int(height * 0.3) + shelf * 48
            draw.rounded_rectangle((x0 + 8, y, x0 + 34, y + 10), radius=4, fill=_fallback_hex_to_rgba(palette['accent'], 88 if idx % 2 == 0 else 54))

    draw.line((width * 0.48, height * 0.36, width - 74, height * 0.36), fill=line, width=6)
    draw.line((width * 0.48, height * 0.38, width - 74, height * 0.38), fill=_fallback_hex_to_rgba(palette['bg0'], 160), width=2)

    floor_y = int(height * 0.66)
    draw.polygon(
        [(0, floor_y), (width, floor_y - 24), (width, height), (0, height)],
        fill=_fallback_hex_to_rgba('#845E49', 220),
    )
    for idx in range(9):
        x = 12 + idx * 82
        draw.line((x, floor_y + 18, x + 38, height - 10), fill=_fallback_hex_to_rgba('#A77A58', 104), width=3)

    draw.rounded_rectangle((104, height - 222, 286, height - 120), radius=28, fill=panel)
    draw.rounded_rectangle((286, height - 204, 394, height - 120), radius=24, fill=panel2)
    draw.rounded_rectangle((318, height - 248, 436, height - 160), radius=22, fill=_fallback_hex_to_rgba(palette['panel2'], 198))
    draw.ellipse((width * 0.76, height - 224, width * 0.82, height - 170), fill=_fallback_hex_to_rgba('#63B67A', 184))
    draw.ellipse((width * 0.72, height - 188, width * 0.79, height - 132), fill=_fallback_hex_to_rgba('#4C9E67', 188))

    if topic == 'housing':
        draw.rounded_rectangle((width * 0.56, height * 0.48, width - 134, height * 0.56), radius=16, fill=_fallback_hex_to_rgba('#A17053', 168))
        draw.rounded_rectangle((width * 0.6, height * 0.52, width * 0.68, height * 0.64), radius=12, fill=_fallback_hex_to_rgba(palette['panel'], 180))


def _fallback_draw_operations_scene(scene, draw, palette, rnd, topic):
    width, height = scene.size
    line = _fallback_hex_to_rgba(palette['line'], 50)
    panel = _fallback_hex_to_rgba(palette['panel'], 220)
    panel2 = _fallback_hex_to_rgba(palette['panel2'], 240)
    accent = _fallback_hex_to_rgba(palette['accent'], 188)
    accent2 = _fallback_hex_to_rgba(palette['accent2'], 176)

    draw.rounded_rectangle((64, 78, width - 64, height - 124), radius=36, fill=panel)
    draw.rounded_rectangle((96, 112, width * 0.42, height * 0.36), radius=26, fill=panel2)
    draw.rounded_rectangle((width * 0.48, 112, width - 100, height * 0.56), radius=30, fill=panel2)
    draw.rounded_rectangle((96, height * 0.46, width * 0.42, height - 168), radius=28, fill=_fallback_hex_to_rgba(palette['panel2'], 188))
    draw.rounded_rectangle((width * 0.48, height * 0.62, width - 100, height - 140), radius=28, fill=_fallback_hex_to_rgba(palette['accent'], 36))

    for idx in range(4):
        x0 = 124 + idx * 62
        draw.rounded_rectangle((x0, 148, x0 + 42, 216), radius=12, fill=_fallback_hex_to_rgba(palette['accent'], 62 if idx % 2 == 0 else 34))
    for idx in range(5):
        x0 = int(width * 0.53) + idx * 64
        draw.rounded_rectangle((x0, 166, x0 + 42, 42 + 166), radius=12, fill=line)
    draw.rounded_rectangle((int(width * 0.53), 238, width - 134, 252), radius=7, fill=line)
    draw.rounded_rectangle((int(width * 0.53), 276, width - 214, 288), radius=6, fill=_fallback_hex_to_rgba(palette['accent'], 72))

    draw.line((148, height - 176, width - 156, height - 238), fill=accent, width=10)
    draw.line((152, height - 132, width - 124, height - 178), fill=accent2, width=8)
    for point in [(148, height - 176), (332, height - 210), (528, height - 192), (width - 156, height - 238)]:
        draw.ellipse((point[0] - 12, point[1] - 12, point[0] + 12, point[1] + 12), fill=accent2)

    if topic == 'restaurant':
        for idx in range(5):
            x0 = int(width * 0.56) + idx * 54
            draw.rounded_rectangle((x0, 352, x0 + 36, 456), radius=10, fill=_fallback_hex_to_rgba('#B1805B', 166))
        draw.rounded_rectangle((width * 0.6, height - 214, width * 0.82, height - 132), radius=24, fill=panel)
    elif topic == 'logistics':
        for idx in range(3):
            base_x = 148 + idx * 110
            draw.rounded_rectangle((base_x, height - 284, base_x + 82, height - 210), radius=18, fill=_fallback_hex_to_rgba('#9A704C', 176))
            draw.rounded_rectangle((base_x + 16, height - 254, base_x + 66, height - 240), radius=6, fill=line)


def _fallback_draw_lab_scene(scene, draw, palette, rnd):
    width, height = scene.size
    line = _fallback_hex_to_rgba(palette['line'], 42)
    panel = _fallback_hex_to_rgba(palette['panel'], 214)
    panel2 = _fallback_hex_to_rgba(palette['panel2'], 236)
    accent = _fallback_hex_to_rgba(palette['accent'], 176)
    accent2 = _fallback_hex_to_rgba(palette['accent2'], 162)

    draw.ellipse((68, 160, 376, 468), fill=_fallback_hex_to_rgba(palette['accent'], 28))
    draw.ellipse((122, 214, 322, 414), fill=panel)
    draw.ellipse((182, 274, 262, 354), fill=accent2)
    draw.rounded_rectangle((420, 110, width - 96, height - 120), radius=38, fill=panel)
    draw.rounded_rectangle((472, 176, width - 140, 196), radius=10, fill=line)
    draw.rounded_rectangle((472, 228, width - 214, 242), radius=7, fill=line)
    draw.rounded_rectangle((472, 266, width - 246, 280), radius=7, fill=_fallback_hex_to_rgba(palette['accent'], 74))
    draw.ellipse((540, 364, 812, 656), fill=_fallback_hex_to_rgba(palette['accent'], 34))
    draw.ellipse((588, 408, 764, 620), fill=_fallback_hex_to_rgba(palette['accent2'], 28))
    draw.rounded_rectangle((618, 662, 714, 682), radius=10, fill=line)


def _fallback_draw_arena_scene(scene, draw, palette, rnd):
    width, height = scene.size
    line = _fallback_hex_to_rgba(palette['line'], 44)
    panel = _fallback_hex_to_rgba(palette['panel'], 208)
    panel2 = _fallback_hex_to_rgba(palette['panel2'], 234)
    accent = _fallback_hex_to_rgba(palette['accent'], 182)
    accent2 = _fallback_hex_to_rgba(palette['accent2'], 176)

    draw.rounded_rectangle((78, 110, width - 78, height - 140), radius=42, fill=panel)
    draw.rounded_rectangle((112, 146, width - 112, height - 176), radius=30, fill=panel2)
    draw.rounded_rectangle((152, 182, width - 152, height - 218), radius=28, outline=line, width=4)
    draw.line((width // 2, 182, width // 2, height - 218), fill=line, width=4)
    draw.ellipse((width * 0.42, height * 0.38, width * 0.58, height * 0.54), outline=line, width=4)
    draw.arc((width - 284, 182, width - 124, 322), start=214, end=30, fill=accent, width=16)
    draw.line((182, height - 202, width - 176, 240), fill=accent2, width=12)
    for point in [(182, height - 202), (402, height - 168), (604, 322), (width - 176, 240)]:
        draw.ellipse((point[0] - 12, point[1] - 12, point[0] + 12, point[1] + 12), fill=accent)


def _fallback_draw_studio_scene(scene, draw, palette, rnd):
    width, height = scene.size
    line = _fallback_hex_to_rgba(palette['line'], 42)
    panel = _fallback_hex_to_rgba(palette['panel'], 220)
    panel2 = _fallback_hex_to_rgba(palette['panel2'], 236)

    draw.rounded_rectangle((126, 124, width - 126, height - 180), radius=38, fill=panel)
    draw.rounded_rectangle((198, 182, width - 198, 256), radius=24, fill=line)
    draw.rounded_rectangle((198, 308, width - 244, 382), radius=24, fill=_fallback_hex_to_rgba(palette['accent'], 52))
    draw.rounded_rectangle((198, 438, width - 356, 512), radius=24, fill=_fallback_hex_to_rgba(palette['accent2'], 46))
    draw.ellipse((286, 216, 430, 360), fill=_fallback_hex_to_rgba(palette['accent'], 36))
    draw.ellipse((width - 356, height - 330, width - 148, height - 122), fill=_fallback_hex_to_rgba(palette['accent2'], 28))


def _fallback_generate_scene_image(topic, slide_type, palette, size, seed):
    scene = _fallback_gradient_canvas(size, palette['bg1'], palette['bg0'], palette['bg2']).convert('RGBA')
    rnd = random.Random(seed)
    draw = ImageDraw.Draw(scene)
    _fallback_add_glow(scene, (int(size[0] * 0.82), int(size[1] * 0.18)), int(size[0] * 0.18), palette['accent'], alpha=62, blur=54)
    _fallback_add_glow(scene, (int(size[0] * 0.18), int(size[1] * 0.82)), int(size[0] * 0.16), palette['accent2'], alpha=42, blur=48)

    family = _fallback_scene_family(topic)
    if family == 'workspace':
        _fallback_draw_workspace_scene(scene, draw, palette, rnd, topic)
    elif family == 'interior':
        _fallback_draw_interior_scene(scene, draw, palette, rnd, topic)
    elif family == 'operations':
        _fallback_draw_operations_scene(scene, draw, palette, rnd, topic)
    elif family == 'lab':
        _fallback_draw_lab_scene(scene, draw, palette, rnd)
    elif family == 'arena':
        _fallback_draw_arena_scene(scene, draw, palette, rnd)
    else:
        _fallback_draw_studio_scene(scene, draw, palette, rnd)

    return _fallback_add_scene_finish(scene, palette)


def _fallback_build_abstract_illustration(topic, slide_type, palette, seed):
    topic = _fallback_clean_text(topic, 'generic').lower()
    slide_type = _fallback_clean_text(slide_type, 'story').lower()
    rnd = random.Random(seed)
    canvas = _fallback_gradient_canvas((1024, 1024), palette['bg1'], palette['bg0'], palette['bg2']).convert('RGBA')
    draw = ImageDraw.Draw(canvas)
    border_color = _fallback_hex_to_rgba(palette['line'], 54)
    line_color = _fallback_hex_to_rgba(palette['line'], 94)
    line_soft = _fallback_hex_to_rgba(palette['line'], 40)
    accent = _fallback_hex_to_rgba(palette['accent'], 214)
    accent2 = _fallback_hex_to_rgba(palette['accent2'], 196)
    panel_fill = _fallback_hex_to_rgba(palette['panel'], 214)
    panel_soft = _fallback_hex_to_rgba(palette['panel2'], 186)

    draw.rounded_rectangle((28, 28, 996, 996), radius=60, outline=border_color, width=2)
    draw.line((122, 108, 760, 108), fill=line_color, width=3)
    for x, fill in [(80, accent), (104, accent2), (124, _fallback_hex_to_rgba(palette['line'], 84))]:
        draw.ellipse((x, 92, x + 14, 106), fill=fill)

    for _ in range(44):
        x = rnd.randint(86, 936)
        y = rnd.randint(74, 940)
        r = rnd.randint(2, 5)
        fill = accent if rnd.random() > 0.5 else accent2
        draw.ellipse((x - r, y - r, x + r, y + r), fill=_fallback_hex_to_rgba('#FFFFFF', 96 if fill == accent2 else 72))

    def pill(box, fill, outline=None, width=0, radius=20):
        draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=width)

    def neon_line(points, fill, width=12, glow=64):
        glow_layer = Image.new('RGBA', canvas.size, (0, 0, 0, 0))
        glow_draw = ImageDraw.Draw(glow_layer)
        glow_draw.line(points, fill=_fallback_hex_to_rgba('#FFFFFF', glow), width=width + 10, joint='curve')
        glow_layer = glow_layer.filter(ImageFilter.GaussianBlur(radius=18))
        canvas.alpha_composite(glow_layer)
        draw.line(points, fill=fill, width=width, joint='curve')
        draw.line(points, fill=_fallback_hex_to_rgba(palette['line'], 108), width=max(2, width // 4), joint='curve')

    def circuit_grid():
        for _ in range(16):
            x = rnd.randint(88, 936)
            y0 = rnd.randint(120, 280)
            y1 = rnd.randint(620, 930)
            color = accent if rnd.random() > 0.45 else accent2
            draw.line((x, y0, x, y1), fill=_fallback_hex_to_rgba('#FFFFFF', 28), width=rnd.randint(1, 3))
            step_y = rnd.randint(y0 + 60, y1 - 60)
            branch = rnd.randint(-140, 160)
            draw.line((x, step_y, x + branch, step_y), fill=_fallback_hex_to_rgba(palette['line'], 18), width=2)
            for ny in (y0 + 24, step_y, y1 - 24):
                rr = rnd.randint(4, 8)
                node_fill = accent if rnd.random() > 0.5 else accent2
                draw.ellipse((x - rr, ny - rr, x + rr, ny + rr), fill=node_fill)

    def floating_card(x, y, w, h, border_fill, inner_fill=None, with_chart=False):
        pill((x, y, x + w, y + h), border_fill, outline=_fallback_hex_to_rgba(palette['line'], 44), width=2, radius=28)
        if inner_fill:
            pill((x + 18, y + 18, x + w - 18, y + h - 18), inner_fill, radius=22)
        if with_chart:
            neon_line(
                [
                    (x + 52, y + h - 88),
                    (x + 116, y + h - 126),
                    (x + 188, y + h - 92),
                    (x + 262, y + h - 152),
                    (x + w - 46, y + h - 114),
                ],
                accent,
                width=9,
                glow=72,
            )

    def draw_dashboard_shell(box, highlight=None, accent_fill=None):
        x1, y1, x2, y2 = box
        highlight = highlight or accent
        accent_fill = accent_fill or accent2
        pill((x1, y1, x2, y2), panel_soft, outline=_fallback_hex_to_rgba(palette['line'], 52), width=2, radius=32)
        pill((x1 + 20, y1 + 20, x2 - 20, y2 - 20), _fallback_hex_to_rgba('#0E1331', 220), radius=24)
        for idx, width_ratio in enumerate((0.28, 0.22, 0.16)):
            bar_w = int((x2 - x1) * width_ratio)
            bar_x2 = x1 + 64 + bar_w
            pill((x1 + 48, y1 + 42 + idx * 38, bar_x2, y1 + 56 + idx * 38), line_soft, radius=8)
        neon_line(
            [
                (x1 + 64, y2 - 120),
                (x1 + 156, y2 - 178),
                (x1 + 252, y2 - 132),
                (x1 + 346, y2 - 204),
                (x2 - 54, y2 - 152),
            ],
            highlight,
            width=10,
            glow=74,
        )
        draw.ellipse((x2 - 170, y1 + 54, x2 - 64, y1 + 160), fill=_fallback_hex_to_rgba(palette['accent2'], 64))
        draw.ellipse((x2 - 132, y1 + 82, x2 - 88, y1 + 126), fill=accent_fill)

    def draw_controller(box):
        x1, y1, x2, y2 = box
        cx = (x1 + x2) / 2
        cy = y1 + (y2 - y1) * 0.74
        body_w = (x2 - x1) * 0.42
        grip_w = (x2 - x1) * 0.22
        grip_h = (y2 - y1) * 0.18
        pill((cx - body_w / 2, cy - grip_h / 2, cx + body_w / 2, cy + grip_h / 2), _fallback_hex_to_rgba(palette['panel'], 206), radius=28)
        draw.ellipse((cx - body_w / 2 - grip_w * 0.78, cy - grip_h * 0.84, cx - body_w / 2 + grip_w * 0.28, cy + grip_h * 0.84), fill=_fallback_hex_to_rgba(palette['panel2'], 218), outline=_fallback_hex_to_rgba(palette['line'], 82), width=2)
        draw.ellipse((cx + body_w / 2 - grip_w * 0.28, cy - grip_h * 0.84, cx + body_w / 2 + grip_w * 0.78, cy + grip_h * 0.84), fill=_fallback_hex_to_rgba(palette['panel2'], 218), outline=_fallback_hex_to_rgba(palette['line'], 82), width=2)
        draw.ellipse((cx - 74, cy - 20, cx - 34, cy + 20), fill=accent2)
        draw.ellipse((cx + 34, cy - 20, cx + 74, cy + 20), fill=accent)
        draw.ellipse((cx + 82, cy - 36, cx + 100, cy - 18), fill=accent2)
        draw.ellipse((cx + 106, cy - 12, cx + 124, cy + 6), fill=accent2)
        draw.ellipse((cx + 82, cy + 12, cx + 100, cy + 30), fill=accent2)

    def draw_building_cluster(box):
        x1, y1, x2, y2 = box
        base_y = y2 - 42
        specs = [
            (x1 + 58, y1 + 110, 118, base_y),
            (x1 + 212, y1 + 74, 146, base_y),
            (x1 + 392, y1 + 136, 104, base_y),
        ]
        for idx, (left, top, width_px, bottom) in enumerate(specs):
            pill((left, top, left + width_px, bottom), _fallback_hex_to_rgba(palette['panel'], 214 if idx == 1 else 196), outline=_fallback_hex_to_rgba(palette['line'], 62), width=2, radius=26)
            for row in range(4):
                for col in range(2 if idx != 1 else 3):
                    wx = left + 24 + col * 36
                    wy = top + 30 + row * 56
                    pill((wx, wy, wx + 16, wy + 22), _fallback_hex_to_rgba(palette['line'], 26), radius=6)
        pin_x = x1 + (x2 - x1) * 0.74
        pin_y = y1 + 74
        draw.ellipse((pin_x - 46, pin_y - 46, pin_x + 46, pin_y + 46), fill=_fallback_hex_to_rgba(palette['accent'], 90))
        draw.ellipse((pin_x - 20, pin_y - 20, pin_x + 20, pin_y + 20), fill=_fallback_hex_to_rgba(palette['panel'], 214))
        draw.polygon([(pin_x, pin_y + 64), (pin_x - 22, pin_y + 20), (pin_x + 22, pin_y + 20)], fill=_fallback_hex_to_rgba(palette['accent'], 158))
        neon_line([(x1 + 88, y2 - 84), (x1 + 220, y2 - 144), (x1 + 360, y2 - 122), (pin_x, y2 - 192)], accent2, width=8, glow=66)

    def draw_health_signal(box):
        x1, y1, x2, y2 = box
        draw_dashboard_shell((x1, y1 + 34, x2 - 120, y2 - 30), highlight=accent2, accent_fill=accent)
        heart_box = (x2 - 250, y1 + 86, x2 - 88, y1 + 248)
        pill(heart_box, _fallback_hex_to_rgba(palette['panel'], 210), outline=_fallback_hex_to_rgba(palette['line'], 54), width=2, radius=30)
        hx = (heart_box[0] + heart_box[2]) / 2
        hy = (heart_box[1] + heart_box[3]) / 2 + 10
        draw.ellipse((hx - 48, hy - 84, hx - 4, hy - 40), fill=_fallback_hex_to_rgba(palette['accent'], 170))
        draw.ellipse((hx + 4, hy - 84, hx + 48, hy - 40), fill=_fallback_hex_to_rgba(palette['accent'], 170))
        draw.polygon([(hx - 62, hy - 56), (hx + 62, hy - 56), (hx, hy + 56)], fill=_fallback_hex_to_rgba(palette['accent2'], 176))
        neon_line([(x1 + 86, y2 - 152), (x1 + 182, y2 - 152), (x1 + 230, y2 - 192), (x1 + 272, y2 - 110), (x1 + 348, y2 - 110), (x1 + 412, y2 - 164), (x2 - 146, y2 - 164)], accent, width=8, glow=64)

    def draw_finance_signal(box):
        x1, y1, x2, y2 = box
        draw_dashboard_shell((x1, y1, x2, y2), highlight=accent, accent_fill=accent2)
        ring_cx = x2 - 136
        ring_cy = y1 + 138
        draw.ellipse((ring_cx - 62, ring_cy - 62, ring_cx + 62, ring_cy + 62), outline=_fallback_hex_to_rgba(palette['accent2'], 180), width=12)
        draw.ellipse((ring_cx - 34, ring_cy - 34, ring_cx + 34, ring_cy + 34), outline=_fallback_hex_to_rgba(palette['line'], 148), width=8)
        draw.line((ring_cx, ring_cy - 24, ring_cx, ring_cy + 24), fill=_fallback_hex_to_rgba(palette['line'], 188), width=8)
        draw.arc((ring_cx - 28, ring_cy - 54, ring_cx + 28, ring_cy), start=200, end=28, fill=_fallback_hex_to_rgba(palette['line'], 188), width=8)
        draw.arc((ring_cx - 28, ring_cy, ring_cx + 28, ring_cy + 54), start=200, end=28, fill=_fallback_hex_to_rgba(palette['line'], 188), width=8)

    def draw_education_signal(box):
        x1, y1, x2, y2 = box
        bulb_cx = x1 + 208
        bulb_cy = y1 + 232
        draw.ellipse((bulb_cx - 116, bulb_cy - 116, bulb_cx + 116, bulb_cy + 116), fill=_fallback_hex_to_rgba(palette['panel'], 198), outline=_fallback_hex_to_rgba(palette['line'], 76), width=3)
        draw.ellipse((bulb_cx - 68, bulb_cy - 68, bulb_cx + 68, bulb_cy + 68), fill=_fallback_hex_to_rgba(palette['accent'], 120))
        pill((bulb_cx - 42, bulb_cy + 88, bulb_cx + 42, bulb_cy + 138), _fallback_hex_to_rgba(palette['panel2'], 214), radius=16)
        draw.line((bulb_cx, bulb_cy + 18, bulb_cx, bulb_cy + 86), fill=_fallback_hex_to_rgba(palette['line'], 188), width=8)
        for angle in range(0, 360, 45):
            theta = math.radians(angle)
            x0 = bulb_cx + math.cos(theta) * 134
            y0 = bulb_cy + math.sin(theta) * 134
            x1_line = bulb_cx + math.cos(theta) * 170
            y1_line = bulb_cy + math.sin(theta) * 170
            draw.line((x0, y0, x1_line, y1_line), fill=_fallback_hex_to_rgba(palette['line'], 98), width=5)
        pill((x1 + 420, y1 + 142, x2 - 88, y2 - 120), _fallback_hex_to_rgba(palette['panel2'], 198), outline=_fallback_hex_to_rgba(palette['line'], 52), width=2, radius=30)
        neon_line([(x1 + 470, y2 - 210), (x1 + 566, y2 - 246), (x1 + 658, y2 - 176), (x2 - 128, y2 - 228)], accent2, width=8, glow=62)

    def draw_logistics_signal(box):
        x1, y1, x2, y2 = box
        for idx, (bx, by, bw, bh) in enumerate(((x1 + 84, y1 + 142, 126, 110), (x1 + 260, y1 + 214, 154, 126), (x1 + 470, y1 + 154, 124, 110))):
            pill((bx, by, bx + bw, by + bh), _fallback_hex_to_rgba(palette['panel'], 204 if idx != 1 else 226), outline=_fallback_hex_to_rgba(palette['line'], 58), width=2, radius=24)
            draw.line((bx + 20, by + 36, bx + bw - 20, by + 36), fill=_fallback_hex_to_rgba(palette['line'], 92), width=4)
            draw.line((bx + 20, by + 64, bx + bw - 42, by + 64), fill=_fallback_hex_to_rgba(palette['accent'], 122), width=5)
        neon_line([(x1 + 136, y2 - 190), (x1 + 324, y2 - 138), (x1 + 496, y2 - 216), (x2 - 116, y2 - 152)], accent, width=10, glow=74)
        for cx, cy in ((x1 + 136, y2 - 190), (x1 + 324, y2 - 138), (x1 + 496, y2 - 216), (x2 - 116, y2 - 152)):
            draw.ellipse((cx - 16, cy - 16, cx + 16, cy + 16), fill=accent2)

    def draw_stadium_signal(box):
        x1, y1, x2, y2 = box
        pill((x1 + 78, y1 + 116, x2 - 78, y2 - 90), _fallback_hex_to_rgba(palette['panel'], 202), outline=_fallback_hex_to_rgba(palette['line'], 56), width=2, radius=42)
        pill((x1 + 132, y1 + 166, x2 - 132, y2 - 140), _fallback_hex_to_rgba(palette['panel2'], 212), outline=_fallback_hex_to_rgba(palette['line'], 42), width=2, radius=34)
        draw.line((x1 + 132, (y1 + y2) / 2, x2 - 132, (y1 + y2) / 2), fill=_fallback_hex_to_rgba(palette['line'], 76), width=4)
        draw.ellipse((((x1 + x2) / 2) - 72, ((y1 + y2) / 2) - 72, ((x1 + x2) / 2) + 72, ((y1 + y2) / 2) + 72), outline=_fallback_hex_to_rgba(palette['line'], 70), width=4)
        neon_line([(x1 + 206, y2 - 162), (x1 + 326, y2 - 214), (x1 + 468, y2 - 196), (x2 - 192, y1 + 212)], accent2, width=10, glow=68)
        draw.ellipse((x2 - 242, y1 + 152, x2 - 146, y1 + 248), fill=_fallback_hex_to_rgba(palette['accent'], 120))

    def draw_topic_focus(box):
        if topic == 'gaming':
            draw_dashboard_shell(box, highlight=accent, accent_fill=accent2)
            draw_controller(box)
        elif topic == 'software':
            draw_dashboard_shell(box, highlight=accent2, accent_fill=accent)
            pill((box[0] + 66, box[1] + 82, box[0] + 178, box[1] + 308), _fallback_hex_to_rgba(palette['panel'], 198), outline=_fallback_hex_to_rgba(palette['line'], 64), width=2, radius=26)
            draw.ellipse((box[0] + 86, box[1] + 108, box[0] + 126, box[1] + 148), fill=accent2)
        elif topic == 'housing':
            draw_building_cluster(box)
        elif topic == 'health':
            draw_health_signal(box)
        elif topic == 'finance':
            draw_finance_signal(box)
        elif topic == 'education':
            draw_education_signal(box)
        elif topic == 'logistics':
            draw_logistics_signal(box)
        elif topic == 'sports':
            draw_stadium_signal(box)
        elif topic == 'restaurant':
            draw_dashboard_shell(box, highlight=accent2, accent_fill=accent)
            for idx in range(4):
                left = box[0] + 64 + idx * 104
                pill((left, box[1] + 84, left + 72, box[1] + 200), _fallback_hex_to_rgba(palette['panel'], 188 if idx % 2 else 214), radius=18)
            draw.ellipse((box[2] - 246, box[1] + 188, box[2] - 126, box[1] + 308), fill=_fallback_hex_to_rgba(palette['accent'], 74))
            draw.ellipse((box[2] - 214, box[1] + 220, box[2] - 158, box[1] + 276), outline=_fallback_hex_to_rgba(palette['line'], 176), width=6)
            neon_line([(box[0] + 78, box[3] - 108), (box[0] + 220, box[3] - 166), (box[0] + 380, box[3] - 136), (box[2] - 114, box[3] - 194)], accent, width=8, glow=60)
        else:
            draw_dashboard_shell(box, highlight=accent, accent_fill=accent2)

    circuit_grid()
    _fallback_add_glow(canvas, (786, 210), 234, palette['accent'], alpha=108, blur=92)
    _fallback_add_glow(canvas, (238, 804), 218, palette['accent2'], alpha=84, blur=92)
    _fallback_add_glow(canvas, (510, 510), 174, palette['accent'], alpha=44, blur=120)

    if slide_type in {'hook', 'solution', 'business_model'}:
        floating_card(178, 240, 660, 454, panel_fill, inner_fill=_fallback_hex_to_rgba('#0E1331', 228), with_chart=True)
        floating_card(90, 414, 144, 248, panel_soft, inner_fill=_fallback_hex_to_rgba('#0E1331', 204))
        floating_card(724, 268, 168, 300, panel_soft, inner_fill=_fallback_hex_to_rgba('#10173A', 212))
        for y in (320, 366, 412):
            pill((272, y, 432, y + 14), line_soft, radius=7)
        for x, y, rr, fill in [(778, 344, 60, accent2), (814, 338, 18, accent), (784, 452, 14, accent)]:
            draw.ellipse((x - rr, y - rr, x + rr, y + rr), fill=fill)
        neon_line([(248, 728), (394, 650), (542, 674), (734, 586), (900, 612)], accent, width=14, glow=84)
        draw_topic_focus((224, 280, 768, 650))
    elif slide_type in {'problem', 'stakes'}:
        neon_line([(120, 850), (302, 734), (442, 748), (618, 606), (810, 518), (928, 408)], accent, width=18, glow=92)
        neon_line([(154, 768), (286, 706), (444, 644), (624, 578), (786, 432)], accent2, width=10, glow=68)
        for cx, cy, rr, fill in [(190, 738, 74, accent), (760, 334, 86, accent), (562, 560, 38, accent2)]:
            draw.ellipse((cx - rr, cy - rr, cx + rr, cy + rr), fill=_fallback_hex_to_rgba(palette['panel2'], 44))
            draw.ellipse((cx - rr // 2, cy - rr // 2, cx + rr // 2, cy + rr // 2), fill=fill)
        for box in [(126, 206, 402, 430), (586, 188, 876, 448)]:
            pill(box, panel_soft, outline=_fallback_hex_to_rgba(palette['line'], 52), width=2, radius=34)
        draw_topic_focus((152, 224, 376, 412))
        draw_topic_focus((612, 206, 850, 430))
    elif slide_type == 'how_it_works':
        centers = [(252, 534), (512, 370), (778, 534)]
        for idx, (cx, cy) in enumerate(centers):
            outer = 104 if idx == 1 else 88
            inner = outer - 26
            fill = accent if idx != 1 else accent2
            draw.ellipse((cx - outer, cy - outer, cx + outer, cy + outer), outline=fill, width=18)
            draw.ellipse((cx - inner, cy - inner, cx + inner, cy + inner), outline=_fallback_hex_to_rgba(palette['line'], 92), width=8)
            draw.ellipse((cx - 22, cy - 22, cx + 22, cy + 22), fill=fill)
            for spoke in range(0, 360, 45):
                angle = math.radians(spoke)
                x0 = cx + math.cos(angle) * (outer - 8)
                y0 = cy + math.sin(angle) * (outer - 8)
                x1 = cx + math.cos(angle) * (outer + 34)
                y1 = cy + math.sin(angle) * (outer + 34)
                draw.line((x0, y0, x1, y1), fill=_fallback_hex_to_rgba(palette['line'], 86), width=4)
        neon_line([(338, 504), (430, 432), (512, 370), (612, 438), (704, 504)], accent2, width=10, glow=72)
        for idx, box in enumerate(((138, 424, 366, 662), (398, 256, 626, 494), (664, 424, 892, 662))):
            draw_topic_focus(box)
            badge_cx = box[0] + 38
            badge_cy = box[1] + 38
            draw.ellipse((badge_cx - 18, badge_cy - 18, badge_cx + 18, badge_cy + 18), fill=_fallback_hex_to_rgba(palette['panel'], 216), outline=_fallback_hex_to_rgba(palette['line'], 104), width=2)
            draw.ellipse((badge_cx - 8, badge_cy - 8, badge_cx + 8, badge_cy + 8), fill=accent if idx != 1 else accent2)
    elif slide_type in {'impact', 'proof'}:
        pill((246, 212, 776, 730), panel_fill, outline=_fallback_hex_to_rgba(palette['line'], 48), width=2, radius=38)
        for idx, (x, h, fill) in enumerate([(330, 168, accent), (456, 272, accent2), (584, 368, accent), (710, 446, accent2)]):
            pill((x, 678 - h, x + 74, 678), fill, radius=18)
        neon_line([(292, 518), (394, 472), (500, 504), (618, 390), (748, 310)], accent, width=10, glow=82)
        for cx, cy, fill in [(292, 518, accent2), (500, 504, accent), (748, 310, accent2)]:
            draw.ellipse((cx - 14, cy - 14, cx + 14, cy + 14), fill=fill)
        draw_topic_focus((306, 246, 728, 486))
    elif slide_type in {'vision', 'call_to_action'}:
        neon_line([(86, 846), (230, 764), (380, 724), (582, 646), (748, 556), (934, 430)], accent, width=16, glow=92)
        pill((186, 286, 842, 720), panel_fill, outline=_fallback_hex_to_rgba(palette['line'], 48), width=2, radius=42)
        draw.ellipse((426, 344, 594, 512), fill=_fallback_hex_to_rgba(palette['accent'], 56))
        draw.ellipse((454, 372, 566, 484), fill=_fallback_hex_to_rgba(palette['accent2'], 62))
        draw.line((510, 392, 510, 464), fill=_fallback_hex_to_rgba(palette['line'], 180), width=12)
        draw.line((474, 428, 546, 428), fill=_fallback_hex_to_rgba(palette['line'], 180), width=12)
        draw_topic_focus((248, 332, 778, 674))
    else:
        pill((214, 240, 820, 706), panel_fill, outline=_fallback_hex_to_rgba(palette['line'], 46), width=2, radius=44)
        floating_card(264, 304, 506, 288, panel_soft, inner_fill=_fallback_hex_to_rgba('#0D1330', 220), with_chart=True)
        for box, fill in [((286, 644, 452, 744), accent), ((480, 644, 610, 744), accent2), ((638, 644, 760, 744), _fallback_hex_to_rgba(palette['line'], 110))]:
            pill(box, fill, radius=24)
        draw_topic_focus((280, 328, 748, 612))

    return canvas.convert('RGB')


def _fallback_apply_scene_treatment(scene_image, palette, variant):
    scene = scene_image.convert('RGB')
    variant = _fallback_clean_text(variant, 'presentation').lower()

    if variant == 'presentation':
        scene = ImageEnhance.Color(scene).enhance(1.10)
        scene = ImageEnhance.Contrast(scene).enhance(1.12)
        scene = ImageEnhance.Brightness(scene).enhance(1.32)
        scene = ImageEnhance.Sharpness(scene).enhance(1.18)
        rgba = scene.convert('RGBA')
        wash = Image.new('RGBA', scene.size, _fallback_hex_to_rgba(palette['line'], 10))
        tint = Image.new('RGBA', scene.size, _fallback_hex_to_rgba(palette['accent'], 10))
        rgba = Image.alpha_composite(rgba, wash)
        scene = Image.alpha_composite(rgba, tint).convert('RGB')
        return scene

    if variant == 'animated':
        scene = ImageEnhance.Color(scene).enhance(1.34)
        scene = ImageEnhance.Contrast(scene).enhance(1.14)
        scene = ImageEnhance.Brightness(scene).enhance(1.24)
        scene = ImageEnhance.Sharpness(scene).enhance(1.22)
        tint = Image.new('RGBA', scene.size, _fallback_hex_to_rgba(palette['accent2'], 18))
        scene = Image.alpha_composite(scene.convert('RGBA'), tint).convert('RGB')
        return scene

    if variant == 'cartoon':
        scene = ImageEnhance.Color(scene).enhance(1.58)
        scene = ImageEnhance.Contrast(scene).enhance(1.18)
        scene = ImageOps.posterize(scene, 4)
        scene = scene.quantize(colors=28, method=Image.MEDIANCUT).convert('RGB')
        scene = scene.filter(ImageFilter.SMOOTH_MORE)
        scene = ImageEnhance.Sharpness(scene).enhance(1.3)
        tint = Image.new('RGBA', scene.size, _fallback_hex_to_rgba(palette['accent'], 16))
        scene = Image.alpha_composite(scene.convert('RGBA'), tint).convert('RGB')
        return scene

    return scene


def _fallback_scene_focus_for_slide_type(slide_type):
    slide_type = _fallback_clean_text(slide_type, 'story').lower()
    return {
        'hook': (0.56, 0.48),
        'problem': (0.52, 0.52),
        'stakes': (0.5, 0.5),
        'solution': (0.58, 0.5),
        'how_it_works': (0.5, 0.46),
        'impact': (0.54, 0.48),
        'proof': (0.52, 0.46),
        'business_model': (0.55, 0.5),
        'vision': (0.56, 0.44),
        'call_to_action': (0.54, 0.46),
    }.get(slide_type, (0.5, 0.5))


def _fallback_build_scene_visual(topic, slide_type, palette, seed, variant='presentation'):
    scene = _fallback_generate_scene_image(topic, slide_type, palette, (1280, 960), seed + 17)
    scene = _fallback_apply_scene_treatment(scene, palette, variant)
    focus_x, focus_y = _fallback_scene_focus_for_slide_type(slide_type)
    return _fallback_resize_cover(scene, (1024, 1024), focus_x=focus_x, focus_y=focus_y)


def _fallback_build_scene_poster(topic, slide_type, palette, seed, variant='presentation'):
    topic = _fallback_clean_text(topic, 'generic').lower()
    slide_type = _fallback_clean_text(slide_type, 'story').lower()
    variant = _fallback_clean_text(variant, 'presentation').lower()
    rnd = random.Random(seed)
    canvas = _fallback_gradient_canvas((1024, 1024), palette['bg1'], palette['bg0'], palette['bg2']).convert('RGBA')
    draw = ImageDraw.Draw(canvas)
    border_color = _fallback_hex_to_rgba(palette['line'], 60)
    line_color = _fallback_hex_to_rgba(palette['line'], 92)
    accent = _fallback_hex_to_rgba(palette['accent'], 196)
    accent2 = _fallback_hex_to_rgba(palette['accent2'], 188)
    panel_fill = _fallback_hex_to_rgba(palette['panel'], 176)
    panel_soft = _fallback_hex_to_rgba(palette['panel2'], 156)

    _fallback_add_glow(canvas, (824, 184), 226, palette['accent'], alpha=56 if variant == 'presentation' else 82, blur=84)
    _fallback_add_glow(canvas, (190, 822), 204, palette['accent2'], alpha=42 if variant == 'presentation' else 72, blur=78)
    draw.rounded_rectangle((24, 24, 1000, 1000), radius=58, outline=border_color, width=2)
    draw.line((122, 108, 770, 108), fill=line_color, width=3)
    for x, fill in ((82, accent), (106, accent2), (126, _fallback_hex_to_rgba(palette['line'], 96))):
        draw.ellipse((x, 94, x + 14, 108), fill=fill)

    base_scene = _fallback_generate_scene_image(topic, slide_type, palette, (1280, 960), seed + 17)
    base_scene = _fallback_apply_scene_treatment(base_scene, palette, variant)

    def make_scene_card(card_size, focus_x=0.5, focus_y=0.5, radius=36):
        crop = _fallback_resize_cover(base_scene, card_size, focus_x=focus_x, focus_y=focus_y)
        return _fallback_round_image(crop, radius=radius, border=border_color)

    hero_focus = {
        'hook': (0.56, 0.48),
        'problem': (0.52, 0.52),
        'stakes': (0.5, 0.5),
        'solution': (0.58, 0.5),
        'how_it_works': (0.5, 0.46),
        'impact': (0.54, 0.48),
        'proof': (0.52, 0.46),
        'business_model': (0.55, 0.5),
        'vision': (0.56, 0.44),
        'call_to_action': (0.54, 0.46),
    }.get(slide_type, (0.5, 0.5))

    hero_size = (764, 646)
    hero_card = make_scene_card(hero_size, focus_x=hero_focus[0], focus_y=hero_focus[1], radius=44)
    _fallback_paste_card(canvas, hero_card, (130, 182), shadow_alpha=48, shadow_offset=(0, 18))

    if variant == 'presentation':
        left_card = make_scene_card((172, 238), focus_x=0.28, focus_y=0.56, radius=28)
        right_card = make_scene_card((176, 240), focus_x=0.74, focus_y=0.42, radius=28)
        _fallback_paste_card(canvas, left_card, (78, 540), shadow_alpha=28, shadow_offset=(0, 10))
        _fallback_paste_card(canvas, right_card, (770, 222), shadow_alpha=28, shadow_offset=(0, 10))
        draw.rounded_rectangle((106, 146, 248, 194), radius=24, fill=panel_fill)
        draw.rounded_rectangle((758, 824, 924, 874), radius=24, fill=panel_fill)
        draw.rounded_rectangle((792, 846, 886, 856), radius=5, fill=line_color)
        draw.line((154, 876, 872, 876), fill=_fallback_hex_to_rgba(palette['line'], 40), width=3)
    elif variant == 'animated':
        back_left = make_scene_card((216, 280), focus_x=0.24, focus_y=0.4, radius=30)
        back_right = make_scene_card((216, 280), focus_x=0.78, focus_y=0.58, radius=30)
        _fallback_paste_card(canvas, back_left, (94, 230), shadow_alpha=22, shadow_offset=(0, 8))
        _fallback_paste_card(canvas, back_right, (714, 252), shadow_alpha=22, shadow_offset=(0, 8))
        for _ in range(16):
            x = rnd.randint(118, 918)
            y = rnd.randint(150, 902)
            r = rnd.randint(2, 5)
            draw.ellipse((x - r, y - r, x + r, y + r), fill=accent if rnd.random() > 0.5 else accent2)
        draw.line((188, 786, 378, 702, 594, 736, 842, 628), fill=accent, width=10)
        draw.line((214, 706, 364, 640, 566, 686, 806, 566), fill=accent2, width=8)
        draw.rounded_rectangle((106, 146, 232, 194), radius=24, fill=panel_fill)
        draw.rounded_rectangle((792, 150, 920, 198), radius=24, fill=panel_soft)
    elif variant == 'cartoon':
        sticker = Image.new('RGBA', canvas.size, (0, 0, 0, 0))
        sticker_draw = ImageDraw.Draw(sticker)
        sticker_draw.ellipse((724, 160, 928, 364), fill=_fallback_hex_to_rgba(palette['accent'], 84))
        sticker_draw.ellipse((748, 184, 904, 340), fill=_fallback_hex_to_rgba(palette['panel'], 210), outline=_fallback_hex_to_rgba(palette['line'], 132), width=4)
        sticker_draw.ellipse((790, 222, 862, 294), fill=accent2)
        sticker_draw.rounded_rectangle((110, 786, 364, 882), radius=34, fill=panel_fill, outline=_fallback_hex_to_rgba(palette['line'], 112), width=3)
        sticker_draw.rounded_rectangle((160, 826, 314, 844), radius=9, fill=line_color)
        sticker = sticker.filter(ImageFilter.GaussianBlur(radius=0.2))
        canvas.alpha_composite(sticker)
        draw.arc((164, 120, 862, 786), start=208, end=328, fill=accent, width=12)
        for cx, cy, rr in ((188, 226, 20), (866, 408, 16), (774, 758, 22)):
            draw.ellipse((cx - rr, cy - rr, cx + rr, cy + rr), fill=accent2)

    return canvas.convert('RGB')


def _fallback_build_editorial_illustration(topic, slide_type, palette, seed):
    canvas = _fallback_gradient_canvas((1024, 1024), palette['bg1'], palette['bg0'], palette['bg2']).convert('RGBA')
    draw = ImageDraw.Draw(canvas)
    border_color = _fallback_hex_to_rgba(palette['line'], 60)
    line_color = _fallback_hex_to_rgba(palette['line'], 86)
    accent = _fallback_hex_to_rgba(palette['accent'], 188)
    accent2 = _fallback_hex_to_rgba(palette['accent2'], 180)
    panel_fill = _fallback_hex_to_rgba(palette['panel'], 160)

    _fallback_add_glow(canvas, (790, 212), 218, palette['accent'], alpha=76, blur=84)
    _fallback_add_glow(canvas, (184, 832), 196, palette['accent2'], alpha=48, blur=72)
    draw.rounded_rectangle((26, 26, 998, 998), radius=58, outline=border_color, width=2)
    draw.line((126, 112, 772, 112), fill=line_color, width=3)
    draw.ellipse((82, 96, 96, 110), fill=accent)
    draw.ellipse((108, 100, 116, 108), fill=_fallback_hex_to_rgba(palette['line'], 108))
    draw.ellipse((126, 100, 134, 108), fill=_fallback_hex_to_rgba(palette['line'], 62))

    def make_scene_card(card_size, card_seed):
        scene_image = _fallback_generate_scene_image(topic, slide_type, palette, (1024, 1024), card_seed)
        scene_image = _fallback_resize_cover(scene_image, card_size, focus_x=0.5, focus_y=0.5)
        return _fallback_round_image(scene_image, radius=36, border=border_color)

    if slide_type == 'how_it_works':
        positions = [(82, 286), (362, 142), (642, 286)]
        centers = []
        for idx, pos in enumerate(positions):
            card = make_scene_card((240, 324), seed + idx * 13 + 5)
            _fallback_paste_card(canvas, card, pos, shadow_alpha=34, shadow_offset=(0, 10))
            centers.append((pos[0] + 120, pos[1] + 162))
        draw.line((centers[0][0] + 82, centers[0][1], centers[1][0] - 82, centers[1][1]), fill=accent, width=8)
        draw.line((centers[1][0] + 82, centers[1][1], centers[2][0] - 82, centers[2][1]), fill=accent2, width=8)
        for cx, cy in centers:
            draw.ellipse((cx - 14, cy - 14, cx + 14, cy + 14), fill=panel_fill, outline=accent, width=3)
        draw.rounded_rectangle((104, 738, 922, 828), radius=28, fill=panel_fill)
        draw.rounded_rectangle((136, 774, 742, 786), radius=6, fill=line_color)
        draw.rounded_rectangle((136, 804, 656, 814), radius=5, fill=_fallback_hex_to_rgba(palette['accent'], 104))
    elif slide_type in {'impact', 'proof'}:
        stripe_specs = [
            (116, 140, 208, 612, 7),
            (404, 112, 220, 670, 19),
            (700, 168, 192, 558, 31),
        ]
        trend_points = []
        for x, y, w, h, offset in stripe_specs:
            card = make_scene_card((w, h), seed + offset)
            _fallback_paste_card(canvas, card, (x, y), shadow_alpha=30, shadow_offset=(0, 10))
            trend_points.append((x + w // 2, y + h - 76 - (offset % 5) * 18))
        draw.line(trend_points, fill=accent, width=8)
        for point in trend_points:
            draw.ellipse((point[0] - 13, point[1] - 13, point[0] + 13, point[1] + 13), fill=accent2)
        draw.rounded_rectangle((750, 786, 914, 848), radius=26, fill=panel_fill)
        draw.rounded_rectangle((788, 812, 876, 822), radius=5, fill=line_color)
    elif slide_type in {'solution', 'business_model'}:
        main_card = make_scene_card((540, 430), seed + 3)
        support_left = make_scene_card((220, 238), seed + 11)
        support_right = make_scene_card((220, 238), seed + 17)
        _fallback_paste_card(canvas, main_card, (242, 152), shadow_alpha=44)
        _fallback_paste_card(canvas, support_left, (84, 622), shadow_alpha=32, shadow_offset=(0, 10))
        _fallback_paste_card(canvas, support_right, (720, 622), shadow_alpha=32, shadow_offset=(0, 10))
        draw.line((246, 618, 352, 502), fill=accent, width=7)
        draw.line((778, 618, 672, 502), fill=accent2, width=7)
        for point in [(246, 618), (352, 502), (672, 502), (778, 618)]:
            draw.ellipse((point[0] - 12, point[1] - 12, point[0] + 12, point[1] + 12), fill=panel_fill, outline=accent, width=3)
    else:
        hero_card = make_scene_card((760, 650), seed + 9)
        _fallback_paste_card(canvas, hero_card, (132, 188), shadow_alpha=46)
        draw.rounded_rectangle((102, 140, 246, 192), radius=24, fill=panel_fill)
        draw.ellipse((128, 160, 142, 174), fill=accent)
        draw.rounded_rectangle((736, 808, 920, 868), radius=28, fill=panel_fill)
        draw.rounded_rectangle((772, 834, 886, 844), radius=5, fill=line_color)

    return canvas.convert('RGB')


def _fallback_generate_vector_image_meta(idea, slide, index, image_options=None, hosted_error=''):
    slide = slide or {}
    style_key = _fallback_clean_text((image_options or {}).get('style'), 'deck-illustration').lower()
    palette = _fallback_image_palette(style_key)
    slide_type = _fallback_clean_text(slide.get('type'), 'story').lower()
    seed = _fallback_hash_seed(idea, slide.get('title'), slide.get('subtitle'), slide_type, style_key, index)
    suggestion = _fallback_clean_text(slide.get('visual_suggestion'), slide.get('title') or 'Concept illustration')
    topic = _fallback_topic_bucket(idea, slide)
    style_family = _fallback_style_family(style_key)
    model_label, repo_id = _fallback_style_model_meta(style_key)

    if style_family == 'abstract':
        illustration = _fallback_build_abstract_illustration(topic, slide_type, palette, seed)
    elif style_family == 'presentation':
        illustration = _fallback_build_scene_visual(topic, slide_type, palette, seed, variant='presentation')
    elif style_family == 'animated':
        illustration = _fallback_build_scene_visual(topic, slide_type, palette, seed, variant='animated')
    elif style_family == 'cartoon':
        illustration = _fallback_build_scene_visual(topic, slide_type, palette, seed, variant='cartoon')
    else:
        illustration = _fallback_build_editorial_illustration(topic, slide_type, palette, seed)

    return {
        'image_url': _fallback_image_to_data_url(illustration, fmt='JPEG', quality=92),
        'image_prompt': suggestion,
        'image_model': model_label,
        'image_repo_id': repo_id,
        'image_status': 'generated',
        'image_error': hosted_error,
    }


def _fallback_generate_slide_image_meta(idea, slide, index, image_options=None):
    model_key = _fallback_clean_text((image_options or {}).get('model_key'), FALLBACK_IMAGE_MODEL_KEY).lower()
    if model_key == 'ventureos-editorial':
        return _best_available_fallback_image_meta(
            idea,
            slide,
            index,
            image_options=image_options,
        )
    try:
        return _fallback_fetch_hosted_image_meta(idea, slide, index, image_options=image_options)
    except Exception as exc:
        return _best_available_fallback_image_meta(
            idea,
            slide,
            index,
            image_options=image_options,
            hosted_error=str(exc),
        )


def _fallback_prepare_progressive_image_meta(idea, slide, index, model_key=None, image_options=None):
    slide = slide or {}
    style_key = _fallback_clean_text((image_options or {}).get('style'), 'deck-illustration').lower()
    variation_key = _fallback_clean_text((image_options or {}).get('variation_key'))
    model_variant = _fallback_model_variant(model_key)
    provider_candidates = _fallback_provider_candidates(model_key, slide.get('type'), index)
    prompt = _fallback_hosted_image_prompt(idea, slide, style_key, index=index)
    image_size = _fallback_hosted_image_size(slide, index, image_options=image_options)
    seed = _fallback_hash_seed(
        idea,
        slide.get('title'),
        slide.get('subtitle'),
        slide.get('type'),
        style_key,
        variation_key,
        model_key,
        index,
    )
    fallback_meta = _best_available_fallback_image_meta(
        idea=idea,
        slide=slide,
        index=index,
        image_options=image_options,
    )
    return {
        'image_url': '',
        'image_fallback_url': fallback_meta.get('image_url', ''),
        'image_fallback_model': fallback_meta.get('image_model', FALLBACK_VECTOR_MODEL_LABEL),
        'image_fallback_repo_id': fallback_meta.get('image_repo_id', 'builtin/ventureos-neon-editorial'),
        'image_prompt': prompt,
        'image_model': model_variant['label'],
        'image_repo_id': model_variant['repo_id'],
        'image_status': 'queued',
        'image_error': '',
        'image_provider_candidates': provider_candidates,
        'image_seed': seed,
        'image_render_size': image_size,
    }


def _fallback_prepare_slides_for_progressive_images(idea, slides, model_key=None, image_options=None):
    enriched = []
    selected_indices = sorted(_fallback_select_image_slide_indices(slides, image_options))
    selected_set = set(selected_indices)
    for index, slide in enumerate(slides or []):
        item = dict(slide or {})
        if index in selected_set:
            item.update(_fallback_prepare_progressive_image_meta(
                idea=idea,
                slide=item,
                index=index,
                model_key=model_key,
                image_options=image_options,
            ))
        else:
            item.update(_best_available_fallback_image_meta(
                idea=idea,
                slide=item,
                index=index,
                image_options=image_options,
            ))
        enriched.append(item)
    return enriched, selected_indices


def _fallback_enrich_slides_with_images(idea, slides, market_research=None, model_key=None, image_options=None):
    enriched = []
    selected_indices = sorted(_fallback_select_image_slide_indices(slides, image_options))
    selected_set = set(selected_indices)
    lead_meta = None

    if selected_indices:
        lead_index = selected_indices[0]
        lead_slide = dict((slides or [])[lead_index] or {})
        lead_meta = _fallback_generate_slide_image_meta(
            idea=idea,
            slide=lead_slide,
            index=lead_index,
            image_options=image_options,
        )

    for index, slide in enumerate(slides or []):
        item = dict(slide or {})
        if index in selected_set:
            if lead_meta and index == selected_indices[0]:
                item.update(lead_meta)
            else:
                use_derived_layout = bool(
                    lead_meta
                    and str(lead_meta.get('image_repo_id', '')).startswith('pollinations/')
                    and lead_meta.get('image_url')
                )
                if use_derived_layout:
                    try:
                        item.update(_fallback_derive_image_meta(
                            lead_meta,
                            item,
                            index,
                            image_options=image_options,
                        ))
                    except Exception:
                        item.update(_best_available_fallback_image_meta(
                            idea=idea,
                            slide=item,
                            index=index,
                            image_options=image_options,
                            hosted_error='Hosted lead image remix failed, so VentureOS used a slide-specific illustration.',
                        ))
                else:
                    item.update(_best_available_fallback_image_meta(
                        idea=idea,
                        slide=item,
                        index=index,
                        image_options=image_options,
                        hosted_error='Using slide-specific illustration because no strong lead image was available.',
                    ))
        else:
            item.setdefault('image_status', 'skipped')
        enriched.append(item)
    return enriched


DEFAULT_MODEL_KEY = os.getenv('VENTUREOS_IMAGE_MODEL', FALLBACK_IMAGE_MODEL_KEY).strip() or FALLBACK_IMAGE_MODEL_KEY
enrich_slides_with_images = _fallback_enrich_slides_with_images
get_image_coverage_options = _fallback_coverage_options
get_image_style_options = _fallback_style_options
get_supported_models = _fallback_supported_models

if not IMAGE_GENERATION_DISABLED:
    try:
        from tools.image_generation import (
            DEFAULT_MODEL_KEY as IMAGEGEN_DEFAULT_MODEL_KEY,
            enrich_slides_with_images as imagegen_enrich_slides_with_images,
            get_image_coverage_options as imagegen_get_image_coverage_options,
            get_image_style_options as imagegen_get_image_style_options,
            get_supported_models as imagegen_get_supported_models,
        )
        DEFAULT_MODEL_KEY = IMAGEGEN_DEFAULT_MODEL_KEY
        enrich_slides_with_images = imagegen_enrich_slides_with_images
        get_image_coverage_options = imagegen_get_image_coverage_options
        get_image_style_options = imagegen_get_image_style_options
        get_supported_models = imagegen_get_supported_models
        IMAGE_GENERATION_AVAILABLE = True
        IMAGE_GENERATION_IMPORT_ERROR = ''
        IMAGE_GENERATION_IS_FALLBACK = False
    except Exception as image_import_error:
        IMAGE_GENERATION_IMPORT_ERROR = str(image_import_error)
else:
    IMAGE_GENERATION_IMPORT_ERROR = 'Using built-in illustrated slide visuals in this live deployment.'


app = Flask(__name__)


def _resolve_reports_dir():
    repo_reports_dir = os.path.join(os.path.dirname(__file__), 'reports')
    if os.getenv('VERCEL'):
        temp_reports_dir = os.path.join(tempfile.gettempdir(), 'ventureos-reports')
        os.makedirs(temp_reports_dir, exist_ok=True)
        return temp_reports_dir
    try:
        os.makedirs(repo_reports_dir, exist_ok=True)
        return repo_reports_dir
    except OSError:
        temp_reports_dir = os.path.join(tempfile.gettempdir(), 'ventureos-reports')
        os.makedirs(temp_reports_dir, exist_ok=True)
        return temp_reports_dir


REPORTS_DIR = _resolve_reports_dir()


DECK_TEMPLATE_PRESETS = [
    {
        'template_id': 'editorial-midnight',
        'theme_name': 'Scholar Noir',
        'palette': {
            'primary': '#23262F',
            'secondary': '#2E313B',
            'accent': '#F4728A',
            'surface': '#333745',
            'text': '#F7EEF4',
        },
        'style_notes': [
            'Dark academic-editorial canvas with rose-coral accent discipline',
            'Split-layout storytelling with contextual illustration or chart panel',
            'Feels like a premium research deck or polished conference talk',
        ]
    },
    {
        'template_id': 'boardroom-ivory',
        'theme_name': 'Boardroom Ivory',
        'palette': {
            'primary': '#F6F2EA',
            'secondary': '#E9E3D8',
            'accent': '#1D4ED8',
            'surface': '#FFFFFF',
            'text': '#111827',
        },
        'style_notes': [
            'Consulting-style light canvas with precise alignment',
            'Minimal borders, disciplined spacing, executive clarity',
            'Feels like a premium strategy memo turned into slides',
        ]
    },
    {
        'template_id': 'kinetic-ember',
        'theme_name': 'Kinetic Ember',
        'palette': {
            'primary': '#111827',
            'secondary': '#1F2937',
            'accent': '#F97316',
            'surface': '#18212F',
            'text': '#FFF7ED',
        },
        'style_notes': [
            'Dark executive canvas with warmer highlight energy',
            'Bold accents reserved for proof points and movement cues',
            'Feels more keynote-driven and launch-oriented',
        ]
    },
    {
        'template_id': 'atlas-sapphire',
        'theme_name': 'Atlas Sapphire',
        'palette': {
            'primary': '#081A3A',
            'secondary': '#0F274F',
            'accent': '#60A5FA',
            'surface': '#0E2243',
            'text': '#EFF6FF',
        },
        'style_notes': [
            'High-contrast cobalt system for data-heavy investor stories',
            'Sharp hierarchy with crisp metric framing',
            'Feels analytical, premium, and globally scalable',
        ]
    },
    {
        'template_id': 'monochrome-luxe',
        'theme_name': 'Monochrome Luxe',
        'palette': {
            'primary': '#111111',
            'secondary': '#1E1E1E',
            'accent': '#D4B483',
            'surface': '#202020',
            'text': '#F5F5F4',
        },
        'style_notes': [
            'Luxury monochrome composition with warm metallic accent',
            'Large type, sparse copy, elevated investor-brand mood',
            'Feels premium, polished, and slightly fashion-editorial',
        ]
    },
    {
        'template_id': 'forest-venture',
        'theme_name': 'Forest Venture',
        'palette': {
            'primary': '#0D1F14',
            'secondary': '#152B1D',
            'accent': '#4ADE80',
            'surface': '#112018',
            'text': '#F0FDF4',
        },
        'style_notes': [
            'Deep forest canvas with electric green proof points',
            'Evokes sustainability, growth, and long-horizon thinking',
            'Strong for impact, climate, health, and B-corp stories',
        ]
    },
    {
        'template_id': 'crimson-capital',
        'theme_name': 'Crimson Capital',
        'palette': {
            'primary': '#1A0A0A',
            'secondary': '#2A1010',
            'accent': '#F87171',
            'surface': '#220D0D',
            'text': '#FFF1F2',
        },
        'style_notes': [
            'Bold dark-red canvas — confident, high-conviction energy',
            'Accent drives urgency and decisive proof point emphasis',
            'Ideal for fintech, security, and high-stakes market narratives',
        ]
    },
    {
        'template_id': 'arctic-clarity',
        'theme_name': 'Arctic Clarity',
        'palette': {
            'primary': '#F8FAFF',
            'secondary': '#EEF2FB',
            'accent': '#3B82F6',
            'surface': '#FFFFFF',
            'text': '#0F172A',
        },
        'style_notes': [
            'Crisp polar white canvas — airy, trustworthy, effortlessly clean',
            'Blue accent anchors data, CTAs, and hierarchy with precision',
            'Feels like a top-tier SaaS product deck or Series A raise',
        ]
    },
    {
        'template_id': 'violet-epoch',
        'theme_name': 'Violet Epoch',
        'palette': {
            'primary': '#0F0A1E',
            'secondary': '#18102E',
            'accent': '#A78BFA',
            'surface': '#130D26',
            'text': '#F5F3FF',
        },
        'style_notes': [
            'Deep violet — distinctly premium for AI, crypto, and frontier tech',
            'Lavender accent evokes intelligence, creativity, and depth',
            'Feels visionary, slightly futuristic, and quietly luxurious',
        ]
    },
]

DEFAULT_DECK_THEME = DECK_TEMPLATE_PRESETS[0]


def _get_template_preset(template_id=None, exclude_template_id=None):
    if template_id:
        for preset in DECK_TEMPLATE_PRESETS:
            if preset['template_id'] == template_id:
                return preset
    candidates = [
        preset for preset in DECK_TEMPLATE_PRESETS
        if preset['template_id'] != exclude_template_id
    ] or DECK_TEMPLATE_PRESETS
    return random.choice(candidates)


def _clean_text(value, fallback=''):
    if value is None:
        return fallback
    text = str(value).strip()
    return text or fallback


def _clean_list(value, limit=5):
    if not isinstance(value, list):
        return []
    items = []
    for item in value:
        text = _clean_text(item)
        if text:
            items.append(text)
        if len(items) >= limit:
            break
    return items


def _normalize_hex(value, fallback):
    text = _clean_text(value, fallback).replace('#', '')
    if re.fullmatch(r'[0-9a-fA-F]{6}', text):
        return f"#{text.upper()}"
    return fallback


def _default_layout(slide_type: str, index: int) -> str:
    st = (slide_type or '').lower()
    if st in {'hook', 'opening', 'title', 'cover'}:
        return 'cover'
    if st in {'market', 'impact', 'value', 'traction', 'proof', 'data'}:
        return 'metrics'
    if st in {'how_it_works', 'business_model', 'process', 'roadmap'}:
        return 'roadmap'
    if st in {'vision', 'future', 'ask', 'call_to_action', 'cta'}:
        return 'closing'
    if st in {'competition', 'competitive_landscape'}:
        return 'comparison'
    return 'spotlight' if index else 'cover'


def _normalize_stats(raw_slide: dict) -> list:
    stats = raw_slide.get('stats')
    normalized = []
    if isinstance(stats, list):
        for item in stats[:3]:
            if not isinstance(item, dict):
                continue
            value = _clean_text(item.get('value') or item.get('num'))
            label = _clean_text(item.get('label'))
            if value or label:
                normalized.append({'value': value, 'label': label})
    else:
        legacy = [
            (raw_slide.get('stat1_num'), raw_slide.get('stat1_label')),
            (raw_slide.get('stat2_num'), raw_slide.get('stat2_label')),
            (raw_slide.get('stat3_num'), raw_slide.get('stat3_label')),
        ]
        for value, label in legacy:
            value = _clean_text(value)
            label = _clean_text(label)
            if value or label:
                normalized.append({'value': value, 'label': label})
    return normalized[:3]


def _normalize_slide(raw_slide: dict, index: int) -> dict:
    raw_slide = raw_slide or {}
    animation = raw_slide.get('animation_plan') or {}
    content = _clean_list(
        raw_slide.get('content')
        or raw_slide.get('points')
        or raw_slide.get('key_points')
        or [],
        limit=5,
    )
    sequence = _clean_list(
        animation.get('sequence')
        or raw_slide.get('sequence')
        or [],
        limit=4,
    )
    slide_type = _clean_text(
        raw_slide.get('type') or raw_slide.get('slide_type'),
        'story'
    ).lower()

    return {
        'slide_number': raw_slide.get('slide_number') or index + 1,
        'type': slide_type,
        'layout': _clean_text(
            raw_slide.get('layout'),
            _default_layout(slide_type, index)
        ).lower(),
        'title': _clean_text(
            raw_slide.get('title')
            or raw_slide.get('headline')
            or raw_slide.get('slide_title'),
            f"Slide {index + 1}",
        ),
        'subtitle': _clean_text(
            raw_slide.get('subtitle')
            or raw_slide.get('subheadline')
        ),
        'objective': _clean_text(raw_slide.get('objective')),
        'content': content,
        'stats': _normalize_stats(raw_slide),
        'visual_suggestion': _clean_text(
            raw_slide.get('visual_suggestion')
            or raw_slide.get('visual')
            or raw_slide.get('chart_suggestion')
        ),
        'image_url': _clean_text(
            raw_slide.get('image_url')
            or raw_slide.get('visual_image_url')
        ),
        'image_prompt': _clean_text(raw_slide.get('image_prompt')),
        'image_model': _clean_text(raw_slide.get('image_model')),
        'image_repo_id': _clean_text(raw_slide.get('image_repo_id')),
        'image_status': _clean_text(raw_slide.get('image_status')),
        'image_error': _clean_text(raw_slide.get('image_error')),
        'design_notes': _clean_text(raw_slide.get('design_notes')),
        'animation_plan': {
            'entry': _clean_text(animation.get('entry'), 'Fade'),
            'sequence': sequence or ['Title', 'Core message', 'Proof point'],
            'transition': _clean_text(animation.get('transition'), 'Smooth fade'),
            'emphasis': _clean_text(animation.get('emphasis')),
        }
    }


def _extract_json_payload(text: str):
    text = text.strip()
    text = re.sub(r'```json', '', text)
    text = re.sub(r'```', '', text)
    text = text.strip()

    start_obj = text.find('{')
    end_obj = text.rfind('}') + 1
    if start_obj != -1 and end_obj > start_obj:
        return json.loads(text[start_obj:end_obj])

    start_arr = text.find('[')
    end_arr = text.rfind(']') + 1
    if start_arr != -1 and end_arr > start_arr:
        return json.loads(text[start_arr:end_arr])

    raise ValueError('No JSON object or array found in model response')


def _normalize_deck_payload(payload, idea: str, template_preset=None) -> dict:
    payload = {'slides': payload} if isinstance(
        payload, list) else (payload or {})
    design = payload.get('design_system') or {}
    palette = design.get('palette') or {}
    template_preset = template_preset or DEFAULT_DECK_THEME

    slides = payload.get('slides') or payload.get('deck') or []
    slides = [
        _normalize_slide(slide, i)
        for i, slide in enumerate(slides[:12])
    ]

    return {
        'presentation_title': _clean_text(
            payload.get('presentation_title'),
            idea[:72] or 'Investor Presentation'
        ),
        'presentation_subtitle': _clean_text(
            payload.get('presentation_subtitle'),
            'Premium investor-ready story deck'
        ),
        'design_system': {
            'template_id': _clean_text(
                design.get('template_id'),
                template_preset['template_id']
            ),
            'theme_name': _clean_text(
                design.get('theme_name')
                or template_preset.get('theme_name'),
                DEFAULT_DECK_THEME['theme_name']
            ),
            'palette': {
                'primary': _normalize_hex(
                    template_preset['palette'].get('primary'),
                    DEFAULT_DECK_THEME['palette']['primary']
                ),
                'secondary': _normalize_hex(
                    template_preset['palette'].get('secondary'),
                    DEFAULT_DECK_THEME['palette']['secondary']
                ),
                'accent': _normalize_hex(
                    template_preset['palette'].get('accent'),
                    DEFAULT_DECK_THEME['palette']['accent']
                ),
                'surface': _normalize_hex(
                    template_preset['palette'].get('surface'),
                    DEFAULT_DECK_THEME['palette']['surface']
                ),
                'text': _normalize_hex(
                    template_preset['palette'].get('text'),
                    DEFAULT_DECK_THEME['palette']['text']
                ),
            },
            'style_notes': _clean_list(
                template_preset.get('style_notes')
                or design.get('style_notes'),
                limit=3
            ) or DEFAULT_DECK_THEME['style_notes'],
        },
        'slides': slides,
    }


def _build_local_deck_payload(idea: str, context: dict, template_preset=None) -> dict:
    template_preset = template_preset or DEFAULT_DECK_THEME
    context = context or {}
    market = context.get('market_research', {}) or {}
    score = context.get('scorecard', {}) or {}
    product = context.get('product_strategy', {}) or {}
    pitch = context.get('pitch', {}) or {}
    monetization = product.get('monetization', []) if isinstance(product, dict) else []
    deck = pitch.get('deck', []) if isinstance(pitch, dict) else []

    default_titles = [
        'The Category Is Breaking',
        'The Problem Is Structural',
        'Why Now Matters',
        'A Better Operating Layer',
        'How It Works',
        'Value Delivered Fast',
        'Proof That It Resonates',
        'A Scalable Revenue Engine',
        'What Comes Next',
        'Join The Build',
    ]
    type_map = [
        'hook',
        'problem',
        'stakes',
        'solution',
        'how_it_works',
        'impact',
        'proof',
        'business_model',
        'vision',
        'call_to_action',
    ]
    default_visuals = [
        'Cinematic hero image with one clear subject and strong negative space.',
        'Tension visual showing the current workflow friction in a realistic setting.',
        'Market context image that reinforces scale, urgency, and timing.',
        'Product-led editorial scene showing the solution in action.',
        'Step-based visual that suggests a clean, modern workflow.',
        'Outcome-led image with strong whitespace and premium composition.',
        'Proof-oriented scene that feels credible, sharp, and investor-ready.',
        'Business model image with an operating-system feel and product polish.',
        'Forward-looking aspirational image that hints at the long-term vision.',
        'Confident closing visual that supports a decisive call to action.',
    ]
    default_design_notes = [
        'Use whitespace aggressively and keep the title dominant.',
        'Keep the composition sparse so the pain reads instantly.',
        'Lead with one proof point before explanation.',
        'Balance a bold statement with one supporting visual.',
        'Use structured sequencing and clear progression.',
        'Make the slide feel premium, not dashboard-like.',
        'Keep the tone analytical and confidence-building.',
        'Use clean hierarchy with no unnecessary clutter.',
        'Let the future-state visual breathe with strong rhythm.',
        'End with a clear investor ask and restrained confidence.',
    ]
    default_objectives = [
        'Hook the audience with a sharp opening.',
        'Show the real friction clearly.',
        'Frame the urgency and timing.',
        'Position the product as the answer.',
        'Explain the workflow simply.',
        'Show measurable upside quickly.',
        'Give confidence through proof and traction.',
        'Show how the business compounds.',
        'Extend the story into the future.',
        'Close with a crisp investor ask.',
    ]

    slides = []
    for index in range(10):
        source = deck[index] if index < len(deck) and isinstance(deck[index], dict) else {}
        title = _clean_text(source.get('title'), default_titles[index])
        points = _clean_list(
            source.get('key_points')
            or source.get('content')
            or source.get('points')
            or [],
            limit=4,
        )
        slide_type = type_map[index]
        stats = []

        if index == 2:
            if market.get('market_size'):
                stats.append({'value': _clean_text(market.get('market_size')), 'label': 'Market Size'})
            if market.get('growth_rate'):
                stats.append({'value': _clean_text(market.get('growth_rate')), 'label': 'Growth Rate'})
        if index == 5 and score.get('total'):
            stats.append({'value': f"{score.get('total')}/100", 'label': 'Fundability Score'})
            if score.get('verdict'):
                stats.append({'value': _clean_text(score.get('verdict')), 'label': 'Current Verdict'})
        if index == 7 and isinstance(monetization, list):
            for item in monetization[:2]:
                if isinstance(item, dict) and item.get('model'):
                    stats.append({'value': _clean_text(item.get('model')), 'label': 'Revenue Model'})

        slides.append({
            'slide_number': index + 1,
            'type': slide_type,
            'layout': _default_layout(slide_type, index),
            'title': title,
            'subtitle': _clean_text(
                source.get('subtitle')
                or (points[0] if points else '')
                or market.get('opportunity_summary')
                or market.get('pain_point')
                or idea
            ),
            'objective': default_objectives[index],
            'content': points or _clean_list([
                market.get('pain_point'),
                market.get('target_customer'),
                score.get('biggest_strength'),
                score.get('biggest_risk'),
            ], limit=4),
            'stats': stats,
            'visual_suggestion': default_visuals[index],
            'design_notes': default_design_notes[index],
            'animation_plan': {
                'entry': ['Fade', 'Slide Up', 'Zoom'][index % 3],
                'sequence': _clean_list(points, limit=3) or ['Title', 'Core message', 'Proof point'],
                'transition': ['Smooth fade', 'Directional push', 'Soft zoom'][index % 3],
                'emphasis': 'Subtle stat emphasis' if index in {5, 7} else '',
            }
        })

    return _normalize_deck_payload({
        'presentation_title': idea[:72] or 'Investor Presentation',
        'presentation_subtitle': 'Locally generated from VentureOS analysis',
        'design_system': {
            'template_id': template_preset['template_id'],
            'theme_name': template_preset['theme_name'],
            'palette': template_preset['palette'],
            'style_notes': template_preset['style_notes'],
        },
        'slides': slides,
    }, idea, template_preset)


# ── MAIN PAGE ──────────────────────────────────────────────────────────────
@app.route('/')
def index():
    return render_template('index.html')


# ── ANALYZE — streams SSE from all agents ──────────────────────────────────
@app.route('/analyze', methods=['POST'])
def analyze():
    data = request.json
    idea = data.get('idea', '')
    if not idea:
        return jsonify({'error': 'No idea provided'}), 400

    def generate():
        try:
            for chunk in run_ventureOS_stream(idea):
                yield chunk
        except Exception as e:
            yield f"data: {json.dumps({'event': 'error', 'message': str(e)})}\n\n"

    return Response(
        stream_with_context(generate()),
        mimetype='text/event-stream',
        headers={
            'Cache-Control': 'no-cache',
            'X-Accel-Buffering': 'no',
            'Connection': 'keep-alive'
        }
    )


# ── PROTOTYPE ──────────────────────────────────────────────────────────────
@app.route('/prototype', methods=['POST'])
def prototype():
    data = request.json
    idea = data.get('idea', '')
    product_strategy = data.get('product_strategy', {})
    context = data.get('context', {})
    seed = data.get('seed', None)
    if not idea:
        return jsonify({'error': 'No idea provided'}), 400
    try:
        llm = get_llm()
        html, theme_name = run_prototype_generator(
            idea, product_strategy, llm, seed)
        return jsonify({'html': html, 'theme': theme_name})
    except Exception as e:
        import traceback
        return jsonify({'error': str(e), 'trace': traceback.format_exc()}), 500


# ── CHAT ───────────────────────────────────────────────────────────────────
@app.route('/chat', methods=['POST'])
def chat():
    data = request.json
    idea = data.get('idea', '')
    context = data.get('context', {})
    messages = data.get('messages', [])
    user_message = data.get('user_message', '')

    if not user_message:
        return jsonify({'error': 'No message provided'}), 400

    try:
        llm = get_llm()
        system_prompt = f"""You are VentureOS, an expert startup advisor and co-founder AI.

The user has already analyzed this startup idea: "{idea}"

Here is the full analysis context:
- Market Research: {json.dumps(context.get('market_research', {}))}
- Competitor Analysis: {json.dumps(context.get('competitor_analysis', {}))}
- Product Strategy: {json.dumps(context.get('product_strategy', {}))}
- Fundability Score: {json.dumps(context.get('scorecard', {}))}

Be specific, honest, and actionable. Keep responses concise — 2-4 paragraphs max."""

        from langchain_core.messages import HumanMessage, SystemMessage, AIMessage
        lc_messages = [SystemMessage(content=system_prompt)]
        for msg in messages:
            if msg['role'] == 'user':
                lc_messages.append(HumanMessage(content=msg['content']))
            elif msg['role'] == 'assistant':
                lc_messages.append(AIMessage(content=msg['content']))
        lc_messages.append(HumanMessage(content=user_message))

        response = llm.invoke(lc_messages)
        return jsonify({'reply': response.content})

    except Exception as e:
        import traceback
        return jsonify({'error': str(e), 'trace': traceback.format_exc()}), 500


# ── PIVOT ──────────────────────────────────────────────────────────────────
@app.route('/pivot', methods=['POST'])
def pivot():
    data = request.json
    idea = data.get('idea', '')
    context = data.get('context', {})
    if not idea:
        return jsonify({'error': 'No idea provided'}), 400
    try:
        llm = get_llm()
        result = run_pivot_suggester(
            idea,
            context.get('market_research', {}),
            context.get('competitor_analysis', {}),
            context.get('scorecard', {}),
            llm
        )
        return jsonify(result)
    except Exception as e:
        import traceback
        return jsonify({'error': str(e), 'trace': traceback.format_exc()}), 500


# ── SLIDES — proxy to avoid CORS ───────────────────────────────────────────
@app.route('/slides', methods=['POST'])
def generate_slides():
    data = request.json
    idea = data.get('idea', '')
    context = data.get('context', {})
    template_id = data.get('template_id')
    previous_template_id = data.get('previous_template_id')
    generate_images = bool(data.get('generate_images', True))
    image_options = data.get('image_options', {})
    supported_models = get_supported_models()
    supported_model_keys = {model.get('key') for model in supported_models if model.get('key')}
    image_model = _clean_text(data.get('image_model'), DEFAULT_MODEL_KEY)
    if supported_model_keys and image_model not in supported_model_keys:
        image_model = DEFAULT_MODEL_KEY if DEFAULT_MODEL_KEY in supported_model_keys else next(iter(supported_model_keys))

    if not idea:
        return jsonify({'error': 'No idea provided'}), 400

    m = context.get('market_research', {})
    c = context.get('competitor_analysis', {})
    p = context.get('product_strategy', {})
    sc = context.get('scorecard', {})
    pitch = context.get('pitch', {})
    deck = pitch.get('deck', [])
    template_preset = _get_template_preset(
        template_id=template_id,
        exclude_template_id=previous_template_id,
    )

    prompt = f"""You are an expert Presentation Designer and Storytelling Strategist who creates premium, high-end, investor-grade presentations in the style of McKinsey, BCG, Apple keynote, and top startup pitch decks.

Create a premium 10-slide presentation structure for this startup. The deck must feel intentional, minimal, structured, and highly polished.

Selected visual template:
- Template ID: {template_preset['template_id']}
- Template name: {template_preset['theme_name']}
- Palette: {json.dumps(template_preset['palette'])}
- Style notes: {json.dumps(template_preset['style_notes'])}

Startup idea: "{idea}"
Market size: {m.get('market_size', '—')}
Growth rate: {m.get('growth_rate', '—')}
Target customer: {m.get('target_customer', '—')}
Pain point: {m.get('pain_point', '—')}
Opportunity summary: {m.get('opportunity_summary', '—')}
Competitor whitespace: {c.get('whitespace', '—')}
Fundability score: {sc.get('total', '—')}/100
Verdict: {sc.get('verdict', '—')}
Biggest strength: {sc.get('biggest_strength', '—')}
Biggest risk: {sc.get('biggest_risk', '—')}
Pitch outline: {json.dumps(deck[:10])}
Monetization: {json.dumps(p.get('monetization', [])[:3])}

Strict design rules:
- Premium, high-contrast, minimal layout
- Plenty of whitespace
- Titles short and powerful
- Minimal text: max 6-8 words per line
- Max 3-5 bullets per slide
- Avoid generic filler
- Recommend visuals, charts, icons, or imagery where useful
- Every slide must include a smooth, professional animation plan

Story flow:
1. Hook / Opening
2. Problem
3. Why it matters
4. Solution
5. How it works
6. Value / Impact
7. Proof / Data
8. Business model
9. Future / Vision
10. Call to Action

Return ONLY a valid JSON object with this exact structure:
{{
  "presentation_title": "Deck title",
  "presentation_subtitle": "Short contextual subtitle",
  "design_system": {{
    "theme_name": "Short premium theme name",
    "palette": {{
      "primary": "#0B1020",
      "secondary": "#182033",
      "accent": "#7DD3FC",
      "surface": "#121A2B",
      "text": "#F8FAFC"
    }},
    "style_notes": [
      "Deck-level design principle 1",
      "Deck-level design principle 2",
      "Deck-level design principle 3"
    ]
  }},
  "slides": [
    {{
      "slide_number": 1,
      "type": "hook|problem|stakes|solution|how_it_works|impact|proof|business_model|vision|call_to_action",
      "layout": "cover|spotlight|metrics|roadmap|comparison|closing",
      "title": "Short punchy title",
      "subtitle": "Optional support line",
      "objective": "What this slide is trying to achieve",
      "content": ["Bullet 1", "Bullet 2", "Bullet 3"],
      "stats": [
        {{"value": "$4.2B", "label": "Total Market"}},
        {{"value": "28%", "label": "Annual Growth"}}
      ],
      "visual_suggestion": "Specific visual, chart, or image direction",
      "design_notes": "Layout, alignment, spacing, hierarchy notes",
      "animation_plan": {{
        "entry": "Fade|Zoom|Slide Up",
        "sequence": ["What appears first", "What appears second", "What appears third"],
        "transition": "Fade|Push|Morph-like directional cut",
        "emphasis": "Optional subtle emphasis animation"
      }}
    }}
  ]
}}

Use real numbers from context whenever possible. Keep the tone confident, crisp, and concise."""

    try:
        generation_mode = 'llm'
        generation_notice = ''
        generation_error = ''
        try:
            llm = get_llm()
            from langchain_core.messages import HumanMessage
            response = llm.invoke([HumanMessage(content=prompt)])
            payload = _extract_json_payload(response.content)
            deck_payload = _normalize_deck_payload(payload, idea, template_preset)
        except Exception as llm_error:
            generation_mode = 'local'
            generation_error = str(llm_error)
            generation_notice = 'Using local deck structure because premium text generation is unavailable.'
            deck_payload = _build_local_deck_payload(idea, context, template_preset)

        images_enabled = bool(generate_images and IMAGE_GENERATION_AVAILABLE)
        progressive_hosted_images = bool(
            images_enabled
            and IMAGE_GENERATION_IS_FALLBACK
            and HOSTED_FALLBACK_ENABLED
            and _clean_text(image_model, FALLBACK_IMAGE_MODEL_KEY).lower() != 'ventureos-editorial'
        )
        pending_image_indices = []
        if images_enabled:
            if progressive_hosted_images:
                deck_payload['slides'], pending_image_indices = _fallback_prepare_slides_for_progressive_images(
                    idea=idea,
                    slides=deck_payload.get('slides', []),
                    model_key=image_model,
                    image_options=image_options,
                )
                fallback_notice = (
                    'Generating unique premium visuals progressively so each slide receives its own topic-aware image.'
                )
                generation_notice = f"{generation_notice} {fallback_notice}".strip() if generation_notice else fallback_notice
            else:
                enriched_options = dict(image_options or {})
                enriched_options['model_key'] = image_model
                deck_payload['slides'] = enrich_slides_with_images(
                    idea=idea,
                    slides=deck_payload.get('slides', []),
                    market_research=m,
                    model_key=image_model,
                    image_options=enriched_options,
                )
                if IMAGE_GENERATION_IS_FALLBACK:
                    fallback_notice = (
                        'Using VentureOS built-in scene visuals so every selected slide gets an immediate topic-aware image while premium renders load.'
                        if _clean_text(image_model, FALLBACK_IMAGE_MODEL_KEY).lower() == 'ventureos-editorial'
                        else (
                            'Using hosted scene image generation for topic-aware deck visuals, with a built-in slide-specific fallback if needed.'
                            if HOSTED_FALLBACK_ENABLED
                            else 'Using built-in illustrated scenes so the live deck and PPT still include visuals.'
                        )
                    )
                    generation_notice = f"{generation_notice} {fallback_notice}".strip() if generation_notice else fallback_notice
        deck_payload['image_generation'] = {
            'requested': generate_images,
            'enabled': images_enabled,
            'available': IMAGE_GENERATION_AVAILABLE,
            'message': IMAGE_GENERATION_IMPORT_ERROR,
            'default_model': image_model,
            'selected_model': image_model,
            'supported_models': supported_models,
            'style_options': get_image_style_options(),
            'coverage_options': get_image_coverage_options(),
            'selected_style': _clean_text(
                (image_options or {}).get('style'),
                'deck-illustration'
            ),
            'selected_coverage': _clean_text(
                (image_options or {}).get('coverage'),
                'key-slides'
            ),
            'progressive_enabled': progressive_hosted_images,
            'pending_indices': pending_image_indices,
        }
        deck_payload['generation_mode'] = generation_mode
        if generation_notice:
            deck_payload['generation_notice'] = generation_notice
        if generation_error:
            deck_payload['generation_error'] = generation_error
        return jsonify(deck_payload)
    except Exception as e:
        import traceback
        return jsonify({'error': str(e), 'trace': traceback.format_exc()}), 500


@app.route('/slides/image', methods=['POST'])
def generate_slide_image():
    data = request.json or {}
    prompt = _clean_text(data.get('prompt'))
    provider_candidates = data.get('provider_candidates') or []
    if not isinstance(provider_candidates, list):
        provider_candidates = []
    provider_candidates = [_clean_text(item) for item in provider_candidates if _clean_text(item)]
    provider_candidates = provider_candidates[:4] or ['flux']
    seed = int(data.get('seed') or 1)
    image_size = int(data.get('size') or 512)

    if not prompt:
        return jsonify({'error': 'No image prompt provided'}), 400
    if not GOOGLE_IMAGE_ENABLED and not HOSTED_FALLBACK_ENABLED:
        return jsonify({'error': 'No image provider is enabled'}), 503

    errors = []
    for provider_model in provider_candidates:
        try:
            cache_key = f"{provider_model}|{image_size}|{seed}|{prompt}"
            image_url = _fallback_fetch_provider_image_data_url(
                prompt=prompt,
                provider_model=provider_model,
                seed=seed,
                image_size=image_size,
                cache_key=cache_key,
            )
            return jsonify({
                'image_url': image_url,
                'image_model': _fallback_provider_label(provider_model),
                'image_repo_id': _fallback_provider_repo(provider_model),
                'image_status': 'generated',
            })
        except Exception as image_error:
            errors.append({
                'provider_model': provider_model,
                'message': str(image_error),
            })

    return jsonify({
        'error': 'All hosted image providers failed',
        'attempts': errors,
    }), 502


# ── DOWNLOAD PPTX ──────────────────────────────────────────────────────────
# 8 design systems — each has completely different visual structure per slide

PPTX_DESIGN_SYSTEMS = {
    # ── 1. EDITORIAL MIDNIGHT: left accent bar, right panel, serif titles ──
    'editorial-midnight': {'style': 'editorial', 'FH': 'Georgia',          'FB': 'Calibri'},
    'boardroom-ivory':    {'style': 'boardroom',  'FH': 'Cambria',          'FB': 'Calibri'},
    # ── 2. KINETIC: full-bleed diagonal geometry, bold impact layout ──────
    'kinetic-ember':      {'style': 'kinetic',    'FH': 'Arial Black',      'FB': 'Arial'},
    'crimson-capital':    {'style': 'kinetic',    'FH': 'Georgia',          'FB': 'Calibri'},
    # ── 3. ATLAS: data-forward, left sidebar accent, metric-heavy ─────────
    'atlas-sapphire':     {'style': 'atlas',      'FH': 'Georgia',          'FB': 'Calibri'},
    'violet-epoch':       {'style': 'atlas',      'FH': 'Georgia',          'FB': 'Calibri'},
    # ── 4. LUXE: centred luxury, large serif, minimal lines ───────────────
    'monochrome-luxe':    {'style': 'luxe',       'FH': 'Palatino Linotype', 'FB': 'Calibri'},
    'charcoal-gold':      {'style': 'luxe',       'FH': 'Palatino Linotype', 'FB': 'Calibri'},
    # ── 5. BOARDROOM: light canvas, thin rules, consulting aesthetic ───────
    'arctic-clarity':     {'style': 'boardroom',  'FH': 'Trebuchet MS',     'FB': 'Calibri'},
    # ── 6. FOREST: organic, bottom-anchored stats, nature motif ───────────
    'forest-venture':     {'style': 'forest',     'FH': 'Georgia',          'FB': 'Calibri'},
}


@app.route('/download_pptx', methods=['POST'])
def download_pptx():
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

        data = request.json
        idea = data.get('idea', 'ventureos')
        raw_slides = data.get('slides') or data.get('deck') or []
        slides_data = [_normalize_slide(s, i)
                       for i, s in enumerate(raw_slides)]
        if not slides_data:
            return jsonify({'error': 'No slides provided'}), 400

        design = data.get('design_system') or DEFAULT_DECK_THEME
        palette = design.get('palette') or DEFAULT_DECK_THEME['palette']
        tid = design.get('template_id', DEFAULT_DECK_THEME['template_id'])
        theme_name = design.get('theme_name', 'VentureOS')

        # ── Colour helpers ─────────────────────────────────────────────────
        def _h2r(val, fb):
            v = _normalize_hex(val, fb).replace('#', '')
            return RGBColor(int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))

        def _luma(val, fb='#000000'):
            v = _normalize_hex(val, fb).replace('#', '')
            return 0.299*int(v[0:2], 16)+0.587*int(v[2:4], 16)+0.114*int(v[4:6], 16)

        C_BG = _h2r(palette.get('primary'),
                    DEFAULT_DECK_THEME['palette']['primary'])
        C_BG2 = _h2r(palette.get('secondary'),
                     DEFAULT_DECK_THEME['palette']['secondary'])
        C_SURF = _h2r(palette.get('surface'),
                      DEFAULT_DECK_THEME['palette']['surface'])
        C_ACC = _h2r(palette.get('accent'),
                     DEFAULT_DECK_THEME['palette']['accent'])
        C_TEXT = _h2r(palette.get('text'),
                      DEFAULT_DECK_THEME['palette']['text'])
        is_light = _luma(palette.get('primary', '#000')) > 160
        C_DIM = RGBColor(0x44, 0x55, 0x66) if is_light else RGBColor(
            0xB0, 0xC4, 0xD8)
        C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        C_BLACK = RGBColor(0x0F, 0x0F, 0x0F)

        def _blend(a, b, t):
            return RGBColor(int(a[0]+(b[0]-a[0])*t), int(a[1]+(b[1]-a[1])*t), int(a[2]+(b[2]-a[2])*t))

        def _acc(t): return _blend(C_BG, C_ACC, t)
        def _sur(t): return _blend(C_BG, C_SURF, t)
        def _bg2(t): return _blend(C_BG, C_BG2, t)

        def _acc_light(t): return _blend(
            C_WHITE if is_light else C_BG, C_ACC, t)

        ds = PPTX_DESIGN_SYSTEMS.get(
            tid, {'style': 'editorial', 'FH': 'Georgia', 'FB': 'Calibri'})
        STYLE = ds['style']
        FH = ds['FH']
        FB = ds['FB']

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        BLANK = prs.slide_layouts[6]

        # ── Draw primitives ────────────────────────────────────────────────
        def txt(sl, text, L, T, W, H, sz, col, bold=False, italic=False,
                align=PP_ALIGN.LEFT, fn=None):
            if not text:
                return
            tb = sl.shapes.add_textbox(
                Inches(L), Inches(T), Inches(W), Inches(H))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = align
            r = p.add_run()
            r.text = str(text)
            r.font.size = Pt(sz)
            r.font.color.rgb = col
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = fn or FB

        def rect(sl, L, T, W, H, fill, border=None, bw=0.5, rnd=False):
            st = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rnd else MSO_AUTO_SHAPE_TYPE.RECTANGLE
            sh = sl.shapes.add_shape(st, Inches(
                L), Inches(T), Inches(W), Inches(H))
            sh.fill.solid()
            sh.fill.fore_color.rgb = fill
            if border:
                sh.line.color.rgb = border
                sh.line.width = Pt(bw)
            else:
                sh.line.fill.background()
            return sh

        def oval(sl, L, T, W, H, fill):
            sh = sl.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL,
                                     Inches(L), Inches(T), Inches(W), Inches(H))
            sh.fill.solid()
            sh.fill.fore_color.rgb = fill
            sh.line.fill.background()
            return sh

        def footer(sl, idx, total):
            rect(sl, 0, 7.22, 13.33, 0.28, _bg2(0.8))
            rect(sl, 0, 7.22, 13.33, 0.016, C_ACC)
            txt(sl, 'VentureOS', 0.40, 7.28, 3.0, 0.22, 7, C_DIM)
            txt(sl, theme_name,  4.66, 7.28, 4.0,
                0.22, 7, C_DIM, align=PP_ALIGN.CENTER)
            txt(sl, f'{idx+1:02d} / {total:02d}', 9.50, 7.28, 3.3, 0.22,
                7, C_DIM, bold=True, align=PP_ALIGN.RIGHT)

        total = len(slides_data)

        for i, sd in enumerate(slides_data):
            sl = prs.slides.add_slide(BLANK)
            layout = sd.get('layout', 'spotlight')
            title = sd.get('title') or f'Slide {i+1}'
            sub = sd.get('subtitle') or ''
            pts = (sd.get('content') or sd.get('points') or [])[:5]
            stats = (sd.get('stats') or [])[:3]
            obj = sd.get('objective') or ''
            vis = _presentation_visual_summary(idea, sd, i)
            kicker = sd.get('type', 'story').replace('_', ' ').upper()

            sl.background.fill.solid()
            sl.background.fill.fore_color.rgb = C_BG

            # ══════════════════════════════════════════════════════════════════
            # STYLE A — EDITORIAL  (left accent stripe + right panel)
            # DNA: vertical left bar, right glassmorphic panel, serif headlines
            # ══════════════════════════════════════════════════════════════════
            if STYLE == 'editorial':
                if layout == 'cover':
                    rect(sl, 0, 0, 5.6, 7.5, _bg2(0.88))
                    rect(sl, 5.6, 0, 7.73, 7.5, _bg2(0.35))
                    oval(sl, 5.2, -1.8, 8.0, 8.0, _acc(0.16))
                    oval(sl, 8.5,  3.8, 5.5, 5.5, _acc(0.10))
                    rect(sl, 0, 0, 0.08, 7.5, C_ACC)
                    rect(sl, 0.55, 0.55, 3.0, 0.32, _acc(0.28), rnd=True)
                    txt(sl, kicker, 0.72, 0.59, 2.8, 0.26, 8, C_ACC, bold=True)
                    txt(sl, title, 0.55, 1.10, 4.8, 2.80,
                        38, C_TEXT, bold=True, fn=FH)
                    if sub:
                        txt(sl, sub, 0.55, 3.85, 4.6,
                            0.80, 13, C_DIM, italic=True)
                    rect(sl, 0.55, 4.72, 3.8, 0.04, _acc(0.60))
                    txt(sl, idea.title()[:40], 0.55,
                        4.92, 4.5, 0.40, 10.5, C_DIM)
                    if stats:
                        txt(sl, stats[0].get('value', ''), 6.2, 1.80, 7.0, 1.90,
                            70, _acc(0.88), bold=True, fn=FH, align=PP_ALIGN.CENTER)
                        txt(sl, stats[0].get('label', '').upper(), 6.2, 4.00, 7.0, 0.38,
                            10, C_DIM, bold=True, align=PP_ALIGN.CENTER)
                    else:
                        txt(sl, vis or 'The future starts here.', 6.2, 2.80, 7.0,
                            1.40, 15, C_DIM, italic=True, align=PP_ALIGN.CENTER)
                elif layout == 'metrics':
                    rect(sl, 0, 0, 13.33, 1.38, _bg2(0.88))
                    rect(sl, 0, 1.38, 13.33, 0.05, C_ACC)
                    rect(sl, 0, 0, 0.08, 1.38, C_ACC)
                    txt(sl, kicker, 0.52, 0.18, 6.0, 0.26, 8, C_ACC, bold=True)
                    txt(sl, title,  0.52, 0.50, 9.0, 0.78,
                        26, C_TEXT, bold=True, fn=FH)
                    if stats:
                        n = len(stats)
                        cw = (12.33 - 0.18*(n-1)) / n
                        for bi, st in enumerate(stats):
                            bx = 0.50 + bi*(cw+0.18)
                            rect(sl, bx, 1.65, cw, 3.30, _sur(0.75),
                                 border=_acc(0.25), bw=0.6, rnd=True)
                            rect(sl, bx, 1.65, cw, 0.12, C_ACC)
                            txt(sl, st.get('value', ''), bx+0.18, 1.92, cw-0.36, 1.60,
                                50, C_TEXT, bold=True, fn=FH, align=PP_ALIGN.CENTER)
                            txt(sl, st.get('label', '').upper(), bx+0.18, 3.62, cw-0.36,
                                0.34, 9, C_DIM, bold=True, align=PP_ALIGN.CENTER)
                    for bi, pt in enumerate((pts or [obj])[:3]):
                        ty = 5.22 + bi*0.42
                        oval(sl, 0.52, ty+0.10, 0.11, 0.11, C_ACC)
                        txt(sl, pt, 0.78, ty, 12.0, 0.38, 11.5, C_TEXT)
                elif layout == 'roadmap':
                    rect(sl, 0, 0, 13.33, 1.28, _bg2(0.88))
                    rect(sl, 0, 1.28, 13.33, 0.05, C_ACC)
                    rect(sl, 0, 0, 0.08, 1.28, C_ACC)
                    txt(sl, kicker, 0.52, 0.16, 6.0, 0.26, 8, C_ACC, bold=True)
                    txt(sl, title,  0.52, 0.44, 9.0, 0.74,
                        24, C_TEXT, bold=True, fn=FH)
                    if sub:
                        txt(sl, sub, 0.52, 1.04, 8.0,
                            0.30, 10.5, C_DIM, italic=True)
                    rect(sl, 0.55, 3.05, 12.22, 0.04, _acc(0.32))
                    steps = pts[:4]
                    n = len(steps) or 1
                    sw = 12.22/n
                    for bi, st in enumerate(steps):
                        cx = 0.55 + bi*sw + sw/2
                        oval(sl, cx-0.34, 2.86, 0.68, 0.68, _acc(0.38))
                        oval(sl, cx-0.21, 3.00, 0.42, 0.42, C_ACC)
                        txt(sl, f'{bi+1}', cx-0.21, 3.00, 0.42, 0.42, 10,
                            C_BG if not is_light else C_WHITE, bold=True, align=PP_ALIGN.CENTER)
                        bx = 0.55+bi*sw+0.10
                        cw = sw-0.20
                        rect(sl, bx, 3.82, cw, 2.62, _sur(0.70),
                             border=_acc(0.20), bw=0.5, rnd=True)
                        txt(sl, st, bx+0.14, 3.95, cw-0.28, 2.32, 11.5, C_TEXT)
                elif layout == 'comparison':
                    rect(sl, 0, 0, 13.33, 0.90, _bg2(0.88))
                    rect(sl, 0, 0.90, 13.33, 0.04, C_ACC)
                    rect(sl, 0, 0, 0.08, 0.90, C_ACC)
                    rect(sl, 6.46, 0.94, 0.05, 6.28, _acc(0.38))
                    txt(sl, kicker, 0.48, 0.10, 5.0, 0.24, 8, C_ACC, bold=True)
                    txt(sl, title,  0.48, 0.34, 12.3,
                        0.48, 21, C_TEXT, bold=True, fn=FH)
                    rect(sl, 0.48, 1.06, 5.76, 0.42, _acc(0.20), rnd=True)
                    txt(sl, 'PROBLEM / BEFORE', 0.66, 1.12,
                        5.40, 0.30, 8.5, C_ACC, bold=True)
                    rect(sl, 6.70, 1.06, 6.10, 0.42, _acc(0.38), rnd=True)
                    txt(sl, 'SOLUTION / AFTER',  6.88, 1.12,
                        5.70, 0.30, 8.5, C_ACC, bold=True)
                    lp = pts[:3]
                    rp = pts[3:] or pts[:3]
                    for bi, p in enumerate(lp):
                        ty = 1.74 + bi*0.96
                        oval(sl, 0.52, ty+0.13, 0.13, 0.13, _acc(0.50))
                        txt(sl, p, 0.80, ty, 5.30, 0.84, 11.5, C_TEXT)
                    for bi, p in enumerate(rp):
                        ty = 1.74 + bi*0.96
                        oval(sl, 6.76, ty+0.13, 0.13, 0.13, C_ACC)
                        txt(sl, p, 7.04, ty, 5.82, 0.84, 11.5, C_TEXT)
                    if sub:
                        txt(sl, sub, 0.48, 6.60, 12.3,
                            0.38, 10, C_DIM, italic=True)
                elif layout == 'closing':
                    rect(sl, 0, 0, 13.33, 0.10, C_ACC)
                    oval(sl, -2.0, -1.5, 8.5, 8.5, _acc(0.12))
                    oval(sl, 7.5, 1.0, 7.0, 7.0, _acc(0.10))
                    txt(sl, kicker, 0, 1.10, 13.33, 0.30, 9,
                        C_ACC, bold=True, align=PP_ALIGN.CENTER)
                    txt(sl, title, 0.80, 1.52, 11.73, 1.90, 42, C_TEXT,
                        bold=True, fn=FH, align=PP_ALIGN.CENTER)
                    rect(sl, 4.66, 3.50, 4.0, 0.05, _acc(0.58))
                    if sub:
                        txt(sl, sub, 1.0, 3.68, 11.33, 0.68, 13.5,
                            C_DIM, italic=True, align=PP_ALIGN.CENTER)
                    if stats:
                        n = min(len(stats), 3)
                        cw = 3.5
                        sx = (13.33-(n*cw+(n-1)*0.28))/2
                        for bi, st in enumerate(stats[:3]):
                            bx = sx+bi*(cw+0.28)
                            rect(sl, bx, 4.56, cw, 1.60, _sur(0.75),
                                 border=_acc(0.30), bw=0.6, rnd=True)
                            rect(sl, bx, 4.56, cw, 0.10, C_ACC)
                            txt(sl, st.get('value', ''), bx+0.14, 4.72, cw-0.28, 0.80,
                                30, C_TEXT, bold=True, fn=FH, align=PP_ALIGN.CENTER)
                            txt(sl, st.get('label', '').upper(), bx+0.14, 5.58, cw -
                                0.28, 0.30, 8, C_DIM, bold=True, align=PP_ALIGN.CENTER)
                    rect(sl, 0, 7.40, 13.33, 0.10, C_ACC)
                else:  # spotlight
                    rect(sl, 7.88, 0, 5.45, 7.5, _bg2(0.85))
                    oval(sl, 8.5, -1.0, 5.5, 5.5, _acc(0.18))
                    oval(sl, 9.5,  4.0, 3.8, 3.8, _acc(0.10))
                    rect(sl, 0, 0, 0.08, 7.5, C_ACC)
                    txt(sl, kicker, 0.38, 0.40, 5.5, 0.26, 8, C_ACC, bold=True)
                    txt(sl, title,  0.38, 0.80, 7.2, 1.80,
                        28, C_TEXT, bold=True, fn=FH)
                    if sub:
                        txt(sl, sub, 0.38, 2.54, 7.0,
                            0.58, 12, C_DIM, italic=True)
                    rect(sl, 0.38, 3.20, 5.0, 0.04, _acc(0.42))
                    if pts:
                        for bi, pt in enumerate(pts):
                            ty = 3.40 + bi*0.57
                            rect(sl, 0.38, ty+0.07, 0.22,
                                 0.22, _acc(0.30), rnd=True)
                            txt(sl, f'{bi+1}', 0.38, ty+0.07, 0.22, 0.22,
                                7, C_ACC, bold=True, align=PP_ALIGN.CENTER)
                            txt(sl, pt, 0.72, ty, 6.90, 0.52, 11.5, C_TEXT)
                    elif obj:
                        txt(sl, obj, 0.38, 3.40, 7.0,
                            0.80, 12, C_DIM, italic=True)
                    if stats:
                        for bi, st in enumerate(stats[:3]):
                            ty = 0.68 + bi*2.08
                            rect(sl, 8.08, ty, 4.92, 1.82, _sur(0.70),
                                 border=_acc(0.22), bw=0.5, rnd=True)
                            rect(sl, 8.08, ty, 4.92, 0.12, C_ACC)
                            txt(sl, st.get('value', ''), 8.26, ty+0.17, 4.56, 1.0,
                                34, C_TEXT, bold=True, fn=FH, align=PP_ALIGN.CENTER)
                            txt(sl, st.get('label', '').upper(), 8.26, ty+1.28,
                                4.56, 0.30, 8, C_DIM, bold=True, align=PP_ALIGN.CENTER)
                    else:
                        rect(sl, 8.08, 0.68, 4.92, 5.82, _sur(0.70),
                             border=_acc(0.18), bw=0.5, rnd=True)
                        txt(sl, 'VISUAL DIRECTION', 8.28, 0.86,
                            4.52, 0.26, 7.5, C_ACC, bold=True)
                        rect(sl, 8.08, 1.10, 4.92, 0.04, _acc(0.32))
                        txt(sl, vis or 'Use one dominant visual.',
                            8.28, 1.20, 4.52, 4.92, 11.5, C_TEXT)

            # ══════════════════════════════════════════════════════════════════
            # STYLE B — KINETIC  (diagonal slash, impact typography, bold energy)
            # DNA: angled accent band top-right, full-bleed color blocks, huge type
            # ══════════════════════════════════════════════════════════════════
            elif STYLE == 'kinetic':
                # Shared kinetic BG motif — thick accent slash top-right
                rect(sl, 9.0, 0, 4.33, 0.06, C_ACC)
                rect(sl, 0, 7.44, 13.33, 0.06, C_ACC)

                if layout == 'cover':
                    rect(sl, 0, 0, 13.33, 3.8, _bg2(0.92))
                    oval(sl, 8.0, -1.0, 7.0, 7.0, _acc(0.22))
                    txt(sl, kicker, 0.65, 0.42, 6.0, 0.28, 8, C_ACC, bold=True)
                    txt(sl, title, 0.65, 0.85, 10.0, 2.60,
                        48, C_TEXT, bold=True, fn=FH)
                    if sub:
                        txt(sl, sub, 0.65, 3.60, 8.0,
                            0.70, 14, C_DIM, italic=True)
                    rect(sl, 0.65, 4.42, 4.0, 0.06, C_ACC)
                    txt(sl, idea.title()[:40], 0.65,
                        4.58, 6.0, 0.42, 11, C_DIM)
                    if stats:
                        for bi, st in enumerate(stats[:2]):
                            bx = 7.0 + bi*3.1
                            txt(sl, st.get('value', ''), bx, 1.0, 3.0,
                                1.60, 54, C_ACC, bold=True, fn=FH)
                            txt(sl, st.get('label', '').upper(), bx,
                                2.70, 3.0, 0.36, 9, C_DIM, bold=True)
                elif layout == 'metrics':
                    rect(sl, 0, 0, 13.33, 0.80, C_ACC)
                    txt(sl, kicker, 0.55, 0.10, 6.0, 0.26, 8,
                        C_BG if not is_light else C_WHITE, bold=True)
                    txt(sl, title,  0.55, 0.42, 10.0,
                        0.84, 28, C_TEXT, bold=True, fn=FH)
                    if stats:
                        n = len(stats)
                        cw = (12.33 - 0.2*(n-1)) / n
                        for bi, st in enumerate(stats):
                            bx = 0.50 + bi*(cw+0.20)
                            rect(sl, bx, 1.60, cw, 4.0, _sur(0.72),
                                 border=C_ACC, bw=1.2, rnd=False)
                            txt(sl, st.get('value', ''), bx+0.20, 1.90, cw-0.40, 2.0,
                                56, C_ACC, bold=True, fn=FH, align=PP_ALIGN.CENTER)
                            txt(sl, st.get('label', '').upper(), bx+0.20, 4.10, cw-0.40,
                                0.38, 9.5, C_TEXT, bold=True, align=PP_ALIGN.CENTER)
                    for bi, pt in enumerate((pts or [obj])[:2]):
                        ty = 5.85 + bi*0.46
                        rect(sl, 0.50, ty+0.06, 0.30, 0.22, C_ACC)
                        txt(sl, pt, 0.94, ty, 11.8, 0.40, 11.5, C_TEXT)
                elif layout == 'roadmap':
                    rect(sl, 0, 0, 13.33, 1.0, _bg2(0.90))
                    txt(sl, kicker, 0.55, 0.12, 6.0, 0.24, 8, C_ACC, bold=True)
                    txt(sl, title,  0.55, 0.38, 10.0,
                        0.56, 22, C_TEXT, bold=True, fn=FH)
                    # Vertical timeline (left edge)
                    rect(sl, 1.20, 1.20, 0.05, 5.80, _acc(0.40))
                    steps = pts[:4]
                    for bi, st in enumerate(steps):
                        ty = 1.40 + bi*1.48
                        oval(sl, 0.90, ty, 0.60, 0.60, C_ACC)
                        txt(sl, f'{bi+1}', 0.90, ty, 0.60, 0.60, 11,
                            C_BG if not is_light else C_WHITE, bold=True, align=PP_ALIGN.CENTER)
                        rect(sl, 1.80, ty, 11.0, 1.24, _sur(0.70),
                             border=_acc(0.22), bw=0.5, rnd=False)
                        txt(sl, st, 2.00, ty+0.10, 10.60, 1.0, 12, C_TEXT)
                elif layout == 'comparison':
                    rect(sl, 0, 0, 6.46, 7.5, _bg2(0.90))
                    rect(sl, 6.46, 0, 6.87, 7.5, _acc(0.20))
                    txt(sl, kicker, 0.50, 0.22, 5.5, 0.26, 8, C_ACC, bold=True)
                    txt(sl, title,  0.50, 0.55, 12.0,
                        0.80, 24, C_TEXT, bold=True, fn=FH)
                    txt(sl, 'BEFORE', 1.0, 1.55, 4.5,
                        0.38, 11, C_DIM, bold=True)
                    txt(sl, 'AFTER',  7.0, 1.55, 5.5,
                        0.38, 11, C_ACC, bold=True)
                    lp = pts[:3]
                    rp = pts[3:] or pts[:3]
                    for bi, p in enumerate(lp):
                        ty = 2.10 + bi*1.12
                        rect(sl, 0.50, ty, 5.60, 0.88, _sur(0.65), rnd=False)
                        txt(sl, p, 0.70, ty+0.12, 5.20, 0.72, 11.5, C_TEXT)
                    for bi, p in enumerate(rp):
                        ty = 2.10 + bi*1.12
                        rect(sl, 6.66, ty, 6.30, 0.88, _acc(0.28), rnd=False)
                        txt(sl, p, 6.86, ty+0.12, 5.90, 0.72, 11.5, C_TEXT)
                elif layout == 'closing':
                    rect(sl, 0, 0, 13.33, 7.5, _bg2(0.95))
                    rect(sl, 0, 0, 13.33, 0.12, C_ACC)
                    rect(sl, 0, 7.38, 13.33, 0.12, C_ACC)
                    oval(sl, -1.5, -1.5, 7.5, 7.5, _acc(0.14))
                    oval(sl, 9.0,  3.0, 6.0, 6.0, _acc(0.12))
                    txt(sl, title, 0.80, 0.90, 11.73, 2.20, 50, C_TEXT,
                        bold=True, fn=FH, align=PP_ALIGN.CENTER)
                    rect(sl, 4.66, 3.20, 4.0, 0.06, C_ACC)
                    if sub:
                        txt(sl, sub, 1.0, 3.38, 11.33, 0.70, 14,
                            C_DIM, italic=True, align=PP_ALIGN.CENTER)
                    if stats:
                        n = min(len(stats), 3)
                        cw = 3.5
                        sx = (13.33-(n*cw+(n-1)*0.28))/2
                        for bi, st in enumerate(stats[:3]):
                            bx = sx+bi*(cw+0.28)
                            rect(sl, bx, 4.40, cw, 1.80, _acc(0.22),
                                 border=C_ACC, bw=1.0, rnd=False)
                            txt(sl, st.get('value', ''), bx+0.14, 4.55, cw-0.28, 0.90,
                                32, C_ACC, bold=True, fn=FH, align=PP_ALIGN.CENTER)
                            txt(sl, st.get('label', '').upper(), bx+0.14, 5.58, cw -
                                0.28, 0.34, 8.5, C_TEXT, bold=True, align=PP_ALIGN.CENTER)
                else:  # spotlight
                    rect(sl, 0, 0, 0.18, 7.5, C_ACC)
                    rect(sl, 0.18, 0, 13.15, 0.72, _bg2(0.90))
                    txt(sl, kicker, 0.42, 0.12, 5.5, 0.26, 8, C_ACC, bold=True)
                    txt(sl, title,  0.42, 0.80, 7.8, 2.0,
                        32, C_TEXT, bold=True, fn=FH)
                    if sub:
                        txt(sl, sub, 0.42, 2.76, 7.8,
                            0.58, 12, C_DIM, italic=True)
                    rect(sl, 0.42, 3.48, 0.60, 0.06, C_ACC)
                    if pts:
                        for bi, pt in enumerate(pts):
                            ty = 3.70 + bi*0.62
                            rect(sl, 0.42, ty+0.10, 0.28,
                                 0.28, C_ACC, rnd=False)
                            txt(sl, f'{bi+1}', 0.42, ty+0.10, 0.28, 0.28, 8,
                                C_BG if not is_light else C_WHITE, bold=True, align=PP_ALIGN.CENTER)
                            txt(sl, pt, 0.84, ty, 7.20, 0.56, 12, C_TEXT)
                    elif obj:
                        txt(sl, obj, 0.42, 3.70, 7.8,
                            0.80, 12, C_DIM, italic=True)
                    if stats:
                        for bi, st in enumerate(stats[:3]):
                            ty = 0.20 + bi*2.32
                            rect(sl, 8.70, ty, 4.30, 2.10, _acc(0.22),
                                 border=C_ACC, bw=0.8, rnd=False)
                            txt(sl, st.get('value', ''), 8.88, ty+0.22, 3.94, 1.20,
                                38, C_ACC, bold=True, fn=FH, align=PP_ALIGN.CENTER)
                            txt(sl, st.get('label', '').upper(), 8.88, ty+1.62, 3.94,
                                0.32, 8.5, C_TEXT, bold=True, align=PP_ALIGN.CENTER)
                    else:
                        rect(sl, 8.70, 0.20, 4.30, 6.82, _sur(0.70),
                             border=_acc(0.22), bw=0.6, rnd=False)
                        txt(sl, 'VISUAL DIRECTION', 8.90, 0.38,
                            3.90, 0.26, 7.5, C_ACC, bold=True)
                        rect(sl, 8.70, 0.62, 4.30, 0.05, C_ACC)
                        txt(sl, vis or 'Lead with a bold proof visual.',
                            8.90, 0.72, 3.90, 6.10, 12, C_TEXT)

            # ══════════════════════════════════════════════════════════════════
            # STYLE C — ATLAS  (left sidebar strip + data-centric layout)
            # DNA: narrow left accent sidebar, top header bar, right panel for data
            # ══════════════════════════════════════════════════════════════════
            elif STYLE == 'atlas':
                # Persistent left sidebar
                rect(sl, 0, 0, 0.90, 7.5, _bg2(0.92))
                rect(sl, 0.90, 0, 0.04, 7.5, _acc(0.30))
                # Sidebar kicker text (rotated effect via tall narrow box)
                txt(sl, str(i+1).zfill(2), 0.08, 0.30, 0.74, 0.50, 18,
                    C_ACC, bold=True, fn=FH, align=PP_ALIGN.CENTER)
                txt(sl, kicker[:8], 0.05, 1.0, 0.80, 3.0, 8,
                    C_DIM, bold=True, align=PP_ALIGN.CENTER)

                if layout == 'cover':
                    rect(sl, 0.94, 0, 12.39, 7.5, _bg2(0.40))
                    oval(sl, 7.0, -1.5, 7.5, 7.5, _acc(0.18))
                    oval(sl, 9.5,  3.5, 5.0, 5.0, _acc(0.12))
                    txt(sl, title, 1.10, 0.80, 9.0, 3.20,
                        44, C_TEXT, bold=True, fn=FH)
                    if sub:
                        txt(sl, sub, 1.10, 4.10, 8.0,
                            0.80, 13, C_DIM, italic=True)
                    rect(sl, 1.10, 5.00, 4.0, 0.04, C_ACC)
                    txt(sl, idea.title()[:38], 1.10,
                        5.18, 6.0, 0.38, 10.5, C_DIM)
                    if stats:
                        txt(sl, stats[0].get('value', ''), 7.5, 1.60, 5.5, 2.20, 64, _acc(
                            0.88), bold=True, fn=FH, align=PP_ALIGN.CENTER)
                        txt(sl, stats[0].get('label', '').upper(
                        ), 7.5, 3.92, 5.5, 0.36, 10, C_DIM, bold=True, align=PP_ALIGN.CENTER)
                elif layout == 'metrics':
                    rect(sl, 0.94, 0, 12.39, 1.10, _bg2(0.88))
                    rect(sl, 0.94, 1.10, 12.39, 0.05, C_ACC)
                    txt(sl, title, 1.10, 0.20, 10.0, 0.80,
                        26, C_TEXT, bold=True, fn=FH)
                    if stats:
                        n = len(stats)
                        cw = (11.83 - 0.18*(n-1)) / n
                        for bi, st in enumerate(stats):
                            bx = 1.10 + bi*(cw+0.18)
                            rect(sl, bx, 1.32, cw, 3.60, _sur(0.75),
                                 border=_acc(0.28), bw=0.7, rnd=True)
                            rect(sl, bx, 1.32, cw, 0.14, C_ACC)
                            txt(sl, st.get('value', ''), bx+0.18, 1.58, cw-0.36, 1.80,
                                52, C_TEXT, bold=True, fn=FH, align=PP_ALIGN.CENTER)
                            txt(sl, st.get('label', '').upper(), bx+0.18, 3.50, cw-0.36,
                                0.34, 9, C_DIM, bold=True, align=PP_ALIGN.CENTER)
                    for bi, pt in enumerate((pts or [obj])[:3]):
                        ty = 5.22 + bi*0.40
                        rect(sl, 1.10, ty+0.08, 0.20, 0.20, C_ACC, rnd=True)
                        txt(sl, pt, 1.44, ty, 11.5, 0.38, 11, C_TEXT)
                elif layout == 'roadmap':
                    rect(sl, 0.94, 0, 12.39, 1.0, _bg2(0.88))
                    rect(sl, 0.94, 1.0, 12.39, 0.04, C_ACC)
                    txt(sl, title, 1.10, 0.14, 10.0, 0.76,
                        24, C_TEXT, bold=True, fn=FH)
                    if sub:
                        txt(sl, sub, 1.10, 0.78, 8.0,
                            0.28, 10, C_DIM, italic=True)
                    rect(sl, 1.10, 3.04, 12.0, 0.04, _acc(0.32))
                    steps = pts[:4]
                    n = len(steps) or 1
                    sw = 12.0/n
                    for bi, st in enumerate(steps):
                        cx = 1.10 + bi*sw + sw/2
                        oval(sl, cx-0.32, 2.84, 0.64, 0.64, _acc(0.40))
                        oval(sl, cx-0.20, 2.98, 0.40, 0.40, C_ACC)
                        txt(sl, f'{bi+1}', cx-0.20, 2.98, 0.40, 0.40, 9,
                            C_BG if not is_light else C_WHITE, bold=True, align=PP_ALIGN.CENTER)
                        bx = 1.10+bi*sw+0.08
                        cw = sw-0.16
                        rect(sl, bx, 3.80, cw, 2.70, _sur(0.72),
                             border=_acc(0.20), bw=0.5, rnd=True)
                        txt(sl, st, bx+0.12, 3.92, cw-0.24, 2.42, 11, C_TEXT)
                elif layout == 'comparison':
                    rect(sl, 0.94, 0, 12.39, 0.80, _bg2(0.88))
                    rect(sl, 0.94, 0.80, 12.39, 0.04, C_ACC)
                    rect(sl, 7.16, 0.84, 0.05, 6.38, _acc(0.36))
                    txt(sl, title, 1.10, 0.10, 11.5, 0.64,
                        22, C_TEXT, bold=True, fn=FH)
                    rect(sl, 1.10, 0.98, 5.78, 0.40, _acc(0.20), rnd=True)
                    txt(sl, 'PROBLEM / BEFORE', 1.28, 1.04,
                        5.40, 0.28, 8.5, C_ACC, bold=True)
                    rect(sl, 7.40, 0.98, 5.78, 0.40, _acc(0.36), rnd=True)
                    txt(sl, 'SOLUTION / AFTER', 7.58, 1.04,
                        5.40, 0.28, 8.5, C_ACC, bold=True)
                    lp = pts[:3]
                    rp = pts[3:] or pts[:3]
                    for bi, p in enumerate(lp):
                        ty = 1.58 + bi*1.0
                        oval(sl, 1.14, ty+0.12, 0.12, 0.12, _acc(0.50))
                        txt(sl, p, 1.40, ty, 5.38, 0.88, 11.5, C_TEXT)
                    for bi, p in enumerate(rp):
                        ty = 1.58 + bi*1.0
                        oval(sl, 7.44, ty+0.12, 0.12, 0.12, C_ACC)
                        txt(sl, p, 7.70, ty, 5.58, 0.88, 11.5, C_TEXT)
                elif layout == 'closing':
                    rect(sl, 0.94, 0, 12.39, 7.5, _bg2(0.50))
                    oval(sl, 5.0, -1.0, 9.0, 9.0, _acc(0.12))
                    txt(sl, title, 1.10, 1.20, 11.20, 2.10, 44, C_TEXT,
                        bold=True, fn=FH, align=PP_ALIGN.CENTER)
                    rect(sl, 4.66, 3.42, 4.0, 0.05, C_ACC)
                    if sub:
                        txt(sl, sub, 1.10, 3.60, 11.2, 0.70, 13.5,
                            C_DIM, italic=True, align=PP_ALIGN.CENTER)
                    if stats:
                        n = min(len(stats), 3)
                        cw = 3.4
                        sx = (11.23-(n*cw+(n-1)*0.26))/2 + 1.0
                        for bi, st in enumerate(stats[:3]):
                            bx = sx+bi*(cw+0.26)
                            rect(sl, bx, 4.50, cw, 1.72, _sur(0.75),
                                 border=_acc(0.30), bw=0.7, rnd=True)
                            rect(sl, bx, 4.50, cw, 0.12, C_ACC)
                            txt(sl, st.get('value', ''), bx+0.14, 4.65, cw-0.28, 0.82,
                                30, C_TEXT, bold=True, fn=FH, align=PP_ALIGN.CENTER)
                            txt(sl, st.get('label', '').upper(), bx+0.14, 5.56, cw -
                                0.28, 0.30, 8, C_DIM, bold=True, align=PP_ALIGN.CENTER)
                    rect(sl, 0.94, 7.38, 12.39, 0.10, C_ACC)
                else:  # spotlight
                    rect(sl, 0.94, 0, 7.0, 7.5, _bg2(0.50))
                    rect(sl, 7.94, 0, 5.39, 7.5, _sur(0.80))
                    oval(sl, 8.0, -0.8, 5.2, 5.2, _acc(0.16))
                    txt(sl, title, 1.10, 0.50, 6.60, 2.20,
                        30, C_TEXT, bold=True, fn=FH)
                    if sub:
                        txt(sl, sub, 1.10, 2.68, 6.60,
                            0.56, 12, C_DIM, italic=True)
                    rect(sl, 1.10, 3.34, 4.50, 0.04, C_ACC)
                    if pts:
                        for bi, pt in enumerate(pts):
                            ty = 3.52 + bi*0.60
                            rect(sl, 1.10, ty+0.08, 0.22,
                                 0.22, _acc(0.30), rnd=True)
                            txt(sl, f'{bi+1}', 1.10, ty+0.08, 0.22, 0.22,
                                7, C_ACC, bold=True, align=PP_ALIGN.CENTER)
                            txt(sl, pt, 1.44, ty, 6.30, 0.54, 11.5, C_TEXT)
                    elif obj:
                        txt(sl, obj, 1.10, 3.52, 6.60,
                            0.80, 12, C_DIM, italic=True)
                    if stats:
                        for bi, st in enumerate(stats[:3]):
                            ty = 0.40 + bi*2.26
                            rect(sl, 8.12, ty, 5.0, 2.06, _sur(0.70),
                                 border=_acc(0.22), bw=0.5, rnd=True)
                            rect(sl, 8.12, ty, 5.0, 0.12, C_ACC)
                            txt(sl, st.get('value', ''), 8.30, ty+0.18, 4.64, 1.12,
                                36, C_TEXT, bold=True, fn=FH, align=PP_ALIGN.CENTER)
                            txt(sl, st.get('label', '').upper(), 8.30, ty+1.40,
                                4.64, 0.30, 8, C_DIM, bold=True, align=PP_ALIGN.CENTER)
                    else:
                        rect(sl, 8.12, 0.40, 5.0, 6.32, _sur(0.70),
                             border=_acc(0.18), bw=0.5, rnd=True)
                        txt(sl, 'VISUAL DIRECTION', 8.32, 0.58,
                            4.60, 0.26, 7.5, C_ACC, bold=True)
                        rect(sl, 8.12, 0.82, 5.0, 0.04, _acc(0.32))
                        txt(sl, vis or 'Use one proof visual.',
                            8.32, 0.92, 4.60, 5.50, 12, C_TEXT)

            # ══════════════════════════════════════════════════════════════════
            # STYLE D — LUXE  (centred luxury, large serif, gold rule system)
            # DNA: centred everything, thin horizontal rules, editorial white space
            # ══════════════════════════════════════════════════════════════════
            elif STYLE == 'luxe':
                # Top + bottom thin gold rules
                rect(sl, 0, 0.25, 13.33, 0.016, C_ACC)
                rect(sl, 0, 7.22, 13.33, 0.016, C_ACC)

                if layout == 'cover':
                    oval(sl, -2.0, -2.0, 8.0, 8.0, _acc(0.10))
                    oval(sl, 9.0,  3.0, 7.0, 7.0, _acc(0.08))
                    # Centred rule above headline
                    rect(sl, 5.66, 0.80, 2.0, 0.016, C_ACC)
                    txt(sl, kicker, 0, 0.90, 13.33, 0.28, 9,
                        C_ACC, bold=True, align=PP_ALIGN.CENTER)
                    txt(sl, title, 0.80, 1.28, 11.73, 3.20, 52, C_TEXT,
                        bold=True, fn=FH, align=PP_ALIGN.CENTER)
                    rect(sl, 5.66, 4.58, 2.0, 0.016, C_ACC)
                    if sub:
                        txt(sl, sub, 1.0, 4.74, 11.33, 0.72, 14,
                            C_DIM, italic=True, align=PP_ALIGN.CENTER)
                    txt(sl, idea.title()[:48], 0, 5.60, 13.33,
                        0.38, 10.5, C_DIM, align=PP_ALIGN.CENTER)
                    if stats:
                        txt(sl, stats[0].get('value', ''), 0, 3.00, 13.33, 1.40, 60, _acc(
                            0.80), bold=True, fn=FH, align=PP_ALIGN.CENTER)
                elif layout == 'metrics':
                    rect(sl, 0, 0.92, 13.33, 0.016, C_ACC)
                    txt(sl, kicker, 0, 0.46, 13.33, 0.28, 9,
                        C_ACC, bold=True, align=PP_ALIGN.CENTER)
                    txt(sl, title, 0.80, 1.10, 11.73, 1.40, 30, C_TEXT,
                        bold=True, fn=FH, align=PP_ALIGN.CENTER)
                    if stats:
                        n = len(stats)
                        cw = (11.73 - 0.20*(n-1)) / n
                        sx = (13.33-(n*cw+(n-1)*0.20))/2
                        for bi, st in enumerate(stats):
                            bx = sx + bi*(cw+0.20)
                            rect(sl, bx, 2.72, cw, 3.20, _sur(0.72),
                                 border=_acc(0.22), bw=0.5, rnd=True)
                            rect(sl, bx, 2.72, cw, 0.016, C_ACC)
                            txt(sl, st.get('value', ''), bx+0.18, 2.90, cw-0.36, 1.70,
                                52, C_TEXT, bold=True, fn=FH, align=PP_ALIGN.CENTER)
                            txt(sl, st.get('label', '').upper(), bx+0.18, 4.68, cw-0.36,
                                0.34, 9, C_DIM, bold=True, align=PP_ALIGN.CENTER)
                    for bi, pt in enumerate((pts or [obj])[:2]):
                        txt(sl, pt, 1.0, 6.10+bi*0.42, 11.33, 0.38,
                            11.5, C_TEXT, align=PP_ALIGN.CENTER)
                elif layout == 'roadmap':
                    rect(sl, 0, 0.92, 13.33, 0.016, C_ACC)
                    txt(sl, kicker, 0, 0.46, 13.33, 0.28, 9,
                        C_ACC, bold=True, align=PP_ALIGN.CENTER)
                    txt(sl, title, 0.80, 1.10, 11.73, 1.10, 28, C_TEXT,
                        bold=True, fn=FH, align=PP_ALIGN.CENTER)
                    if sub:
                        txt(sl, sub, 1.0, 2.10, 11.33, 0.38, 11,
                            C_DIM, italic=True, align=PP_ALIGN.CENTER)
                    rect(sl, 0.55, 3.28, 12.22, 0.016, _acc(0.40))
                    steps = pts[:4]
                    n = len(steps) or 1
                    sw = 12.22/n
                    for bi, st in enumerate(steps):
                        cx = 0.55 + bi*sw + sw/2
                        oval(sl, cx-0.28, 3.12, 0.56, 0.56, C_ACC)
                        txt(sl, f'{bi+1}', cx-0.28, 3.12, 0.56, 0.56, 9,
                            C_BG if not is_light else C_WHITE, bold=True, align=PP_ALIGN.CENTER)
                        bx = 0.55+bi*sw+0.10
                        cw = sw-0.20
                        rect(sl, bx, 3.96, cw, 2.56, _sur(0.68),
                             border=_acc(0.18), bw=0.4, rnd=True)
                        txt(sl, st, bx+0.14, 4.10, cw-0.28, 2.28,
                            11.5, C_TEXT, align=PP_ALIGN.CENTER)
                elif layout == 'comparison':
                    rect(sl, 0, 0.92, 13.33, 0.016, C_ACC)
                    rect(sl, 6.66, 1.0, 0.016, 6.22, _acc(0.40))
                    txt(sl, kicker, 0, 0.46, 13.33, 0.28, 9,
                        C_ACC, bold=True, align=PP_ALIGN.CENTER)
                    txt(sl, title, 0.80, 1.10, 11.73, 0.88, 26, C_TEXT,
                        bold=True, fn=FH, align=PP_ALIGN.CENTER)
                    txt(sl, 'BEFORE', 0.55, 2.06, 5.88, 0.32, 10,
                        C_DIM, bold=True, align=PP_ALIGN.CENTER)
                    txt(sl, 'AFTER',  6.88, 2.06, 6.10, 0.32, 10,
                        C_ACC, bold=True, align=PP_ALIGN.CENTER)
                    lp = pts[:3]
                    rp = pts[3:] or pts[:3]
                    for bi, p in enumerate(lp):
                        ty = 2.58 + bi*1.10
                        txt(sl, p, 0.55, ty, 5.88, 0.96, 12,
                            C_TEXT, align=PP_ALIGN.CENTER)
                        rect(sl, 2.0, ty+1.04, 2.96, 0.014, _acc(0.20))
                    for bi, p in enumerate(rp):
                        ty = 2.58 + bi*1.10
                        txt(sl, p, 6.88, ty, 6.10, 0.96, 12,
                            C_TEXT, align=PP_ALIGN.CENTER)
                        rect(sl, 8.0, ty+1.04, 3.10, 0.014, _acc(0.20))
                elif layout == 'closing':
                    oval(sl, -1.5, -1.5, 7.0, 7.0, _acc(0.11))
                    oval(sl, 9.0,  2.5, 6.0, 6.0, _acc(0.09))
                    rect(sl, 5.66, 1.10, 2.0, 0.016, C_ACC)
                    txt(sl, kicker, 0, 1.24, 13.33, 0.30, 9,
                        C_ACC, bold=True, align=PP_ALIGN.CENTER)
                    txt(sl, title, 0.80, 1.60, 11.73, 2.40, 50, C_TEXT,
                        bold=True, fn=FH, align=PP_ALIGN.CENTER)
                    rect(sl, 5.66, 4.10, 2.0, 0.016, C_ACC)
                    if sub:
                        txt(sl, sub, 1.0, 4.26, 11.33, 0.68, 14,
                            C_DIM, italic=True, align=PP_ALIGN.CENTER)
                    if stats:
                        n = min(len(stats), 3)
                        cw = 3.4
                        sx = (13.33-(n*cw+(n-1)*0.26))/2
                        for bi, st in enumerate(stats[:3]):
                            bx = sx+bi*(cw+0.26)
                            rect(sl, bx, 5.10, cw, 1.80, _sur(0.72),
                                 border=_acc(0.24), bw=0.5, rnd=True)
                            rect(sl, bx+cw*0.3, 5.10, cw*0.4, 0.016, C_ACC)
                            txt(sl, st.get('value', ''), bx+0.14, 5.24, cw-0.28, 0.88,
                                30, C_TEXT, bold=True, fn=FH, align=PP_ALIGN.CENTER)
                            txt(sl, st.get('label', '').upper(), bx+0.14, 6.18, cw -
                                0.28, 0.30, 8, C_DIM, bold=True, align=PP_ALIGN.CENTER)
                else:  # spotlight
                    rect(sl, 0, 0.92, 13.33, 0.016, C_ACC)
                    txt(sl, kicker, 0, 0.46, 13.33, 0.28, 9,
                        C_ACC, bold=True, align=PP_ALIGN.CENTER)
                    txt(sl, title, 0.80, 1.10, 11.73, 2.0, 34, C_TEXT,
                        bold=True, fn=FH, align=PP_ALIGN.CENTER)
                    if sub:
                        txt(sl, sub, 1.0, 3.00, 11.33, 0.58, 12.5,
                            C_DIM, italic=True, align=PP_ALIGN.CENTER)
                    rect(sl, 5.66, 3.68, 2.0, 0.016, C_ACC)
                    if pts:
                        for bi, pt in enumerate(pts[:4]):
                            ty = 3.90 + bi*0.68
                            txt(sl, f'— {pt}', 1.50, ty, 10.33, 0.60,
                                12, C_TEXT, align=PP_ALIGN.CENTER)
                    elif obj:
                        txt(sl, obj, 1.0, 3.90, 11.33, 0.80, 12,
                            C_DIM, italic=True, align=PP_ALIGN.CENTER)
                    if stats:
                        n = len(stats)
                        cw = (11.33-(0.20*(n-1)))/n
                        sx = (13.33-(n*cw+(n-1)*0.20))/2
                        for bi, st in enumerate(stats):
                            bx = sx + bi*(cw+0.20)
                            rect(sl, bx, 5.90, cw, 1.10, _sur(0.72),
                                 border=_acc(0.22), bw=0.4, rnd=True)
                            txt(sl, st.get('value', ''), bx+0.12, 5.98, cw-0.24, 0.64,
                                26, C_TEXT, bold=True, fn=FH, align=PP_ALIGN.CENTER)
                            txt(sl, st.get('label', '').upper(), bx+0.12, 6.68, cw -
                                0.24, 0.24, 7.5, C_DIM, bold=True, align=PP_ALIGN.CENTER)

            # ══════════════════════════════════════════════════════════════════
            # STYLE E — BOARDROOM  (light canvas, consulting aesthetic, thin rules)
            # DNA: top thin header rule, serif titles, structured grid, clean spacing
            # ══════════════════════════════════════════════════════════════════
            elif STYLE == 'boardroom':
                # Light bg already set; add structural chrome
                rect(sl, 0, 0, 13.33, 0.60, _bg2(0.70))
                rect(sl, 0, 0.60, 13.33, 0.024, C_ACC)

                if layout == 'cover':
                    rect(sl, 0, 0, 0.06, 7.5, C_ACC)
                    txt(sl, kicker, 0.50, 0.14, 6.0, 0.26, 8, C_ACC, bold=True)
                    txt(sl, title, 0.50, 0.82, 8.5, 3.50,
                        46, C_TEXT, bold=True, fn=FH)
                    if sub:
                        txt(sl, sub, 0.50, 4.40, 7.5,
                            0.80, 14, C_DIM, italic=True)
                    rect(sl, 0.50, 5.30, 3.50, 0.024, C_ACC)
                    txt(sl, idea.title()[:44], 0.50,
                        5.46, 6.0, 0.44, 11, C_DIM)
                    if stats:
                        for bi, st in enumerate(stats[:2]):
                            bx = 9.20 + bi*1.8
                            rect(sl, 9.0, 0.80, 4.0, 6.0, _bg2(
                                0.60), border=_acc(0.30), bw=0.5)
                            txt(sl, st.get('value', ''), 9.10, 1.60, 3.80, 2.0,
                                52, C_ACC, bold=True, fn=FH, align=PP_ALIGN.CENTER)
                            txt(sl, st.get('label', '').upper(), 9.10, 3.70, 3.80,
                                0.36, 9, C_DIM, bold=True, align=PP_ALIGN.CENTER)
                            break
                    else:
                        rect(sl, 9.0, 0.80, 4.0, 6.0, _bg2(
                            0.60), border=_acc(0.30), bw=0.5)
                        txt(sl, vis or 'Lead with precision.', 9.10, 2.60, 3.80,
                            2.0, 13, C_DIM, italic=True, align=PP_ALIGN.CENTER)
                elif layout == 'metrics':
                    txt(sl, kicker, 0.50, 0.14, 6.0, 0.26, 8, C_ACC, bold=True)
                    txt(sl, title,  0.50, 0.72, 10.0,
                        1.30, 30, C_TEXT, bold=True, fn=FH)
                    if stats:
                        n = len(stats)
                        cw = (12.33 - 0.16*(n-1)) / n
                        for bi, st in enumerate(stats):
                            bx = 0.50 + bi*(cw+0.16)
                            rect(sl, bx, 2.20, cw, 3.50, _bg2(
                                0.60), border=_acc(0.28), bw=0.5)
                            rect(sl, bx, 2.20, cw, 0.024, C_ACC)
                            txt(sl, st.get('value', ''), bx+0.18, 2.44, cw-0.36, 1.90,
                                52, C_ACC, bold=True, fn=FH, align=PP_ALIGN.CENTER)
                            txt(sl, st.get('label', '').upper(), bx+0.18, 4.42, cw-0.36,
                                0.36, 9, C_DIM, bold=True, align=PP_ALIGN.CENTER)
                    for bi, pt in enumerate((pts or [obj])[:3]):
                        ty = 6.00 + bi*0.38
                        rect(sl, 0.50, ty+0.10, 0.10, 0.10, C_ACC)
                        txt(sl, pt, 0.74, ty, 12.0, 0.36, 11, C_TEXT)
                elif layout == 'roadmap':
                    txt(sl, kicker, 0.50, 0.14, 6.0, 0.26, 8, C_ACC, bold=True)
                    txt(sl, title,  0.50, 0.72, 9.0, 0.88,
                        26, C_TEXT, bold=True, fn=FH)
                    if sub:
                        txt(sl, sub, 0.50, 1.48, 8.0,
                            0.34, 11, C_DIM, italic=True)
                    rect(sl, 0.50, 3.12, 12.33, 0.024, _acc(0.36))
                    steps = pts[:4]
                    n = len(steps) or 1
                    sw = 12.33/n
                    for bi, st in enumerate(steps):
                        cx = 0.50 + bi*sw + sw/2
                        rect(sl, cx-0.28, 2.96, 0.56, 0.32, C_ACC)
                        txt(sl, f'{bi+1:02d}', cx-0.28, 2.96, 0.56, 0.32, 10,
                            C_BG if not is_light else C_WHITE, bold=True, align=PP_ALIGN.CENTER)
                        bx = 0.50+bi*sw+0.08
                        cw = sw-0.16
                        rect(sl, bx, 3.30, cw, 3.42, _bg2(
                            0.65), border=_acc(0.22), bw=0.4)
                        txt(sl, st, bx+0.14, 3.44, cw-0.28, 3.12, 12, C_TEXT)
                elif layout == 'comparison':
                    rect(sl, 6.46, 0.64, 0.024, 6.58, _acc(0.36))
                    txt(sl, kicker, 0.50, 0.14, 5.0, 0.26, 8, C_ACC, bold=True)
                    txt(sl, title,  0.50, 0.72, 12.3,
                        0.88, 24, C_TEXT, bold=True, fn=FH)
                    rect(sl, 0.50, 1.72, 5.70, 0.36, _bg2(
                        0.70), border=_acc(0.28), bw=0.4)
                    txt(sl, 'STATUS QUO', 0.68, 1.78,
                        5.32, 0.24, 8.5, C_DIM, bold=True)
                    rect(sl, 6.70, 1.72, 6.30, 0.36, _bg2(
                        0.70), border=_acc(0.28), bw=0.4)
                    txt(sl, 'WITH THIS SOLUTION', 6.88, 1.78,
                        5.92, 0.24, 8.5, C_ACC, bold=True)
                    lp = pts[:4]
                    rp = pts[4:] or pts[:4]
                    for bi, p in enumerate(lp[:3]):
                        ty = 2.24 + bi*1.14
                        rect(sl, 0.50, ty, 5.70, 0.94, _bg2(
                            0.60), border=_acc(0.14), bw=0.4)
                        txt(sl, p, 0.66, ty+0.12, 5.38, 0.70, 11.5, C_TEXT)
                    for bi, p in enumerate(rp[:3]):
                        ty = 2.24 + bi*1.14
                        rect(sl, 6.70, ty, 6.30, 0.94, _bg2(
                            0.60), border=_acc(0.14), bw=0.4)
                        txt(sl, p, 6.86, ty+0.12, 5.98, 0.70, 11.5, C_TEXT)
                elif layout == 'closing':
                    oval(sl, 2.0, 0.5, 9.0, 7.0, _acc(0.07))
                    txt(sl, kicker, 0, 1.10, 13.33, 0.30, 9,
                        C_ACC, bold=True, align=PP_ALIGN.CENTER)
                    txt(sl, title, 0.80, 1.50, 11.73, 2.50, 46, C_TEXT,
                        bold=True, fn=FH, align=PP_ALIGN.CENTER)
                    rect(sl, 4.66, 4.10, 4.0, 0.024, C_ACC)
                    if sub:
                        txt(sl, sub, 1.0, 4.28, 11.33, 0.72, 13.5,
                            C_DIM, italic=True, align=PP_ALIGN.CENTER)
                    if stats:
                        n = min(len(stats), 3)
                        cw = 3.5
                        sx = (13.33-(n*cw+(n-1)*0.28))/2
                        for bi, st in enumerate(stats[:3]):
                            bx = sx+bi*(cw+0.28)
                            rect(sl, bx, 5.10, cw, 1.70, _bg2(
                                0.60), border=_acc(0.30), bw=0.5)
                            rect(sl, bx, 5.10, cw, 0.024, C_ACC)
                            txt(sl, st.get('value', ''), bx+0.14, 5.24, cw-0.28, 0.84,
                                30, C_ACC, bold=True, fn=FH, align=PP_ALIGN.CENTER)
                            txt(sl, st.get('label', '').upper(), bx+0.14, 6.18, cw -
                                0.28, 0.30, 8, C_DIM, bold=True, align=PP_ALIGN.CENTER)
                else:  # spotlight
                    txt(sl, kicker, 0.50, 0.14, 5.5, 0.26, 8, C_ACC, bold=True)
                    txt(sl, title,  0.50, 0.72, 7.5, 2.20,
                        32, C_TEXT, bold=True, fn=FH)
                    if sub:
                        txt(sl, sub, 0.50, 2.88, 7.5,
                            0.60, 12.5, C_DIM, italic=True)
                    rect(sl, 0.50, 3.58, 5.0, 0.024, C_ACC)
                    if pts:
                        for bi, pt in enumerate(pts):
                            ty = 3.76 + bi*0.62
                            rect(sl, 0.50, ty+0.14, 0.10, 0.10, C_ACC)
                            txt(sl, pt, 0.74, ty, 7.0, 0.56, 12, C_TEXT)
                    elif obj:
                        txt(sl, obj, 0.50, 3.76, 7.5,
                            0.80, 12, C_DIM, italic=True)
                    rect(sl, 8.30, 0.64, 4.60, 6.58, _bg2(
                        0.65), border=_acc(0.24), bw=0.5)
                    if stats:
                        for bi, st in enumerate(stats[:3]):
                            ty = 1.0 + bi*2.14
                            rect(sl, 8.30, ty, 4.60, 1.90, _bg2(
                                0.60), border=_acc(0.22), bw=0.4)
                            txt(sl, st.get('value', ''), 8.46, ty+0.18, 4.28, 1.10,
                                36, C_ACC, bold=True, fn=FH, align=PP_ALIGN.CENTER)
                            txt(sl, st.get('label', '').upper(), 8.46, ty+1.38,
                                4.28, 0.30, 8, C_DIM, bold=True, align=PP_ALIGN.CENTER)
                    else:
                        txt(sl, 'VISUAL DIRECTION', 8.50, 0.82,
                            4.20, 0.26, 7.5, C_ACC, bold=True)
                        rect(sl, 8.30, 1.06, 4.60, 0.024, C_ACC)
                        txt(sl, vis or 'Use a clean structured visual.',
                            8.50, 1.16, 4.20, 5.70, 12, C_TEXT)

            # ══════════════════════════════════════════════════════════════════
            # STYLE F — FOREST  (organic, bottom-anchored stats, nature motif)
            # DNA: large bottom accent bar, organic orbs, bottom-up composition
            # ══════════════════════════════════════════════════════════════════
            else:  # forest
                # Bottom anchor bar (green ground)
                rect(sl, 0, 6.70, 13.33, 0.80, _bg2(0.88))
                rect(sl, 0, 6.70, 13.33, 0.05, C_ACC)

                if layout == 'cover':
                    oval(sl, -1.0, -1.0, 7.0, 7.0, _acc(0.14))
                    oval(sl, 8.5,  2.0, 6.5, 6.5, _acc(0.10))
                    rect(sl, 0, 0, 0.08, 6.70, C_ACC)
                    txt(sl, kicker, 0.48, 0.40, 5.0, 0.28, 8, C_ACC, bold=True)
                    txt(sl, title, 0.48, 0.88, 7.8, 3.20,
                        42, C_TEXT, bold=True, fn=FH)
                    if sub:
                        txt(sl, sub, 0.48, 4.18, 7.0,
                            0.80, 13.5, C_DIM, italic=True)
                    rect(sl, 0.48, 5.08, 3.6, 0.04, C_ACC)
                    txt(sl, idea.title()[:40], 0.48,
                        5.26, 6.0, 0.40, 10.5, C_DIM)
                    if stats:
                        s0 = stats[0]
                        txt(sl, s0.get('value', ''), 7.5, 1.60, 5.5, 2.40, 68, _acc(
                            0.88), bold=True, fn=FH, align=PP_ALIGN.CENTER)
                        txt(sl, s0.get('label', '').upper(), 7.5, 4.10, 5.5,
                            0.40, 10.5, C_DIM, bold=True, align=PP_ALIGN.CENTER)
                    txt(sl, 'VentureOS', 0.48, 6.78, 4.0, 0.22, 7, C_DIM)
                    txt(sl, theme_name, 4.66, 6.78, 4.0,
                        0.22, 7, C_DIM, align=PP_ALIGN.CENTER)
                    txt(sl, f'{i+1:02d} / {total:02d}', 9.50, 6.78, 3.3,
                        0.22, 7, C_DIM, bold=True, align=PP_ALIGN.RIGHT)
                    footer(sl, i, total)
                    continue
                elif layout == 'metrics':
                    oval(sl, 9.0, -1.0, 5.5, 5.5, _acc(0.14))
                    rect(sl, 0, 0, 13.33, 0.80, _bg2(0.88))
                    txt(sl, kicker, 0.50, 0.12, 6.0, 0.26, 8, C_ACC, bold=True)
                    txt(sl, title,  0.50, 0.42, 9.0, 0.80,
                        26, C_TEXT, bold=True, fn=FH)
                    if stats:
                        n = len(stats)
                        cw = (12.33 - 0.18*(n-1)) / n
                        for bi, st in enumerate(stats):
                            bx = 0.50 + bi*(cw+0.18)
                            rect(sl, bx, 1.20, cw, 3.80, _sur(0.72),
                                 border=_acc(0.22), bw=0.5, rnd=True)
                            rect(sl, bx, 1.20, cw, 0.12, C_ACC)
                            txt(sl, st.get('value', ''), bx+0.18, 1.46, cw-0.36, 1.80,
                                50, C_TEXT, bold=True, fn=FH, align=PP_ALIGN.CENTER)
                            txt(sl, st.get('label', '').upper(), bx+0.18, 3.38, cw-0.36,
                                0.36, 9, C_DIM, bold=True, align=PP_ALIGN.CENTER)
                    for bi, pt in enumerate((pts or [obj])[:3]):
                        ty = 5.42 + bi*0.38
                        oval(sl, 0.50, ty+0.10, 0.10, 0.10, C_ACC)
                        txt(sl, pt, 0.74, ty, 12.0, 0.34, 11, C_TEXT)
                elif layout == 'roadmap':
                    oval(sl, 10.0, -0.5, 4.5, 4.5, _acc(0.12))
                    rect(sl, 0, 0, 13.33, 0.80, _bg2(0.88))
                    txt(sl, kicker, 0.50, 0.12, 6.0, 0.26, 8, C_ACC, bold=True)
                    txt(sl, title,  0.50, 0.42, 9.0, 0.72,
                        24, C_TEXT, bold=True, fn=FH)
                    rect(sl, 0.55, 2.88, 12.22, 0.04, _acc(0.32))
                    steps = pts[:4]
                    n = len(steps) or 1
                    sw = 12.22/n
                    for bi, st in enumerate(steps):
                        cx = 0.55+bi*sw+sw/2
                        oval(sl, cx-0.30, 2.70, 0.60, 0.60, _acc(0.38))
                        oval(sl, cx-0.18, 2.84, 0.36, 0.36, C_ACC)
                        txt(sl, f'{bi+1}', cx-0.18, 2.84, 0.36, 0.36, 9,
                            C_BG if not is_light else C_WHITE, bold=True, align=PP_ALIGN.CENTER)
                        bx = 0.55+bi*sw+0.10
                        cw = sw-0.20
                        rect(sl, bx, 3.54, cw, 2.86, _sur(0.70),
                             border=_acc(0.20), bw=0.5, rnd=True)
                        txt(sl, st, bx+0.14, 3.68, cw-0.28, 2.54, 11.5, C_TEXT)
                elif layout == 'comparison':
                    oval(sl, 11.0, -0.5, 4.0, 4.0, _acc(0.12))
                    rect(sl, 0, 0, 13.33, 0.72, _bg2(0.88))
                    rect(sl, 6.46, 0.72, 0.05, 5.98, _acc(0.34))
                    txt(sl, kicker, 0.50, 0.12, 5.0, 0.26, 8, C_ACC, bold=True)
                    txt(sl, title,  0.50, 0.42, 12.0,
                        0.72, 22, C_TEXT, bold=True, fn=FH)
                    rect(sl, 0.50, 0.88, 5.68, 0.40, _acc(0.20), rnd=True)
                    txt(sl, 'BEFORE', 0.68, 0.94, 5.30,
                        0.28, 8.5, C_ACC, bold=True)
                    rect(sl, 6.68, 0.88, 6.12, 0.40, _acc(0.36), rnd=True)
                    txt(sl, 'AFTER', 6.86, 0.94, 5.74,
                        0.28, 8.5, C_ACC, bold=True)
                    lp = pts[:3]
                    rp = pts[3:] or pts[:3]
                    for bi, p in enumerate(lp):
                        ty = 1.48 + bi*1.0
                        oval(sl, 0.54, ty+0.12, 0.12, 0.12, _acc(0.50))
                        txt(sl, p, 0.80, ty, 5.28, 0.88, 11.5, C_TEXT)
                    for bi, p in enumerate(rp):
                        ty = 1.48 + bi*1.0
                        oval(sl, 6.72, ty+0.12, 0.12, 0.12, C_ACC)
                        txt(sl, p, 6.98, ty, 5.80, 0.88, 11.5, C_TEXT)
                elif layout == 'closing':
                    oval(sl, -1.5, -1.0, 7.0, 7.0, _acc(0.12))
                    oval(sl, 9.0,  1.5, 6.0, 6.0, _acc(0.09))
                    rect(sl, 0, 0, 13.33, 0.08, C_ACC)
                    txt(sl, kicker, 0, 1.0, 13.33, 0.30, 9,
                        C_ACC, bold=True, align=PP_ALIGN.CENTER)
                    txt(sl, title, 0.80, 1.44, 11.73, 2.20, 44, C_TEXT,
                        bold=True, fn=FH, align=PP_ALIGN.CENTER)
                    rect(sl, 4.66, 3.74, 4.0, 0.05, C_ACC)
                    if sub:
                        txt(sl, sub, 1.0, 3.92, 11.33, 0.68, 13.5,
                            C_DIM, italic=True, align=PP_ALIGN.CENTER)
                    if stats:
                        n = min(len(stats), 3)
                        cw = 3.5
                        sx = (13.33-(n*cw+(n-1)*0.28))/2
                        for bi, st in enumerate(stats[:3]):
                            bx = sx+bi*(cw+0.28)
                            rect(sl, bx, 4.80, cw, 1.70, _sur(0.75),
                                 border=_acc(0.30), bw=0.6, rnd=True)
                            rect(sl, bx, 4.80, cw, 0.10, C_ACC)
                            txt(sl, st.get('value', ''), bx+0.14, 4.95, cw-0.28, 0.84,
                                30, C_TEXT, bold=True, fn=FH, align=PP_ALIGN.CENTER)
                            txt(sl, st.get('label', '').upper(), bx+0.14, 5.90, cw -
                                0.28, 0.30, 8, C_DIM, bold=True, align=PP_ALIGN.CENTER)
                else:  # spotlight
                    oval(sl, 9.5, -0.5, 5.0, 5.0, _acc(0.14))
                    rect(sl, 0, 0, 0.08, 6.70, C_ACC)
                    rect(sl, 7.88, 0, 5.45, 6.70, _bg2(0.85))
                    txt(sl, kicker, 0.38, 0.38, 5.5, 0.26, 8, C_ACC, bold=True)
                    txt(sl, title,  0.38, 0.76, 7.2, 1.80,
                        28, C_TEXT, bold=True, fn=FH)
                    if sub:
                        txt(sl, sub, 0.38, 2.52, 7.0,
                            0.58, 12, C_DIM, italic=True)
                    rect(sl, 0.38, 3.20, 4.8, 0.04, C_ACC)
                    if pts:
                        for bi, pt in enumerate(pts):
                            ty = 3.38 + bi*0.58
                            oval(sl, 0.38, ty+0.12, 0.10, 0.10, C_ACC)
                            txt(sl, pt, 0.62, ty, 6.90, 0.52, 11.5, C_TEXT)
                    elif obj:
                        txt(sl, obj, 0.38, 3.38, 7.0,
                            0.80, 12, C_DIM, italic=True)
                    if stats:
                        for bi, st in enumerate(stats[:3]):
                            ty = 0.50 + bi*2.02
                            rect(sl, 8.08, ty, 4.92, 1.82, _sur(0.70),
                                 border=_acc(0.22), bw=0.5, rnd=True)
                            rect(sl, 8.08, ty, 4.92, 0.12, C_ACC)
                            txt(sl, st.get('value', ''), 8.26, ty+0.17, 4.56, 1.0,
                                34, C_TEXT, bold=True, fn=FH, align=PP_ALIGN.CENTER)
                            txt(sl, st.get('label', '').upper(), 8.26, ty+1.28,
                                4.56, 0.30, 8, C_DIM, bold=True, align=PP_ALIGN.CENTER)
                    else:
                        rect(sl, 8.08, 0.50, 4.92, 5.90, _sur(0.70),
                             border=_acc(0.18), bw=0.5, rnd=True)
                        txt(sl, 'VISUAL DIRECTION', 8.28, 0.68,
                            4.52, 0.26, 7.5, C_ACC, bold=True)
                        rect(sl, 8.08, 0.92, 4.92, 0.04, _acc(0.32))
                        txt(sl, vis or 'One compelling nature-inspired visual.',
                            8.28, 1.02, 4.52, 5.10, 12, C_TEXT)

            footer(sl, i, total)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        fname = re.sub(r'[^a-z0-9]+', '-', idea.lower()
                       )[:40] + '-pitch-deck.pptx'
        return send_file(buf,
                         mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                         as_attachment=True, download_name=fname)

    except Exception as e:
        import traceback
        return jsonify({'error': str(e), 'trace': traceback.format_exc()}), 500
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

        data = request.json
        idea = data.get('idea', 'ventureos')
        raw_slides = data.get('slides') or data.get('deck') or []
        slides_data = [_normalize_slide(s, i)
                       for i, s in enumerate(raw_slides)]
        if not slides_data:
            return jsonify({'error': 'No slides provided'}), 400

        # ── Theme & colour setup ───────────────────────────────────────────
        design = data.get('design_system') or DEFAULT_DECK_THEME
        palette = design.get('palette') or DEFAULT_DECK_THEME['palette']
        template_id = design.get(
            'template_id', DEFAULT_DECK_THEME['template_id'])
        theme_name = design.get('theme_name', 'VentureOS')

        def _h2r(val, fallback):
            v = _normalize_hex(val, fallback).replace('#', '')
            return RGBColor(int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))

        def _luma(val, fb='#000000'):
            v = _normalize_hex(val, fb).replace('#', '')
            return 0.299*int(v[0:2], 16) + 0.587*int(v[2:4], 16) + 0.114*int(v[4:6], 16)

        C_BG = _h2r(palette.get('primary'),
                    DEFAULT_DECK_THEME['palette']['primary'])
        C_BG2 = _h2r(palette.get('secondary'),
                     DEFAULT_DECK_THEME['palette']['secondary'])
        C_SURF = _h2r(palette.get('surface'),
                      DEFAULT_DECK_THEME['palette']['surface'])
        C_ACC = _h2r(palette.get('accent'),
                     DEFAULT_DECK_THEME['palette']['accent'])
        C_TEXT = _h2r(palette.get('text'),
                      DEFAULT_DECK_THEME['palette']['text'])
        is_light = _luma(palette.get('primary', '#000')) > 160
        C_DIM = RGBColor(0x44, 0x55, 0x66) if is_light else RGBColor(
            0xB0, 0xC4, 0xD8)
        C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        C_BLACK = RGBColor(0x0F, 0x0F, 0x0F)

        def _blend(a: RGBColor, b: RGBColor, t: float) -> RGBColor:
            return RGBColor(int(a[0]+(b[0]-a[0])*t),
                            int(a[1]+(b[1]-a[1])*t),
                            int(a[2]+(b[2]-a[2])*t))

        def _acc(t): return _blend(C_BG, C_ACC, t)
        def _sur(t): return _blend(C_BG, C_SURF, t)
        def _bg2(t): return _blend(C_BG, C_BG2, t)

        font_pairs = {
            'editorial-midnight': ('Georgia',          'Calibri'),
            'boardroom-ivory':    ('Cambria',           'Calibri'),
            'kinetic-ember':      ('Arial Black',       'Arial'),
            'atlas-sapphire':     ('Georgia',           'Calibri'),
            'monochrome-luxe':    ('Palatino Linotype', 'Calibri'),
            'forest-venture':     ('Georgia',           'Calibri'),
            'crimson-capital':    ('Georgia',           'Calibri'),
            'arctic-clarity':     ('Trebuchet MS',      'Calibri'),
            'violet-epoch':       ('Georgia',           'Calibri'),
        }
        FH, FB = font_pairs.get(template_id, ('Georgia', 'Calibri'))

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        BLANK = prs.slide_layouts[6]

        # ── Primitive helpers ──────────────────────────────────────────────
        def txt(sl, text, L, T, W, H, sz, col,
                bold=False, italic=False, align=PP_ALIGN.LEFT, fn=None):
            if not text:
                return
            tb = sl.shapes.add_textbox(
                Inches(L), Inches(T), Inches(W), Inches(H))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = align
            r = p.add_run()
            r.text = str(text)
            r.font.size = Pt(sz)
            r.font.color.rgb = col
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = fn or FB
            return tb

        def rect(sl, L, T, W, H, fill, border=None, bw=0.5, rnd=False):
            st = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rnd else MSO_AUTO_SHAPE_TYPE.RECTANGLE
            sh = sl.shapes.add_shape(st, Inches(
                L), Inches(T), Inches(W), Inches(H))
            sh.fill.solid()
            sh.fill.fore_color.rgb = fill
            if border:
                sh.line.color.rgb = border
                sh.line.width = Pt(bw)
            else:
                sh.line.fill.background()
            return sh

        def oval(sl, L, T, W, H, fill):
            sh = sl.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.OVAL, Inches(L), Inches(T), Inches(W), Inches(H))
            sh.fill.solid()
            sh.fill.fore_color.rgb = fill
            sh.line.fill.background()
            return sh

        # ── Shared chrome: footer ──────────────────────────────────────────
        def footer(sl, idx, total):
            rect(sl, 0, 7.22, 13.33, 0.28, _bg2(0.8), border=None)
            rect(sl, 0, 7.22, 13.33, 0.014, _acc(0.35), border=None)
            txt(sl, 'VentureOS',   0.40, 7.28, 3.0, 0.22, 7.5, C_DIM)
            txt(sl, theme_name,    4.66, 7.28, 4.0, 0.22, 7.5, C_DIM,
                align=PP_ALIGN.CENTER)
            txt(sl, f'{idx+1:02d} / {total:02d}',
                9.50, 7.28, 3.3, 0.22, 7.5, C_DIM,
                bold=True, align=PP_ALIGN.RIGHT)

        total = len(slides_data)

        for i, sd in enumerate(slides_data):
            sl = prs.slides.add_slide(BLANK)
            layout = sd.get('layout', 'spotlight')
            title = sd.get('title') or f'Slide {i+1}'
            sub = sd.get('subtitle') or ''
            pts = (sd.get('content') or sd.get('points') or [])[:5]
            stats = (sd.get('stats') or [])[:3]
            obj = sd.get('objective') or ''
            vis = _presentation_visual_summary(idea, sd, i)
            kicker = sd.get('type', 'story').replace('_', ' ').upper()

            # ── Full background ────────────────────────────────────────────
            sl.background.fill.solid()
            sl.background.fill.fore_color.rgb = C_BG

            # ══════════════════════════════════════════════════════════════
            # LAYOUT 1 — COVER  (full bleed, big centred hero)
            # ══════════════════════════════════════════════════════════════
            if layout == 'cover':
                # Left accent panel (40% wide)
                rect(sl, 0, 0, 5.4, 7.5, _bg2(0.9))
                # Decorative large orb bottom-left
                oval(sl, -1.2, 3.8, 5.5, 5.5, _acc(0.14))
                oval(sl, -0.4, 4.4, 3.2, 3.2, _acc(0.20))
                # Right panel slightly lighter
                rect(sl, 5.4, 0, 7.93, 7.5, _bg2(0.4))
                # Big accent orb top-right
                oval(sl, 7.5, -1.5, 7.0, 7.0, _acc(0.18))
                oval(sl, 9.5,  3.5, 4.5, 4.5, _acc(0.10))

                # Thick left accent stripe
                rect(sl, 0, 0, 0.07, 7.5, C_ACC)

                # Kicker pill
                rect(sl, 0.55, 0.55, 2.8, 0.34, _acc(0.28), rnd=True)
                txt(sl, kicker, 0.72, 0.59, 2.6, 0.26, 8, C_ACC, bold=True)

                # Giant title — left panel
                txt(sl, title, 0.55, 1.10, 4.6, 2.8, 40, C_TEXT,
                    bold=True, fn=FH)

                # Subtitle
                if sub:
                    txt(sl, sub, 0.55, 3.80, 4.5, 0.9, 13.5, C_DIM,
                        italic=True)

                # Horizontal rule
                rect(sl, 0.55, 4.72, 3.8, 0.04, _acc(0.55))

                # Idea / company label bottom-left
                txt(sl, idea.title()[:48], 0.55, 4.90, 4.4, 0.5, 11, C_DIM)

                # Right side — big decorative stat or visual cue
                if stats:
                    s0 = stats[0]
                    txt(sl, s0.get('value', ''), 6.0, 1.8, 7.0, 2.0, 72,
                        _acc(0.9), bold=True, fn=FH, align=PP_ALIGN.CENTER)
                    txt(sl, s0.get('label', '').upper(), 6.0, 4.2, 7.0, 0.4, 11,
                        C_DIM, bold=True, align=PP_ALIGN.CENTER)
                else:
                    txt(sl, vis or 'The future starts here.', 6.0, 2.6, 7.0,
                        1.6, 16, C_DIM, italic=True, align=PP_ALIGN.CENTER)

            # ══════════════════════════════════════════════════════════════
            # LAYOUT 2 — METRICS  (big numbers dominate the slide)
            # ══════════════════════════════════════════════════════════════
            elif layout == 'metrics':
                # Top accent band
                rect(sl, 0, 0, 13.33, 1.40, _bg2(0.85))
                rect(sl, 0, 1.40, 13.33, 0.05, _acc(0.50))

                # Kicker top-left
                txt(sl, kicker, 0.55, 0.20, 6.0, 0.28, 8, C_ACC, bold=True)

                # Title in top band
                txt(sl, title, 0.55, 0.52, 8.5, 0.80, 28, C_TEXT,
                    bold=True, fn=FH)

                # Stat cards — 3-up large format
                n_stats = len(stats) if stats else 0
                if n_stats:
                    card_w = (12.33 - 0.2*(n_stats-1)) / n_stats
                    for bi, st in enumerate(stats):
                        bx = 0.50 + bi*(card_w+0.20)
                        # Card bg
                        rect(sl, bx, 1.70, card_w, 3.20, _sur(0.75),
                             border=_acc(0.25), bw=0.6, rnd=True)
                        # Top accent bar inside card
                        rect(sl, bx, 1.70, card_w, 0.14, C_ACC, rnd=False)
                        # Huge stat value
                        txt(sl, st.get('value', ''), bx+0.20, 1.98,
                            card_w-0.40, 1.60, 52, C_TEXT,
                            bold=True, fn=FH, align=PP_ALIGN.CENTER)
                        # Label
                        txt(sl, st.get('label', '').upper(),
                            bx+0.20, 3.60, card_w-0.40, 0.38,
                            9, C_DIM, bold=True, align=PP_ALIGN.CENTER)

                # Bullet points below cards
                if pts:
                    for bi, pt in enumerate(pts[:3]):
                        ty = 5.22 + bi*0.42
                        oval(sl, 0.55, ty+0.10, 0.12, 0.12, C_ACC)
                        txt(sl, pt, 0.80, ty, 12.0, 0.38, 11.5, C_TEXT)
                elif obj:
                    txt(sl, obj, 0.55, 5.22, 12.0, 0.60, 12, C_DIM, italic=True)

            # ══════════════════════════════════════════════════════════════
            # LAYOUT 3 — ROADMAP  (numbered horizontal timeline)
            # ══════════════════════════════════════════════════════════════
            elif layout == 'roadmap':
                # Top band
                rect(sl, 0, 0, 13.33, 1.30, _bg2(0.85))
                rect(sl, 0, 1.30, 13.33, 0.05, _acc(0.45))

                txt(sl, kicker, 0.55, 0.18, 6.0, 0.28, 8, C_ACC, bold=True)
                txt(sl, title,  0.55, 0.48, 9.0, 0.76, 26, C_TEXT,
                    bold=True, fn=FH)
                if sub:
                    txt(sl, sub, 0.55, 1.05, 8.0, 0.34, 11, C_DIM, italic=True)

                # Horizontal connector line
                rect(sl, 0.55, 3.08, 12.22, 0.04, _acc(0.30))

                steps = pts[:4] if pts else []
                n = len(steps) or 1
                step_w = 12.22 / n
                for bi, step_txt in enumerate(steps):
                    cx = 0.55 + bi*step_w + step_w/2
                    # Circle on timeline
                    oval(sl, cx-0.35, 2.88, 0.70, 0.70, _acc(0.40))
                    oval(sl, cx-0.22, 3.01, 0.44, 0.44, C_ACC)
                    txt(sl, f'{bi+1}', cx-0.22, 3.01, 0.44, 0.44,
                        10, C_BG if not is_light else C_WHITE,
                        bold=True, align=PP_ALIGN.CENTER)
                    # Card below
                    bx = 0.55 + bi*step_w + 0.10
                    cw = step_w - 0.20
                    rect(sl, bx, 3.85, cw, 2.60, _sur(0.7),
                         border=_acc(0.20), bw=0.5, rnd=True)
                    txt(sl, step_txt, bx+0.16, 3.98, cw-0.32, 2.30,
                        11.5, C_TEXT)

                # Visual suggestion bottom strip
                if vis:
                    txt(sl, f'Visual: {vis}', 0.55, 6.70, 12.22, 0.30,
                        8.5, C_DIM, italic=True)

            # ══════════════════════════════════════════════════════════════
            # LAYOUT 4 — COMPARISON  (two-column split)
            # ══════════════════════════════════════════════════════════════
            elif layout == 'comparison':
                # Vertical divider
                rect(sl, 6.46, 0.9, 0.06, 6.3, _acc(0.35))

                # Header band
                rect(sl, 0, 0, 13.33, 0.90, _bg2(0.88))
                rect(sl, 0, 0.90, 13.33, 0.04, _acc(0.40))
                txt(sl, kicker, 0.50, 0.12, 5.0, 0.26, 8, C_ACC, bold=True)
                txt(sl, title,  0.50, 0.38, 12.3, 0.50, 22, C_TEXT,
                    bold=True, fn=FH)

                # Left column header
                rect(sl, 0.50, 1.08, 5.76, 0.44, _acc(0.22), rnd=True)
                txt(sl, 'PROBLEM / STATUS QUO', 0.68, 1.14, 5.40, 0.32,
                    9, C_ACC, bold=True)

                # Right column header
                rect(sl, 6.72, 1.08, 6.10, 0.44, _acc(0.35), rnd=True)
                txt(sl, 'OUR SOLUTION', 6.90, 1.14, 5.70, 0.32,
                    9, C_ACC, bold=True)

                left_pts = pts[:3]
                right_pts = pts[3:] or pts[:3]
                for bi, p in enumerate(left_pts):
                    ty = 1.76 + bi*0.96
                    oval(sl, 0.55, ty+0.14, 0.14, 0.14, _acc(0.50))
                    txt(sl, p, 0.82, ty, 5.30, 0.84, 12, C_TEXT)

                for bi, p in enumerate(right_pts):
                    ty = 1.76 + bi*0.96
                    oval(sl, 6.78, ty+0.14, 0.14, 0.14, C_ACC)
                    txt(sl, p, 7.06, ty, 5.82, 0.84, 12, C_TEXT)

                if sub:
                    txt(sl, sub, 0.50, 6.60, 12.3, 0.40, 10.5,
                        C_DIM, italic=True)

            # ══════════════════════════════════════════════════════════════
            # LAYOUT 5 — CLOSING / CTA  (centred, bold, minimal)
            # ══════════════════════════════════════════════════════════════
            elif layout == 'closing':
                # Full-bleed accent bar at top
                rect(sl, 0, 0, 13.33, 0.10, C_ACC)
                # Large ambient orbs
                oval(sl, -2.0, -1.5, 8.0, 8.0, _acc(0.12))
                oval(sl, 7.5,  1.0,  7.0, 7.0, _acc(0.10))

                # Centred layout
                txt(sl, kicker, 0, 1.10, 13.33, 0.32, 9, C_ACC,
                    bold=True, align=PP_ALIGN.CENTER)

                # Giant title centred
                txt(sl, title, 0.80, 1.55, 11.73, 2.0, 44, C_TEXT,
                    bold=True, fn=FH, align=PP_ALIGN.CENTER)

                rect(sl, 4.66, 3.55, 4.0, 0.05, _acc(0.55))

                if sub:
                    txt(sl, sub, 1.0, 3.72, 11.33, 0.70, 14, C_DIM,
                        italic=True, align=PP_ALIGN.CENTER)

                # CTA stat callouts (up to 3)
                if stats:
                    n_s = len(stats)
                    sw = min(n_s, 3)
                    cw = 3.6
                    start_x = (13.33 - (sw*cw + (sw-1)*0.30)) / 2
                    for bi, st in enumerate(stats[:3]):
                        bx = start_x + bi*(cw+0.30)
                        rect(sl, bx, 4.60, cw, 1.60, _sur(0.75),
                             border=_acc(0.30), bw=0.6, rnd=True)
                        rect(sl, bx, 4.60, cw, 0.10, C_ACC, rnd=False)
                        txt(sl, st.get('value', ''), bx+0.15, 4.76,
                            cw-0.30, 0.80, 32, C_TEXT,
                            bold=True, fn=FH, align=PP_ALIGN.CENTER)
                        txt(sl, st.get('label', '').upper(),
                            bx+0.15, 5.60, cw-0.30, 0.32,
                            8.5, C_DIM, bold=True, align=PP_ALIGN.CENTER)
                elif pts:
                    combined = '   ·   '.join(pts[:3])
                    txt(sl, combined, 1.0, 4.70, 11.33, 0.50, 12,
                        C_DIM, align=PP_ALIGN.CENTER)

                # Bottom accent bar
                rect(sl, 0, 7.40, 13.33, 0.10, C_ACC)

            # ══════════════════════════════════════════════════════════════
            # LAYOUT 6 — SPOTLIGHT  (default: split left-text / right-card)
            # ══════════════════════════════════════════════════════════════
            else:
                # Right panel background
                rect(sl, 7.90, 0, 5.43, 7.5, _bg2(0.85))
                # Top-right orb
                oval(sl, 8.5, -1.0, 5.5, 5.5, _acc(0.18))
                oval(sl, 9.5,  4.0, 3.8, 3.8, _acc(0.10))

                # Left accent bar
                rect(sl, 0, 0, 0.08, 7.5, C_ACC)

                # Kicker
                txt(sl, kicker, 0.38, 0.42, 5.5, 0.28, 8, C_ACC, bold=True)

                # Title
                txt(sl, title, 0.38, 0.82, 7.2, 1.80, 30, C_TEXT,
                    bold=True, fn=FH)

                # Subtitle
                if sub:
                    txt(sl, sub, 0.38, 2.56, 7.0, 0.60, 12.5, C_DIM,
                        italic=True)

                # Thin rule
                rect(sl, 0.38, 3.22, 5.0, 0.04, _acc(0.40))

                # Bullet points
                if pts:
                    for bi, pt in enumerate(pts):
                        ty = 3.42 + bi*0.58
                        rect(sl, 0.38, ty+0.08, 0.22, 0.22, _acc(0.30),
                             rnd=True)
                        txt(sl, f'{bi+1}', 0.38, ty+0.08, 0.22, 0.22,
                            7, C_ACC, bold=True, align=PP_ALIGN.CENTER)
                        txt(sl, pt, 0.72, ty, 6.90, 0.52, 12, C_TEXT)
                elif obj:
                    txt(sl, obj, 0.38, 3.42, 7.0, 0.80, 12.5, C_DIM,
                        italic=True)

                # Stat boxes stacked on right panel
                if stats:
                    for bi, st in enumerate(stats[:3]):
                        ty = 0.70 + bi*2.10
                        rect(sl, 8.10, ty, 4.90, 1.82, _sur(0.7),
                             border=_acc(0.22), bw=0.5, rnd=True)
                        rect(sl, 8.10, ty, 4.90, 0.12, C_ACC, rnd=False)
                        txt(sl, st.get('value', ''), 8.28, ty+0.18,
                            4.52, 1.0, 36, C_TEXT,
                            bold=True, fn=FH, align=PP_ALIGN.CENTER)
                        txt(sl, st.get('label', '').upper(),
                            8.28, ty+1.30, 4.52, 0.30,
                            8.5, C_DIM, bold=True, align=PP_ALIGN.CENTER)
                else:
                    # Visual cue card
                    rect(sl, 8.10, 0.70, 4.90, 5.80, _sur(0.7),
                         border=_acc(0.18), bw=0.5, rnd=True)
                    txt(sl, 'VISUAL DIRECTION', 8.30, 0.88, 4.50, 0.28,
                        7.5, C_ACC, bold=True)
                    rect(sl, 8.10, 1.12, 4.90, 0.04, _acc(0.30))
                    txt(sl, vis or 'Use one dominant proof visual.',
                        8.30, 1.22, 4.50, 4.90, 12, C_TEXT)

            # ── Shared footer ──────────────────────────────────────────────
            footer(sl, i, total)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        fname = re.sub(r'[^a-z0-9]+', '-', idea.lower()
                       )[:40] + '-pitch-deck.pptx'
        return send_file(buf,
                         mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                         as_attachment=True, download_name=fname)

    except Exception as e:
        import traceback
        return jsonify({'error': str(e), 'trace': traceback.format_exc()}), 500


# ── SHARE / REPORT ─────────────────────────────────────────────────────────
@app.route('/report/save', methods=['POST'])
def save_report():
    data = request.json
    report_id = str(uuid.uuid4())[:8]
    path = os.path.join(REPORTS_DIR, f'{report_id}.json')
    with open(path, 'w') as f:
        json.dump(data, f)
    return jsonify({'id': report_id, 'url': f'/report/{report_id}'})


@app.route('/report/<report_id>')
def view_report(report_id):
    path = os.path.join(REPORTS_DIR, f'{report_id}.json')
    if not os.path.exists(path):
        return 'Report not found', 404
    return render_template('index.html')


# ── RUN ────────────────────────────────────────────────────────────────────
if __name__ == '__main__':
    app.run(debug=True, threaded=True)
